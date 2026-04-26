"""
FastAPI backend for traffic route analysis.
Provides endpoints for autocomplete, route analysis, and serving the frontend.
"""

import os
import json
import logging
import re
import csv
import random
import socket
import socketio
from pathlib import Path
from dotenv import load_dotenv
from typing import Optional, Union, List
from functools import wraps
from fastapi import FastAPI, HTTPException, Query, Depends, status, BackgroundTasks, Request
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse, FileResponse, Response, RedirectResponse
import io
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordRequestForm
from pydantic import BaseModel, Field, EmailStr, field_validator
from jose.exceptions import ExpiredSignatureError
import joblib
from datetime import datetime, UTC, timedelta, date
import secrets
import uuid
import traceback
import asyncio
import time
import threading
from sqlalchemy.exc import IntegrityError
from sqlalchemy import func

from .paths import DATA_DIR, ENV_FILE, STATIC_DIR, TEMPLATES_DIR

load_dotenv(ENV_FILE)

# Import logging and rate limiting
from .logging_config import setup_logging, get_logger
from .rate_limiter import RateLimitMiddleware

# Setup logging
setup_logging()
logger = get_logger(__name__)

from .utils import (
    clean_location,
    format_incident_type,
    format_officer_name,
    tomtom_geocode,
    tomtom_reverse_geocode_area,
    tomtom_autocomplete,
    tomtom_route,
    summarize_route,
    compute_route_cost,
    haversine_m
)
from .db import (
    init_db, get_session, get_db, save_analysis, AnalysisResult,
    User, SavedRoute, RouteRating, Notification, PoliceDispatchAssignment, OfficerDispatchStatus, DispatchLog, SharedAlert, Shift, ShiftAttendance, OfficerIncidentCount, MLFeedback, MLRetrainAudit
)
from sqlalchemy.orm import Session
from .auth import (
    verify_password, get_password_hash, create_access_token,
    get_current_user, get_current_active_user, get_current_admin_user,
    authenticate_user, create_user, get_user_by_username, Token, UserCreate as AuthUserCreate, UserResponse,
    get_optional_user, RoleLoginRequest, RoleToken, UserRole, create_role_access_token, require_role, require_any_role,
    require_police_department_user
)
from .analytics import (
    get_peak_hours_analysis, get_day_of_week_analysis,
    get_seasonal_trends, calculate_route_reliability, predict_future_congestion,
)
from .export_utils import export_to_csv, export_to_excel, export_to_pdf
from .notifications import (
    create_notification, check_traffic_alerts,
    suggest_best_time_to_leave, check_congestion_warnings,
    get_user_notifications, mark_notification_read
)
from .cache_utils import cached, clear_cache, get_cache_stats
from .realtime_utils import get_traffic_incidents, auto_refresh_route, monitor_route_changes
from .dispatch_notifications import send_officer_dispatch_notification

# Initialize FastAPI app
app = FastAPI(
    title="Traffic Route Analysis API",
    description="Real-time traffic congestion analysis with ML predictions",
    version="1.0.0"
)

# Socket.IO server for push updates to the command center frontend.
SOCKETIO_REDIS_URL = os.getenv("SOCKETIO_REDIS_URL", "redis://redis:6379/0")

# Try to use Redis manager for distributed systems, fall back to in-memory for development
try:
    socket_client_manager = socketio.AsyncRedisManager(SOCKETIO_REDIS_URL)
    logger.info(f"✅ Socket.IO using Redis manager: {SOCKETIO_REDIS_URL}")
except Exception as e:
    logger.warning(f"⚠️ Redis connection failed ({e}), using in-memory manager for Socket.IO")
    socket_client_manager = None  # Will use in-memory manager by default


socket_client_manager = None


def _can_reach_socketio_redis(redis_url: str) -> bool:
    """Only enable Redis pub/sub when the configured host is actually reachable."""
    try:
        from urllib.parse import urlparse

        parsed = urlparse(redis_url)
        host = parsed.hostname
        port = parsed.port or 6379
        if not host:
            return False

        with socket.create_connection((host, port), timeout=1.5):
            return True
    except Exception:
        return False


if _can_reach_socketio_redis(SOCKETIO_REDIS_URL):
    try:
        socket_client_manager = socketio.AsyncRedisManager(SOCKETIO_REDIS_URL)
        logger.info(f"Socket.IO using Redis manager: {SOCKETIO_REDIS_URL}")
    except Exception as e:
        logger.warning(f"Redis manager unavailable ({e}); using in-memory Socket.IO manager")
        socket_client_manager = None
else:
    logger.info(f"Redis not reachable at {SOCKETIO_REDIS_URL}; using in-memory Socket.IO manager")

sio = socketio.AsyncServer(
    async_mode="asgi",
    cors_allowed_origins="*",
    client_manager=socket_client_manager,  # None = use in-memory manager
    ping_timeout=60,  # Server expects ping from client within 60 seconds
    ping_interval=25,  # Server sends ping every 25 seconds
    max_http_buffer_size=1000000,  # 1MB max message size
    logger=True,  # Enable Socket.IO logger
    engineio_logger=False,  # Don't spam with engine.io logs
)
# Wrap FastAPI inside the Socket.IO ASGI app so Socket.IO handles its own
# /socket.io/* path and delegates everything else to FastAPI. This is the
# officially recommended pattern for python-socketio + Starlette/FastAPI.
# Do NOT use app.mount() for Socket.IO — it strips the path prefix and causes
# the client to connect to /socket.io/socket.io/... (doubled path → timeout).
asgi_app = socketio.ASGIApp(sio, other_asgi_app=app)

from .socketio_events import (
    register_police_socketio_handlers,
    emit_incident_new,
    emit_incident_updated,
    emit_officer_status_changed,
    emit_officer_dispatched,
)

register_police_socketio_handlers(sio, logger)

# ============================================================================
# REAL-TIME ALERTS INFRASTRUCTURE
# ============================================================================
# In-memory alert store: {district_id: [AlertData, ...]}
_alerts_store = {}
_alert_subscribers = {}  # {district_id: [callback_func, ...]}
_manual_incidents_store: dict[str, list[dict]] = {}
_patrol_user_assignments: dict[str, dict[str, dict[str, object]]] = {}
_patrol_user_assignments_lock = threading.Lock()


def add_alert(district_id: str, severity: str, message: str, incident_id: Optional[str] = None):
    """Add an alert to the store and notify subscribers."""
    from uuid import uuid4
    
    alert = AlertData(
        alert_id=str(uuid4()),
        severity=severity,
        message=message,
        timestamp=datetime.now(UTC).isoformat(),
        district_id=district_id,
        related_incident_id=incident_id,
    )
    
    if district_id not in _alerts_store:
        _alerts_store[district_id] = []
    
    _alerts_store[district_id].append(alert)
    
    # Keep only last 50 alerts per district to prevent memory bloat
    if len(_alerts_store[district_id]) > 50:
        _alerts_store[district_id] = _alerts_store[district_id][-50:]
    
    logger.info(f"Alert added: {severity} - {message} (District: {district_id})")
    return alert


def get_unread_alerts(district_id: str) -> list[AlertData]:
    """Get unread alerts for a district."""
    return _alerts_store.get(district_id, [])


@app.exception_handler(ExpiredSignatureError)
async def expired_signature_exception_handler(request: Request, exc: ExpiredSignatureError):
    accept_header = (request.headers.get("accept") or "").lower()
    if "text/html" in accept_header:
        return RedirectResponse(url="/auth/login?reason=session_expired", status_code=status.HTTP_307_TEMPORARY_REDIRECT)

    return JSONResponse(
        status_code=status.HTTP_401_UNAUTHORIZED,
        content={"detail": "Session expired. Please log in again."},
    )


@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    detail = exc.detail

    if exc.status_code == status.HTTP_403_FORBIDDEN:
        detail = "Access denied. You do not have permission to view this resource."
    elif exc.status_code == status.HTTP_401_UNAUTHORIZED and not detail:
        detail = "Unauthorized access. Please log in to continue."

    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": detail},
        headers=getattr(exc, "headers", None),
    )

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Add rate limiting middleware
app.add_middleware(RateLimitMiddleware)

templates = Jinja2Templates(directory=str(TEMPLATES_DIR))
templates.env.filters["clean_location"] = clean_location
templates.env.filters["format_incident_type"] = format_incident_type
templates.env.filters["format_officer_name"] = format_officer_name


def _resolve_officer_name(officer: Optional[User]) -> Optional[str]:
    if not officer:
        return None

    blocked_role_labels = {
        "police supervisor",
        "supervisor",
        "police officer",
        "officer",
    }

    full_name = (getattr(officer, "full_name", None) or "").strip()
    username = (getattr(officer, "username", None) or "").strip()

    if full_name and full_name.lower() not in blocked_role_labels:
        return full_name
    if username and username.lower() not in blocked_role_labels:
        return format_officer_name(username)


def _format_district_label(district_id: Optional[str]) -> str:
    value = str(district_id or "").strip()
    if not value:
        return "N/A"
    return value.replace("_", "-").upper()


def _resolve_current_user_district_id(current_user: dict) -> Optional[str]:
    district_id = str(current_user.get("district_id") or current_user.get("district") or "").strip()
    if district_id:
        return district_id

    username = str(current_user.get("username") or "").strip()
    if not username:
        return None

    session = get_session()
    try:
        user = get_user_by_username(session, username)
        return str(getattr(user, "district_id", "") or "").strip() or None
    finally:
        session.close()


def _bootstrap_app_state() -> None:
    global ML_MODEL, APP_BOOTSTRAP_DONE, APP_BOOTSTRAP_ERROR

    try:
        init_db()

        if os.path.exists(MODEL_PATH):
            try:
                ML_MODEL = joblib.load(MODEL_PATH)
                logger.info(f"Loaded ML model from {MODEL_PATH}")
            except Exception as model_error:
                logger.warning(f"Failed to load ML model: {model_error}")

        APP_BOOTSTRAP_DONE = True
        APP_BOOTSTRAP_ERROR = None
    except Exception as exc:
        APP_BOOTSTRAP_ERROR = str(exc)
        logger.exception("Application bootstrap failed")


def _start_bootstrap_once() -> None:
    global APP_BOOTSTRAP_STARTED

    with _bootstrap_lock:
        if APP_BOOTSTRAP_STARTED:
            return
        APP_BOOTSTRAP_STARTED = True

    threading.Thread(target=_bootstrap_app_state, daemon=True, name="app-bootstrap").start()


@app.on_event("startup")
async def schedule_app_bootstrap() -> None:
    _start_bootstrap_once()

ML_MODEL = None
_model_path_value = os.getenv("MODEL_PATH")
MODEL_PATH = Path(_model_path_value) if _model_path_value else DATA_DIR / "rf_model.pkl"
if not MODEL_PATH.is_absolute():
    MODEL_PATH = (Path.cwd() / MODEL_PATH).resolve()
APP_BOOTSTRAP_STARTED = False
APP_BOOTSTRAP_DONE = False
APP_BOOTSTRAP_ERROR: Optional[str] = None
_bootstrap_lock = threading.Lock()

# Mount static files if directory exists
if STATIC_DIR.exists():
    app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

# ============================================================================
# ERROR HANDLING DECORATOR
# ============================================================================

def handle_db_errors(func):
    """Decorator to handle database errors gracefully."""
    @wraps(func)
    async def wrapper(*args, **kwargs):
        try:
            return await func(*args, **kwargs)
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Database error in {func.__name__}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Database operation failed: {str(e)}"
            )
    return wrapper


DISTRICT_LOCATIONS = {
    "district_1": {"lat": 13.0827, "lon": 80.2707, "name": "Central District"},
    "district_2": {"lat": 13.0545, "lon": 80.2450, "name": "West District"},
    "district_3": {"lat": 12.9716, "lon": 80.1534, "name": "South District"},
    "district_4": {"lat": 13.1278, "lon": 80.2270, "name": "North District"},
}


def _normalize_severity(raw_severity: Optional[str]) -> str:
    value = str(raw_severity or "").lower()
    if value in {"1", "low", "minor"}:
        return "low"
    if value in {"2", "moderate", "medium"}:
        return "moderate"
    if value in {"3", "high"}:
        return "high"
    return "unknown"


def _severity_color(severity: Optional[str]) -> str:
    palette = {
        "low": "#22c55e",
        "moderate": "#f59e0b",
        "high": "#ef4444",
        "unknown": "#64748b",
    }
    return palette.get(_normalize_severity(severity), "#64748b")


def _load_police_incidents(district_id: str) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    incidents = get_traffic_incidents(district["lat"], district["lon"], radius=7000)

    normalized: list[dict] = []
    for index, incident in enumerate(incidents):
        location = incident.get("location") or []
        latitude = None
        longitude = None
        if isinstance(location, (list, tuple)) and len(location) >= 2:
            longitude, latitude = location[0], location[1]
        normalized.append({
            "id": incident.get("id") or f"incident-{index}",
            "type": incident.get("type") or "traffic",
            "severity": _normalize_severity(incident.get("severity")),
            "severity_color": _severity_color(incident.get("severity")),
            "description": incident.get("description") or "Traffic incident",
            "latitude": latitude,
            "longitude": longitude,
            "start_time": incident.get("start_time"),
            "end_time": incident.get("end_time"),
            "district_id": district_id,
            "response_time": incident.get("response_time") or incident.get("responseTime"),
        })

    for manual_incident in _manual_incidents_store.get(district_id, []):
        normalized.append({
            "id": manual_incident.get("id") or f"manual-{uuid.uuid4().hex[:8]}",
            "type": manual_incident.get("type") or "traffic",
            "severity": _normalize_severity(manual_incident.get("severity")),
            "severity_color": _severity_color(manual_incident.get("severity")),
            "description": manual_incident.get("description") or "Manual incident",
            "latitude": manual_incident.get("latitude"),
            "longitude": manual_incident.get("longitude"),
            "start_time": manual_incident.get("start_time") or datetime.now(UTC).isoformat(),
            "end_time": manual_incident.get("end_time"),
            "district_id": district_id,
            "response_time": manual_incident.get("response_time"),
        })

    return normalized


def _feedback_rows_for_district(
    session: Session,
    district_id: Optional[str],
    *,
    on_date: Optional[date] = None,
    zone_name: Optional[str] = None,
):
    query = session.query(MLFeedback)
    if district_id:
        query = query.filter(MLFeedback.district_id == district_id)
    if on_date is not None:
        query = query.filter(func.date(MLFeedback.created_at) == on_date.isoformat())
    if zone_name:
        query = query.filter(MLFeedback.zone == zone_name)
    return query


def _build_district_summary(incidents: list[dict], district_id: Optional[str] = None) -> dict:
    today = datetime.now(UTC).date()
    today_count = 0
    response_times: list[float] = []

    for incident in incidents:
        start_time = incident.get("start_time")
        if start_time:
            try:
                parsed = datetime.fromisoformat(str(start_time).replace("Z", "+00:00"))
                if parsed.date() == today:
                    today_count += 1
            except ValueError:
                today_count += 1
        else:
            today_count += 1

        response_time = incident.get("response_time")
        if isinstance(response_time, (int, float)):
            response_times.append(float(response_time))

    session = get_session()
    try:
        feedback_rows = _feedback_rows_for_district(session, district_id, on_date=today).all()
    finally:
        session.close()

    if feedback_rows:
        response_times = [float(row.response_time_minutes) for row in feedback_rows if row.response_time_minutes]

    avg_response_time = round(sum(response_times) / len(response_times), 2) if response_times else 0.0
    return {
        "total_incidents_today": today_count,
        "avg_response_time": avg_response_time,
    }


def _format_police_timestamp(value: object) -> str:
    if not value:
        return "-"
    if isinstance(value, datetime):
        return value.astimezone(UTC).strftime("%I:%M %p")
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed.astimezone(UTC).strftime("%I:%M %p")
    except ValueError:
        return str(value)


def _incident_sort_value(incident: dict) -> datetime:
    raw_value = incident.get("start_time")
    if isinstance(raw_value, datetime):
        return raw_value if raw_value.tzinfo else raw_value.replace(tzinfo=UTC)
    if raw_value:
        try:
            parsed = datetime.fromisoformat(str(raw_value).replace("Z", "+00:00"))
            return parsed if parsed.tzinfo else parsed.replace(tzinfo=UTC)
        except ValueError:
            return datetime.min.replace(tzinfo=UTC)
    return datetime.min.replace(tzinfo=UTC)


def _build_incidents_feed(district_id: str, incidents: list[dict], assignments: list[PoliceDispatchAssignment]) -> list[dict]:
    assignment_map = {assignment.incident_id: assignment for assignment in assignments}

    feed_items: list[dict] = []
    for incident in incidents:
        incident_id = incident.get("id")
        assignment = assignment_map.get(incident_id)
        feed_items.append({
            "incident_id": incident_id,
            "type": incident.get("type") or "traffic",
            "location_name": incident.get("description") or "Unknown location",
            "severity": incident.get("severity") or "unknown",
            "time_reported": incident.get("start_time") or "-",
            "assigned_unit": assignment.unit_id if assignment else None,
            "is_assigned": assignment is not None,
        })

    feed_items.sort(
        key=lambda item: _incident_sort_value({"start_time": item.get("time_reported")}),
        reverse=True,
    )
    return feed_items


def _serialize_recent_incidents(district_id: str, incidents: list[dict], assignments: list[PoliceDispatchAssignment]) -> list[dict]:
    """Return recent incidents with stable fallback values and assignment info."""
    assignment_map = {assignment.incident_id: assignment for assignment in assignments}
    now = datetime.now(UTC)

    normalized_items: list[dict] = []
    for index, incident in enumerate(incidents):
        incident_id = str(incident.get("id") or f"incident-{index + 1}")
        assignment = assignment_map.get(incident_id)
        timestamp = _parse_police_datetime(incident.get("start_time")) or now
        incident_type = str(incident.get("type") or "Traffic Incident").strip() or "Traffic Incident"
        location = str(
            incident.get("description")
            or incident.get("location")
            or incident.get("zone")
            or _infer_zone_name(district_id, incident)
            or "District patrol area"
        ).strip() or "District patrol area"
        severity = str(incident.get("severity") or "medium").strip().lower() or "medium"

        normalized_items.append({
            "id": incident_id,
            "location": location,
            "type": incident_type,
            "timestamp": timestamp.isoformat(),
            "assigned_patrol_unit": assignment.unit_id if assignment else None,
            "assigned_unit": assignment.unit_id if assignment else "Unassigned",
            "is_assigned": assignment is not None,
            "severity": severity,
            "lat": incident.get("latitude"),
            "lng": incident.get("longitude"),
        })

    normalized_items.sort(key=lambda item: _parse_police_datetime(item.get("timestamp")) or now, reverse=True)

    recent_items = [
        item for item in normalized_items
        if (now - (_parse_police_datetime(item.get("timestamp")) or now)) <= timedelta(hours=12)
    ]
    if recent_items:
        return recent_items[:12]
    return normalized_items[:12]


def _build_heatmap_points(incidents: list[dict]) -> list[dict]:
    severity_intensity = {
        "low": 0.35,
        "moderate": 0.65,
        "medium": 0.65,
        "high": 1.0,
        "unknown": 0.45,
    }

    points: list[dict] = []
    for incident in incidents:
        lat = incident.get("latitude")
        lng = incident.get("longitude")
        if lat is None or lng is None:
            continue
        try:
            severity = str(incident.get("severity") or "unknown").lower()
            intensity = severity_intensity.get(severity, 0.45)
            if intensity >= 0.8:
                intensity_level = "HIGH"
                marker_color = "#ef4444"
            elif intensity >= 0.5:
                intensity_level = "MEDIUM"
                marker_color = "#f59e0b"
            else:
                intensity_level = "LOW"
                marker_color = "#3b82f6"

            points.append({
                "incident_id": incident.get("id"),
                "location_name": incident.get("description") or "Unknown location",
                "lat": float(lat),
                "lng": float(lng),
                "intensity": float(intensity),
                "intensity_level": intensity_level,
                "marker_color": marker_color,
            })
        except (TypeError, ValueError):
            continue

    return points


def _parse_police_datetime(value: object) -> Optional[datetime]:
    if isinstance(value, datetime):
        return value if value.tzinfo else value.replace(tzinfo=UTC)
    if not value:
        return None
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed if parsed.tzinfo else parsed.replace(tzinfo=UTC)
    except ValueError:
        return None


def _incident_response_minutes(incident: dict) -> float:
    response_value = incident.get("response_time")
    if isinstance(response_value, (int, float)) and float(response_value) > 0:
        return float(response_value)

    severity_default = {
        "low": 6.0,
        "moderate": 8.0,
        "medium": 8.0,
        "high": 11.0,
        "unknown": 9.0,
    }
    severity = str(incident.get("severity") or "unknown").lower()
    return severity_default.get(severity, 9.0)


def _infer_zone_name(district_id: str, incident: dict) -> str:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    lat = incident.get("latitude")
    lon = incident.get("longitude")

    if isinstance(lat, (int, float)) and isinstance(lon, (int, float)):
        delta_lat = float(lat) - float(district["lat"])
        delta_lon = float(lon) - float(district["lon"])
        core_radius = 0.004

        if abs(delta_lat) <= core_radius and abs(delta_lon) <= core_radius:
            return "Central Zone"
        if abs(delta_lat) >= abs(delta_lon):
            return "North Zone" if delta_lat >= 0 else "South Zone"
        return "East Zone" if delta_lon >= 0 else "West Zone"

    description = str(incident.get("description") or "").strip().lower()
    if "north" in description:
        return "North Zone"
    if "south" in description:
        return "South Zone"
    if "east" in description:
        return "East Zone"
    if "west" in description:
        return "West Zone"
    return "Central Zone"


def _extract_affected_roads(incident: dict, zone_name: str) -> list[str]:
    """Extract and sanitize road-impact data without exposing incident details."""
    raw_roads = incident.get("affected_roads")
    if isinstance(raw_roads, list):
        cleaned = []
        for road in raw_roads:
            text_value = str(road or "").strip()
            if text_value:
                cleaned.append(text_value)
        if cleaned:
            return list(dict.fromkeys(cleaned))[:6]

    description = str(incident.get("description") or "")
    road_pattern = re.compile(
        r"([A-Za-z0-9 .'-]{2,80}\\b(?:Road|Rd|Street|St|Avenue|Ave|Highway|Hwy|Expressway|Expwy|Boulevard|Blvd|Lane|Ln|Marg|Flyover))",
        re.IGNORECASE,
    )
    extracted = [" ".join(match.split()) for match in road_pattern.findall(description)]
    deduped = []
    for road_name in extracted:
        if road_name and road_name not in deduped:
            deduped.append(road_name)

    if deduped:
        return deduped[:6]

    return [f"{zone_name} arterial corridors"]


def _shared_alert_expiry(timestamp: datetime, severity: str) -> datetime:
    durations = {
        "critical": 120,
        "high": 90,
    }
    minutes = durations.get(str(severity or "").lower(), 60)
    return timestamp + timedelta(minutes=minutes)


def _create_shared_alert_for_dispatch(session: Session, district_id: str, incident: dict) -> Optional[SharedAlert]:
    """Persist sanitized logistics alert for high-impact dispatches only."""
    severity = _normalize_severity(incident.get("severity"))
    if severity not in {"critical", "high"}:
        return None

    now = datetime.now(UTC)
    zone_name = _infer_zone_name(district_id, incident)
    alert = SharedAlert(
        alert_id=str(uuid.uuid4()),
        zone=zone_name,
        severity=severity,
        timestamp=now,
        affected_roads=_extract_affected_roads(incident, zone_name),
        expires_at=_shared_alert_expiry(now, severity),
    )
    session.add(alert)
    return alert


def _build_response_time_by_zone(
    district_id: str,
    incidents: list[dict],
    assignments: list[PoliceDispatchAssignment],
    target_threshold_minutes: float = 8.0,
) -> list[dict]:
    assignment_map = {assignment.incident_id: assignment for assignment in assignments}
    today = datetime.now(UTC).date()
    district_unit_prefix = district_id.upper().replace("_", "-")

    session = get_session()
    try:
        feedback_rows = _feedback_rows_for_district(session, district_id, on_date=today).all()
        available_units = [
            row.officer_id
            for row in (
                session.query(OfficerDispatchStatus)
                .filter(OfficerDispatchStatus.district_id == district_id)
                .order_by(OfficerDispatchStatus.officer_id.asc())
                .all()
            )
        ]
    finally:
        session.close()

    if feedback_rows:
        feedback_zone_stats: dict[str, dict] = {}
        for row in feedback_rows:
            zone_name = (row.zone or "Central Zone").strip() or "Central Zone"
            bucket = feedback_zone_stats.setdefault(
                zone_name,
                {
                    "zone": zone_name,
                    "response_values": [],
                    "total_incidents": 0,
                },
            )
            bucket["response_values"].append(float(row.response_time_minutes))
            bucket["total_incidents"] += 1

        response_rows = []
        zone_names = sorted(feedback_zone_stats.keys())
        for zone_index, zone_name in enumerate(zone_names):
            values = sorted(feedback_zone_stats[zone_name]["response_values"])
            avg_response = round(sum(values) / len(values), 2)
            fallback_unit = available_units[zone_index % len(available_units)] if available_units else "N/A"
            slowest_unit = available_units[(zone_index + 1) % len(available_units)] if available_units else "N/A"
            response_rows.append({
                "zone": zone_name,
                "avg_response_time": avg_response,
                "fastest_unit": fallback_unit,
                "slowest_unit": slowest_unit,
                "total_incidents": feedback_zone_stats[zone_name]["total_incidents"],
                "exceeds_target": avg_response > float(target_threshold_minutes),
            })

        response_rows.sort(key=lambda item: item["avg_response_time"], reverse=True)
        return response_rows

    def aggregate_rows(source_incidents: list[dict], today_only: bool) -> dict[str, dict]:
        zone_stats: dict[str, dict] = {}
        for incident in source_incidents:
            start_time = _parse_police_datetime(incident.get("start_time"))
            if today_only and start_time and start_time.date() != today:
                continue

            response_minutes = incident.get("response_time")
            if not isinstance(response_minutes, (int, float)) or float(response_minutes) <= 0:
                continue

            zone_name = _infer_zone_name(district_id, incident)
            bucket = zone_stats.setdefault(
                zone_name,
                {
                    "zone": zone_name,
                    "total_incidents": 0,
                    "response_sum": 0.0,
                    "response_count": 0,
                    "unit_times": {},
                },
            )

            bucket["total_incidents"] += 1
            response_minutes = round(float(response_minutes), 2)
            bucket["response_sum"] += response_minutes
            bucket["response_count"] += 1

            assignment = assignment_map.get(incident.get("id"))
            unit_id = assignment.unit_id if assignment else None
            unit_bucket = bucket["unit_times"].setdefault(unit_id, {"sum": 0.0, "count": 0})
            unit_bucket["sum"] += response_minutes
            unit_bucket["count"] += 1

        return zone_stats

    zone_stats = aggregate_rows(incidents, today_only=True)
    if not zone_stats:
        zone_stats = aggregate_rows(incidents, today_only=False)

    response_rows: list[dict] = []
    for zone_name, stats in zone_stats.items():
        avg_response = round(stats["response_sum"] / max(stats["response_count"], 1), 2)

        unit_averages = []
        for unit_id, values in stats["unit_times"].items():
            if values["count"] <= 0:
                continue
            unit_averages.append((unit_id, values["sum"] / values["count"]))

        unit_averages.sort(key=lambda item: item[1])
        units_in_zone = [unit_id for unit_id in available_units if unit_id.startswith(district_unit_prefix)]
        fallback_unit = units_in_zone[0] if units_in_zone else "N/A"
        fastest_unit = unit_averages[0][0] if unit_averages and unit_averages[0][0] else fallback_unit
        slowest_unit = unit_averages[-1][0] if unit_averages and unit_averages[-1][0] else fallback_unit

        response_rows.append({
            "zone": zone_name,
            "avg_response_time": avg_response,
            "fastest_unit": fastest_unit,
            "slowest_unit": slowest_unit,
            "total_incidents": stats["total_incidents"],
            "exceeds_target": avg_response > float(target_threshold_minutes),
        })

    response_rows.sort(key=lambda item: item["avg_response_time"], reverse=True)
    return response_rows


def _normalize_patrol_status(raw_status: object) -> str:
    status_value = str(raw_status or "").strip().lower()
    if status_value in {"responding", "response", "enroute", "en route"}:
        return "Responding"
    if status_value in {"idle", "available", "standby"}:
        return "Idle"
    if status_value == "active":
        return "Active"
    return "Active"


def _get_dispatch_assignments(district_id: str) -> list[PoliceDispatchAssignment]:
    session = get_session()
    try:
        return (
            session.query(PoliceDispatchAssignment)
            .filter(PoliceDispatchAssignment.district_id == district_id)
            .order_by(PoliceDispatchAssignment.assigned_at.desc())
            .all()
        )
    finally:
        session.close()


def _build_patrol_units(
    district_id: str,
    incidents: list[dict],
    supervisor_name: str,
    assignments: Optional[list[PoliceDispatchAssignment]] = None,
) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    status_by_unit: dict[str, str] = {}
    status_by_officer: dict[str, str] = {}
    session = get_session()
    try:
        police_users = (
            session.query(User)
            .filter(
                User.department == "police",
                User.is_active == True,  # noqa: E712
                User.district_id == district_id,
            )
            .order_by(User.id.asc())
            .all()
        )
        named_police_users = [
            {
                "user": user,
                "officer_name": _resolve_officer_name(user),
            }
            for user in police_users
            if _resolve_officer_name(user)
        ]
        status_rows = (
            session.query(OfficerDispatchStatus)
            .filter(OfficerDispatchStatus.district_id == district_id)
            .all()
        )
        status_by_unit = {row.officer_id: str(row.status or '').lower() for row in status_rows}
        status_by_officer = {row.officer_id: str(row.status or '').lower() for row in status_rows}
    finally:
        session.close()

    assignments = assignments or []
    assignment_by_unit = {assignment.unit_id: assignment for assignment in assignments}

    available_officers = [
        item for item in named_police_users
        if status_by_officer.get(item["user"].username, "available") == "available"
    ]

    # Ensure displayed officer labels are unique even when multiple accounts
    # share the same full_name (for example seeded test users named "Officer Raj").
    display_name_counts: dict[str, int] = {}
    for item in available_officers:
        base_name = str(item.get("officer_name") or "Unassigned").strip() or "Unassigned"
        display_name_counts[base_name] = display_name_counts.get(base_name, 0) + 1

    for item in available_officers:
        base_name = str(item.get("officer_name") or "Unassigned").strip() or "Unassigned"
        username = str(getattr(item.get("user"), "username", "") or "").strip()
        if display_name_counts.get(base_name, 0) > 1 and username:
            item["display_name"] = f"{base_name} ({username})"
        else:
            item["display_name"] = base_name

    unit_count = min(max(4, len(available_officers)), 15)
    patrol_units: list[dict] = []

    patrol_name_cycle = [
        "Patrol Alpha",
        "Patrol Bravo",
        "Patrol Charlie",
        "Patrol Delta",
        "Patrol Echo",
        "Patrol Foxtrot",
        "Patrol Golf",
        "Patrol Hotel",
        "Patrol India",
        "Patrol Juliet",
    ]

    for index in range(unit_count):
        incident = incidents[index % len(incidents)] if incidents else None
        officer_assignment = available_officers[index] if index < len(available_officers) else None
        officer = officer_assignment["user"] if officer_assignment else None
        officer_name = officer_assignment["display_name"] if officer_assignment else "Unassigned"
        
        # Use consistent unit ID format: DISTRICT-X-UXX
        district_prefix = district_id.upper().replace("_", "-")
        patrol_id = f"{district_prefix}-U{index + 1:02d}"
        unit_id = patrol_id
        assignment = assignment_by_unit.get(unit_id)
        assignment_override = None
        with _patrol_user_assignments_lock:
            assignment_override = (_patrol_user_assignments.get(district_id) or {}).get(unit_id)
        persisted_status = status_by_unit.get(unit_id)

        if officer_assignment is None:
            # Unassigned units are always available - cannot respond without officer
            status = "available"
        elif persisted_status in {"responding", "enroute"}:
            status = "responding"
        elif persisted_status == "busy":
            status = "busy"
        elif persisted_status == "offline":
            status = "offline"
        elif assignment_override:
            status = "responding"
        elif assignment:
            status = "responding"

        else:
            status = "available"

        location = district["name"]
        last_updated = datetime.now(UTC)

        if assignment:
            incident_match = next((item for item in incidents if item.get("id") == assignment.incident_id), None)
            if incident_match:
                location = incident_match.get("description") or incident_match.get("type") or district["name"]
            else:
                location = f"Assigned to incident {assignment.incident_id}"
            last_updated = assignment.assigned_at or last_updated
        elif assignment_override:
            location = district["name"]
            assigned_at = assignment_override.get("assigned_at")
            if isinstance(assigned_at, datetime):
                last_updated = assigned_at
        elif incident:
            location = incident.get("description") or incident.get("type") or district["name"]
            last_updated = incident.get("start_time") or last_updated

        latitude = district["lat"]
        longitude = district["lon"]
        if assignment:
            incident_match = next((item for item in incidents if item.get("id") == assignment.incident_id), None)
            if incident_match and incident_match.get("latitude") is not None and incident_match.get("longitude") is not None:
                latitude = float(incident_match["latitude"])
                longitude = float(incident_match["longitude"])
        elif incident and incident.get("latitude") is not None and incident.get("longitude") is not None:
            latitude = float(incident["latitude"])
            longitude = float(incident["longitude"])

        fallback_area = clean_location(location)
        current_area = tomtom_reverse_geocode_area(latitude, longitude, fallback=fallback_area)
        unit_name = patrol_name_cycle[index % len(patrol_name_cycle)] if officer_assignment else f"Reserve Unit {index + 1:02d}"

        patrol_units.append({
            "patrol_id": patrol_id,
            "unit_id": unit_id,
            "unit_name": unit_name,
            "officer_username": officer.username if officer else None,
            "officer_name": officer_name,
            "officer": officer_name,
            "status": status,
            "current_area": current_area,
            "current_location": current_area,
            "last_updated": _format_police_timestamp(last_updated),
            "updated_at": last_updated.isoformat() if isinstance(last_updated, datetime) else str(last_updated),
            "district_id": district_id,
            "assigned_incident_id": assignment.incident_id if assignment else None,
            "assigned_user_id": assignment_override.get("user_id") if isinstance(assignment_override, dict) else None,
            "latitude": latitude,
            "longitude": longitude,
        })

    return patrol_units


def _build_police_dashboard_context(current_user: dict, district_id: str) -> dict:
    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    incidents_feed = _build_incidents_feed(district_id, incidents, assignments)
    district_summary = _build_district_summary(incidents, district_id)
    district_info = DISTRICT_LOCATIONS.get(district_id, {"name": district_id or "Unknown District"})
    supervisor_name = current_user.get("username", "Unknown Supervisor")
    patrol_units = _build_patrol_units(district_id, incidents, supervisor_name, assignments)
    response_zones = _build_response_time_by_zone(district_id, incidents, assignments)

    assigned_incident_ids = {assignment.incident_id for assignment in assignments}
    unassigned_incidents = [incident for incident in incidents if incident.get("id") not in assigned_incident_ids]
    available_patrol_units = [unit for unit in patrol_units if unit["status"] == "available"]
    responding_units = [unit for unit in patrol_units if unit["status"] == "responding"]
    deployed_units = [
        unit for unit in patrol_units
        if str(unit.get("status") or "").lower() in {"responding", "busy"}
    ]
    avg_response_time_value = district_summary.get("avg_response_time", 0.0)
    avg_response_time = avg_response_time_value if avg_response_time_value else "N/A"

    return {
        "district": _format_district_label(district_id),
        "district_info": district_info,
        "district_name": district_info.get("name", district_id or "Unknown District"),
        "supervisor_name": supervisor_name,
        "shift_time": datetime.now(UTC).strftime("%I:%M %p UTC"),
        "total_active_incidents": len(incidents),
        "units_deployed": len(deployed_units),
        "units_available": len(available_patrol_units),
        "total": len(patrol_units),
        "available": len(available_patrol_units),
        "responding": len(responding_units),
        "avg_response_time": avg_response_time,
        "patrol_units": patrol_units,
        "available_patrol_units": available_patrol_units,
        "unassigned_incidents": unassigned_incidents,
        "incidents_feed": incidents_feed,
        "incidents": incidents,
        "district_summary": district_summary,
        "zone_count": len(response_zones),
        "response_target_minutes": 8,
    }


def _district_prediction_candidates(district_id: str, incidents: list[dict]) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    candidates: list[dict] = []

    for incident in incidents:
        if incident.get("latitude") is not None and incident.get("longitude") is not None:
            candidates.append({
                "location": incident.get("description") or "Incident hotspot",
                "latitude": incident["latitude"],
                "longitude": incident["longitude"],
                "base_severity": incident.get("severity", "unknown"),
                "source": "live_incident",
            })

    if not candidates:
        offsets = [
            (0.0000, 0.0000),
            (0.0120, 0.0080),
            (-0.0100, 0.0110),
            (0.0090, -0.0090),
            (-0.0080, -0.0100),
        ]
        for index, (lat_offset, lon_offset) in enumerate(offsets):
            candidates.append({
                "location": f"{district['name']} sector {index + 1}",
                "latitude": district["lat"] + lat_offset,
                "longitude": district["lon"] + lon_offset,
                "base_severity": "unknown",
                "source": "district_grid",
            })

    return candidates[:10]


def _clamp(value: float, minimum: float, maximum: float) -> float:
    return max(minimum, min(maximum, float(value)))


def _candidate_zone_name(district_id: str, candidate: dict) -> str:
    return _infer_zone_name(
        district_id,
        {
            "latitude": candidate.get("latitude"),
            "longitude": candidate.get("longitude"),
            "description": candidate.get("location"),
        },
    )


def _candidate_historical_incident_count(district_id: str, candidate: dict, incidents: list[dict]) -> int:
    candidate_lat = candidate.get("latitude")
    candidate_lon = candidate.get("longitude")
    nearby_count = 0

    for incident in incidents:
        lat = incident.get("latitude")
        lon = incident.get("longitude")
        if not isinstance(lat, (int, float)) or not isinstance(lon, (int, float)):
            continue
        distance_m = haversine_m(candidate_lat, candidate_lon, float(lat), float(lon))
        if distance_m <= 1800:
            nearby_count += 1

    zone_name = _candidate_zone_name(district_id, candidate)
    session = get_session()
    try:
        feedback_count = _feedback_rows_for_district(session, district_id, zone_name=zone_name).count()
    finally:
        session.close()

    return int(nearby_count + feedback_count)


def _candidate_current_traffic_speed(candidate: dict, distance_km: float, historical_incident_count: int) -> float:
    severity_penalty = {
        "low": 2.0,
        "moderate": 5.0,
        "medium": 5.0,
        "high": 9.0,
        "unknown": 3.5,
    }.get(str(candidate.get("base_severity") or "unknown").lower(), 3.5)

    lat = float(candidate.get("latitude") or 0.0)
    lon = float(candidate.get("longitude") or 0.0)
    location_variation = abs(lat * 13.0 + lon * 7.0) % 6.0

    estimated_speed = 42.0 - severity_penalty - (historical_incident_count * 2.5) - (distance_km * 1.4) - location_variation
    return round(_clamp(estimated_speed, 12.0, 55.0), 2)


def _hotspot_model_probability(feature_frame: dict) -> float:
    model_prediction = predict_congestion(feature_frame)
    if model_prediction is None:
        return 0.5

    prediction_value = float(model_prediction)
    normalized_score = (prediction_value - 0.85) / 0.55
    return _clamp(normalized_score, 0.0, 1.0)


def _build_rule_based_hotspots(district_id: str, incidents: list[dict], patrol_units: Optional[list[dict]] = None) -> list[dict]:
    """Build dynamic hotspot rankings from current incidents and patrol availability."""
    candidates = _district_prediction_candidates(district_id, incidents)
    if not candidates:
        return []

    now = datetime.now(UTC)
    rng = random.SystemRandom()
    patrol_units = patrol_units or []
    available_patrol_units = [
        unit for unit in patrol_units
        if str(unit.get("status") or "").lower() == "available"
        and isinstance(unit.get("latitude"), (int, float))
        and isinstance(unit.get("longitude"), (int, float))
    ]

    risk_band_weights = {
        "low": 6,
        "moderate": 14,
        "medium": 14,
        "high": 24,
        "critical": 30,
        "unknown": 8,
    }
    incident_type_keywords = {
        "accident": "Accident",
        "crash": "Accident",
        "collision": "Accident",
        "crowd": "Crowd",
        "market": "Crowd",
        "festival": "Crowd",
        "traffic": "Congestion",
        "jam": "Congestion",
        "slow": "Congestion",
        "closure": "Congestion",
    }

    hotspots: list[dict] = []
    for index, candidate in enumerate(candidates):
        nearby_incidents: list[dict] = []
        candidate_lat = float(candidate.get("latitude") or 0.0)
        candidate_lon = float(candidate.get("longitude") or 0.0)
        for incident in incidents:
            lat = incident.get("latitude")
            lon = incident.get("longitude")
            if not isinstance(lat, (int, float)) or not isinstance(lon, (int, float)):
                continue
            distance_m = haversine_m(candidate_lat, candidate_lon, float(lat), float(lon))
            if distance_m <= 2500:
                nearby_incidents.append(incident)

        severity_score = sum(risk_band_weights.get(str(item.get("severity") or "unknown").lower(), 8) for item in nearby_incidents)
        density_score = min(len(nearby_incidents) * 10, 24)
        time_score = 10 if now.hour in {7, 8, 9, 17, 18, 19, 20} else 5
        location_variation = int((abs(candidate_lat * 1000.0) + abs(candidate_lon * 1000.0)) % 11)

        base_score = 40 + min(severity_score // 2, 22) + density_score + time_score + location_variation
        traffic_score = round(_clamp(base_score + rng.randint(-8, 10), 40.0, 95.0), 1)
        if traffic_score >= 80:
            risk_level = "HIGH"
            intensity = 1.0
        elif traffic_score >= 50:
            risk_level = "MEDIUM"
            intensity = 0.65
        else:
            risk_level = "LOW"
            intensity = 0.35

        inferred_type = None
        text_for_type = " ".join([
            str(candidate.get("location") or ""),
            " ".join(str(item.get("type") or "") for item in nearby_incidents),
        ]).lower()
        for keyword, label in incident_type_keywords.items():
            if keyword in text_for_type:
                inferred_type = label
                break
        if inferred_type is None:
            inferred_type = rng.choice(["Congestion", "Accident", "Crowd"])

        confidence_floor = max(60, min(92, int(traffic_score) - 6))
        confidence = round(float(rng.randint(confidence_floor, 95)), 1)
        suggested_patrol = None
        if available_patrol_units:
            suggested_patrol = min(
                available_patrol_units,
                key=lambda unit: haversine_m(
                    candidate_lat,
                    candidate_lon,
                    float(unit.get("latitude") or 0.0),
                    float(unit.get("longitude") or 0.0),
                ),
            )

        hotspots.append({
            "rank": index + 1,
            "location": candidate.get("location") or "Unknown location",
            "latitude": candidate_lat,
            "longitude": candidate_lon,
            "risk_score": round(traffic_score, 1),
            "traffic_score": round(traffic_score, 1),
            "risk_level": risk_level,
            "intensity": intensity,
            "incident_type": inferred_type,
            "confidence": round(confidence, 1),
            "incident_count": len(nearby_incidents),
            "suggested_patrol_id": suggested_patrol.get("patrol_id") if suggested_patrol else None,
            "suggested_patrol_name": suggested_patrol.get("unit_name") if suggested_patrol else None,
            "suggested_officer_name": suggested_patrol.get("officer_name") if suggested_patrol else None,
            "can_assign_patrol": bool(suggested_patrol),
            "updated_at": now.isoformat(),
        })

    hotspots.sort(key=lambda item: item["traffic_score"], reverse=True)
    for idx, hotspot in enumerate(hotspots[:5], start=1):
        hotspot["rank"] = idx
    return hotspots[:5]


def _predict_police_hotspots(district_id: str, incidents: list[dict]) -> list[dict]:
    """Build rule-based hotspot predictions from current incident density and patrol availability."""
    return _build_rule_based_hotspots(district_id, incidents)


def _hotspots_to_geojson(hotspots: list[dict], district_id: str) -> PredictedIncidentResponse:
    features = []
    for hotspot in hotspots:
        score = float(hotspot.get("traffic_score", hotspot.get("likelihood_score", 0)))
        risk_level = str(hotspot.get("risk_level") or "low").lower()
        if risk_level not in {"high", "medium", "low"}:
            risk_level = "high" if score >= 75 else "medium" if score >= 40 else "low"

        district_coords = DISTRICT_LOCATIONS.get(district_id, {})
        center_lat = district_coords.get("lat", 28.6139)
        center_lon = district_coords.get("lon", 77.2090)
        radius = 0.05 * (100 - score) / 100
        lat = float(hotspot.get("latitude", center_lat))
        lon = float(hotspot.get("longitude", center_lon))

        import math
        polygon_coords = []
        for i in range(8):
            angle = (i / 8) * 2 * math.pi
            poly_lat = lat + radius * math.cos(angle)
            poly_lon = lon + radius * math.sin(angle)
            polygon_coords.append([poly_lon, poly_lat])
        polygon_coords.append(polygon_coords[0])

        features.append({
            "type": "Feature",
            "geometry": {
                "type": "Polygon",
                "coordinates": [polygon_coords],
            },
            "properties": {
                "zone_name": hotspot.get("location", f"Zone_{lat:.4f}_{lon:.4f}"),
                "risk_level": risk_level,
                "prediction_score": round(score, 2),
                "incident_count": hotspot.get("incident_count", 0),
                "severity": hotspot.get("risk_level", "low"),
                "suggested_patrol_id": hotspot.get("suggested_patrol_id"),
                "suggested_patrol_name": hotspot.get("suggested_patrol_name"),
                "can_assign_patrol": hotspot.get("can_assign_patrol", False),
                "traffic_score": score,
                "incident_type": hotspot.get("incident_type", "Congestion"),
                "confidence": hotspot.get("confidence", 0),
            },
        })

    return PredictedIncidentResponse(type="FeatureCollection", features=features)


def _safe_ppt_text(value: object) -> str:
    if value is None:
        return ""
    return str(value)


def _format_ppt_timestamp(value: Optional[str]) -> str:
    if not value:
        return "-"
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed.strftime("%Y-%m-%d %H:%M")
    except ValueError:
        return str(value)


def generate_shift_pptx(
    district_id: str,
    officer_name: str,
    incidents: list[dict],
    district_summary: dict,
    ml_predictions: list[dict],
) -> io.BytesIO:
    """Generate a police shift PPTX report and return an in-memory stream."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE

    district_info = DISTRICT_LOCATIONS.get(district_id, {"name": district_id or "Unknown District"})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_body_layout = prs.slide_layouts[5]

    # Slide 1 - Shift summary
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Police Shift Summary"
    slide1.placeholders[1].text = (
        f"District: {district_info.get('name', district_id or 'Unknown District')}\n"
        f"Date: {datetime.now(UTC).strftime('%Y-%m-%d')}\n"
        f"Officer: {officer_name}\n"
        f"Total Incidents: {len(incidents)}\n"
        f"Avg Response Time: {district_summary.get('avg_response_time', 'N/A')}"
    )

    # Slide 2 - Incidents table
    slide2 = prs.slides.add_slide(title_body_layout)
    slide2.shapes.title.text = "Incident Breakdown"
    rows = max(len(incidents), 1) + 1
    cols = 4
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6))
    table = table_shape.table
    headers = ["ID", "Description", "Severity", "Start Time"]
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, incident in enumerate(incidents, start=1):
        table.cell(row_index, 0).text = _safe_ppt_text(incident.get("id"))
        table.cell(row_index, 1).text = _safe_ppt_text(incident.get("description"))
        table.cell(row_index, 2).text = _safe_ppt_text(incident.get("severity"))
        table.cell(row_index, 3).text = _format_ppt_timestamp(incident.get("start_time"))

    # Slide 3 - ML highlights
    slide3 = prs.slides.add_slide(title_body_layout)
    slide3.shapes.title.text = "ML Prediction Highlights"
    text_box = slide3.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
    text_frame = text_box.text_frame
    if ml_predictions:
        for index, prediction in enumerate(ml_predictions, start=1):
            paragraph = text_frame.paragraphs[0] if index == 1 else text_frame.add_paragraph()
            paragraph.text = (
                f"{index}. {prediction['location']} | "
                f"Likelihood: {prediction['likelihood_score']}% | "
                f"Confidence: {prediction['confidence']}% | "
                f"Type: {prediction['predicted_type']}"
            )
            for run in paragraph.runs:
                run.font.size = Pt(18)
    else:
        text_frame.text = "No ML predictions available for this district."

    # Slide 4 - Map placeholder
    slide4 = prs.slides.add_slide(title_body_layout)
    slide4.shapes.title.text = "Map Screenshot Placeholder"
    placeholder = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(1.6), Inches(11.2), Inches(4.6)
    )
    placeholder.text_frame.text = (
        "Insert map screenshot here before final delivery.\n\n"
        f"District: {district_info.get('name', district_id or 'Unknown District')}\n"
        f"Incidents on map: {len(incidents)}"
    )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ============================================================================
# VALIDATION MODELS (Pydantic V2 Compatible)
# ============================================================================

class ValidatedCoordinates(BaseModel):
    """Validated coordinate pair."""
    lat: float
    lon: float
    
    @field_validator('lat')
    @classmethod
    def validate_latitude(cls, v: float) -> float:
        """Validate latitude is within -90 to 90."""
        if v is None:
            raise ValueError('Latitude is required')
        if not -90 <= v <= 90:
            raise ValueError(f'Invalid latitude: {v}. Must be between -90 and 90')
        return v
    
    @field_validator('lon')
    @classmethod
    def validate_longitude(cls, v: float) -> float:
        """Validate longitude is within -180 to 180."""
        if v is None:
            raise ValueError('Longitude is required')
        if not -180 <= v <= 180:
            raise ValueError(f'Invalid longitude: {v}. Must be between -180 and 180')
        return v


class RouteAnalysisRequest(BaseModel):
    """Request model for route analysis with validation."""
    origin: Union[str, dict] = Field(..., description="Origin as place name or {lat, lon}")
    destination: Union[str, dict] = Field(..., description="Destination as place name or {lat, lon}")
    maxAlternatives: int = Field(3, ge=1, le=5, description="Maximum route alternatives")
    alpha: float = Field(1.0, ge=0, le=10, description="Weight for travel time in cost calculation")
    beta: float = Field(0.5, ge=0, le=10, description="Weight for delay in cost calculation")
    gamma: float = Field(0.001, ge=0, le=1, description="Weight for distance in cost calculation")
    avoid_tolls: bool = Field(False, description="Avoid toll roads")
    avoid_ferries: bool = Field(False, description="Avoid ferries")
    avoid_highways: bool = Field(False, description="Avoid highways")
    
    @field_validator('maxAlternatives')
    @classmethod
    def validate_max_alternatives(cls, v: int) -> int:
        if v < 1 or v > 5:
            raise ValueError('maxAlternatives must be between 1 and 5')
        return v
    
    @field_validator('alpha', 'beta', 'gamma')
    @classmethod
    def validate_weights(cls, v: float, info) -> float:
        if v < 0:
            raise ValueError(f'{info.field_name} cannot be negative')
        return v


class UserCreate(BaseModel):
    """User registration model."""
    email: EmailStr
    username: str = Field(..., min_length=3, max_length=50)
    password: str = Field(..., min_length=6)
    full_name: Optional[str] = None
    role: str = Field("user", description="user | police_supervisor | logistics_manager")


class UserLogin(BaseModel):
    """User login model."""
    username: str
    password: str


class SavedRouteCreate(BaseModel):
    """Create saved route model."""
    route_name: str
    origin: Union[str, dict]
    destination: Union[str, dict]
    route_preferences: Optional[dict] = None


class RouteRatingCreate(BaseModel):
    """Create route rating model."""
    route_id: str
    rating: int = Field(..., ge=1, le=5)
    review: Optional[str] = None


class ShareRouteRequest(BaseModel):
    """Share route request model."""
    route_id: str
    route_index: Optional[int] = None


class UserUpdate(BaseModel):
    """User update model for admin editing."""
    username: Optional[str] = None
    email: Optional[EmailStr] = None
    full_name: Optional[str] = None
    is_active: Optional[bool] = None
    is_admin: Optional[bool] = None
    password: Optional[str] = Field(None, min_length=8)


class DispatchRequest(BaseModel):
    """Dispatch assignment request."""
    incident_id: str = Field(..., min_length=1)
    officer_id: Optional[str] = Field(None, min_length=1)
    unit_id: Optional[str] = Field(None, min_length=1)


class PatrolSelectionDispatchRequest(BaseModel):
    """Dispatch a selected patrol unit from the command center."""
    patrol_id: str = Field(..., min_length=1)
    incident_id: Optional[str] = Field(None, min_length=1)
    officer_id: Optional[str] = Field(None, min_length=1)
    unit_id: Optional[str] = Field(None, min_length=1)
    user_id: Optional[str] = Field(None, min_length=1)


class ShiftAttendanceRequest(BaseModel):
    """Mark officer attendance during shift."""
    officers: list[dict] = Field(..., description="List of {officer_username, officer_name, status}")


class ShiftEndRequest(BaseModel):
    """End shift and optionally export report."""
    notes: Optional[str] = None
    export_pptx: bool = True


class IncidentResolveRequest(BaseModel):
    """Incident resolution payload used for feedback logging and analytics."""
    incident_id: str = Field(..., min_length=1)
    incident_type: str = Field(..., min_length=1)
    zone: str = Field(..., min_length=1)
    response_time_minutes: float = Field(..., gt=0)
    severity: str = Field(..., pattern="^(critical|high|medium|low|moderate|unknown)$")
    outcome: str = Field(..., min_length=1)
    resolved_at: Optional[datetime] = None


class NewIncidentRequest(BaseModel):
    """Create a new command-center incident for live dispatch workflows."""
    incident_type: str = Field(..., min_length=1)
    severity: str = Field(..., pattern="^(critical|high|medium|low|moderate|unknown)$")
    description: str = Field(..., min_length=1)
    latitude: Optional[float] = None
    longitude: Optional[float] = None
    zone: Optional[str] = None


class AlertData(BaseModel):
    """Real-time alert for police supervisor."""
    alert_id: str
    severity: str  # critical, high, medium, low
    message: str
    timestamp: str  # ISO format
    district_id: Optional[str] = None
    related_incident_id: Optional[str] = None


class PredictedHotspotPolygon(BaseModel):
    """GeoJSON polygon for predicted incident hotspot."""
    zone_name: str
    risk_level: str  # high, medium, low
    prediction_score: float  # 0-100
    coordinates: list[list[list[float]]]  # GeoJSON Polygon coordinates


class PredictedIncidentResponse(BaseModel):
    """Response containing predicted incident hotspots as GeoJSON."""
    type: str = "FeatureCollection"
    features: list[dict]  # GeoJSON features


class OfficerWorkloadData(BaseModel):
    """Officer incident counts and rotation status."""
    officer_username: str
    officer_name: str
    total_incidents: int
    critical_incidents: int
    high_incidents: int
    medium_incidents: int
    low_incidents: int
    needs_rotation: bool


class OfficerWorkloadResponse(BaseModel):
    """Response containing officer workload data for current shift."""
    shift_id: int
    officers: list[OfficerWorkloadData]
    officers_needing_rotation: list[str]  # Usernames of officers needing rotation


class AlertListResponse(BaseModel):
    """Response for unread alerts list."""
    alerts: list[AlertData]
    unread_count: int

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def decode_polyline(encoded: str) -> list[tuple[float, float]]:
    """
    Decode Google-style polyline string to list of (lat, lon) tuples.
    Handles TomTom's coordinate formats properly.
    
    Args:
        encoded: Encoded polyline string or coordinate list
        
    Returns:
        List of (latitude, longitude) tuples
    """
    if not encoded:
        return []
    
    # Handle if encoded is already a list of coordinates
    if isinstance(encoded, list):
        coordinates = []
        for p in encoded:
            try:
                if isinstance(p, dict):
                    if "lat" in p and "lon" in p:
                        coordinates.append((float(p["lat"]), float(p["lon"])))
                    elif "latitude" in p and "longitude" in p:
                        coordinates.append((float(p["latitude"]), float(p["longitude"])))
                elif isinstance(p, (list, tuple)) and len(p) >= 2:
                    coordinates.append((float(p[0]), float(p[1])))
            except (ValueError, TypeError):
                continue
        return coordinates
    
    # Try to use the polyline library if available
    try:
        import polyline
        return polyline.decode(encoded)
    except ImportError:
        # Fallback to custom decoder if polyline not installed
        pass
    except Exception as e:
        logger.warning(f"Polyline library decoding failed: {e}")
    
    # Custom decoder implementation
    coordinates = []
    index = 0
    lat = 0
    lon = 0
    
    try:
        while index < len(encoded):
            # Decode latitude
            shift = 0
            result = 0
            while index < len(encoded):
                b = ord(encoded[index]) - 63
                index += 1
                result |= (b & 0x1f) << shift
                shift += 5
                if b < 0x20:
                    break
            dlat = ~(result >> 1) if (result & 1) else (result >> 1)
            lat += dlat
            
            # Check if we have enough characters for longitude
            if index >= len(encoded):
                break
                
            # Decode longitude
            shift = 0
            result = 0
            while index < len(encoded):
                b = ord(encoded[index]) - 63
                index += 1
                result |= (b & 0x1f) << shift
                shift += 5
                if b < 0x20:
                    break
            dlon = ~(result >> 1) if (result & 1) else (result >> 1)
            lon += dlon
            
            coordinates.append((lat / 1e5, lon / 1e5))
    except (IndexError, ValueError, TypeError) as e:
        logger.error(f"Polyline decoding error: {e}")
        return []
    
    return coordinates


def extract_route_geometry(route_json: dict) -> list[tuple[float, float]]:
    """
    Extract route geometry from TomTom route JSON.
    
    Args:
        route_json: Route object from TomTom API
        
    Returns:
        List of (lat, lon) tuples for route path
    """
    geometry = []
    
    try:
        legs = route_json.get("legs", [])
        for leg in legs:
            points = leg.get("points", [])
            for point in points:
                if isinstance(point, dict):
                    if "latitude" in point and "longitude" in point:
                        try:
                            geometry.append((float(point["latitude"]), float(point["longitude"])))
                        except (ValueError, TypeError):
                            continue
                    elif "lat" in point and "lon" in point:
                        try:
                            geometry.append((float(point["lat"]), float(point["lon"])))
                        except (ValueError, TypeError):
                            continue
        
        if not geometry:
            guidance = route_json.get("guidance", {})
            instructions = guidance.get("instructions", [])
            for instruction in instructions:
                point = instruction.get("point", {})
                if "latitude" in point and "longitude" in point:
                    try:
                        geometry.append((float(point["latitude"]), float(point["longitude"])))
                    except (ValueError, TypeError):
                        continue
    except Exception as e:
        logger.error(f"Error extracting route geometry: {e}")
        return []
    
    return geometry


def predict_congestion(features: dict) -> Optional[float]:
    """
    Predict route score using RF model.
    Lower score = better route.
    """
    if ML_MODEL is None:
        return None
    try:
        import pandas as pd
        now = datetime.now(UTC)

        travel_time_s = features.get("travel_time_s", 0)
        no_traffic_s = features.get("no_traffic_s", 0)
        if travel_time_s is None:
            travel_time_s = 0
        if no_traffic_s is None:
            no_traffic_s = 0

        congestion_ratio = (
            float(travel_time_s) / float(no_traffic_s)
            if no_traffic_s and float(no_traffic_s) > 0
            else 1.0
        )

        hour = int(features.get("hour", now.hour))
        weekday = int(features.get("weekday", now.weekday()))

        base_feature_dict = {
            "hour": hour,
            "weekday": weekday,
            "is_weekend": int(features.get("is_weekend", 1 if weekday >= 5 else 0)),
            "distance_km": float(features.get("distance_km", 0) or 0),
            "route_index": int(features.get("route_index", 0) or 0),
            "travel_time_s": float(travel_time_s),
            "no_traffic_s": float(no_traffic_s),
            "delay_s": float(features.get("delay_s", 0) or 0),
            "rolling_mean_congestion": float(features.get("rolling_mean_congestion", 1.0) or 1.0),
            "rolling_std_congestion": float(features.get("rolling_std_congestion", 0.0) or 0.0),
            "congestion_ratio": float(congestion_ratio),
        }

        model_feature_names = getattr(ML_MODEL, "feature_names_in_", None)
        if model_feature_names is not None:
            expected_columns = [str(col) for col in model_feature_names]
        else:
            expected_columns = list(base_feature_dict.keys())

        aligned_feature_dict = {
            col: base_feature_dict.get(col, 0.0) for col in expected_columns
        }
        X = pd.DataFrame([aligned_feature_dict], columns=expected_columns).apply(pd.to_numeric, errors="coerce").fillna(0.0)
        prediction = ML_MODEL.predict(X)[0]
        return round(float(prediction), 2)
    except Exception as e:
        logger.error(f"ML prediction error: {e}")
        return None


# ============================================================================
# FRONTEND ROUTES
# ============================================================================

@app.get("/", response_class=HTMLResponse)
async def serve_index(request: Request):
    """Serve index page. If not logged in, redirect to login."""
    # Check if user is authenticated
    token = request.cookies.get("token") or request.cookies.get("access_token")
    
    # If no token, redirect to login
    if not token:
        return RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    
    # User is logged in, serve index/route optimizer page
    index_path = TEMPLATES_DIR / "index.html"
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Traffic Route Analysis API</h1><p>Frontend not found. Please create templates/index.html</p>",
        status_code=404
    )


@app.get("/favicon.ico")
async def serve_favicon():
    """Serve the favicon."""
    favicon_path = STATIC_DIR / "favicon.svg"
    if os.path.exists(favicon_path):
        with open(favicon_path, "r", encoding="utf-8") as f:
            return Response(content=f.read(), media_type="image/svg+xml")
    # Return a default empty response to prevent 404 logs
    return Response(content="", status_code=204)


@app.get("/login", response_class=HTMLResponse)
async def serve_login():
    """Serve the login/registration page."""
    login_path = TEMPLATES_DIR / "login.html"
    if os.path.exists(login_path):
        with open(login_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Login</h1><p>Login page not found.</p>",
        status_code=404
    )


@app.get("/logout")
async def logout(response: Response):
    """Clear auth cookies and redirect to login."""
    redirect = RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    redirect.delete_cookie(key="token", path="/")
    redirect.delete_cookie(key="access_token", path="/")
    return redirect


@app.get("/auth/login")
async def serve_auth_login_alias(request: Request):
    """Alias route so auth redirects can target /auth/login for browser clients."""
    reason = request.query_params.get("reason")
    if reason:
        return RedirectResponse(url=f"/login?reason={reason}", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    return RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)


@app.get("/admin", response_class=HTMLResponse)
async def serve_admin(current_user: dict = Depends(require_role("admin"))):
    """Serve the admin dashboard page."""
    admin_path = TEMPLATES_DIR / "admin.html"
    if os.path.exists(admin_path):
        with open(admin_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Admin Dashboard</h1><p>Admin page not found.</p>",
        status_code=404
    )


@app.get("/account", response_class=HTMLResponse)
async def serve_account(current_user: dict = Depends(get_current_user)):
    """Serve the user account page."""
    role = current_user.get("role")
    if role == UserRole.police_supervisor.value:
        return RedirectResponse(url="/police/dashboard", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    if role == UserRole.admin.value:
        return RedirectResponse(url="/admin", status_code=status.HTTP_307_TEMPORARY_REDIRECT)

    account_path = TEMPLATES_DIR / "account.html"
    if os.path.exists(account_path):
        with open(account_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>My Account</h1><p>Account page not found.</p>",
        status_code=404
    )


@app.get("/password-toggle-demo", response_class=HTMLResponse)
async def serve_password_toggle_demo():
    """Serve the password toggle demo page."""
    demo_path = TEMPLATES_DIR / "password_toggle_demo.html"
    if os.path.exists(demo_path):
        with open(demo_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Password Toggle Demo</h1><p>Demo page not found.</p>",
        status_code=404
    )


@app.get("/analysis-report", response_class=HTMLResponse)
async def serve_analysis_report():
    """Serve the analysis report HTML page."""
    report_path = TEMPLATES_DIR / "analysis_report.html"
    if os.path.exists(report_path):
        with open(report_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Analysis Report</h1><p>Report page not found.</p>",
        status_code=404
    )


@app.get("/static/manifest.json")
async def get_manifest():
    """Serve PWA manifest."""
    manifest_path = STATIC_DIR / "manifest.json"
    if os.path.exists(manifest_path):
        return FileResponse(manifest_path, media_type="application/json")
    return JSONResponse({"error": "Manifest not found"}, status_code=404)


# ============================================================================
# API ROUTES
# ============================================================================

@app.get("/health")
async def health():
    """Health check endpoint."""
    _start_bootstrap_once()
    return {
        "status": "healthy" if APP_BOOTSTRAP_DONE and not APP_BOOTSTRAP_ERROR else "starting",
        "bootstrap_done": APP_BOOTSTRAP_DONE,
        "bootstrap_error": APP_BOOTSTRAP_ERROR,
        "model_loaded": ML_MODEL is not None,
        "timestamp": datetime.now(UTC).isoformat()
    }


@app.get("/api/stats")
async def get_stats():
    """Get real statistics from database for stats bar."""
    session = None
    try:
        session = get_session()
        
        # Count routes analyzed today
        today_start = datetime.now(UTC).replace(hour=0, minute=0, second=0, microsecond=0)
        today_count = session.query(AnalysisResult).filter(
            AnalysisResult.timestamp >= today_start
        ).count()
        
        # Get average congestion from last 10 records
        recent_records = session.query(AnalysisResult).order_by(
            AnalysisResult.timestamp.desc()
        ).limit(10).all()
        
        avg_congestion = 1.0
        if recent_records:
            ratios = []
            for r in recent_records:
                if r.travel_time_s and r.no_traffic_s and r.no_traffic_s > 0:
                    ratios.append(r.travel_time_s / r.no_traffic_s)
            if ratios:
                avg_congestion = round(sum(ratios) / len(ratios), 2)
        
        # Traffic status based on congestion
        if avg_congestion < 1.15:
            traffic_status = "Light"
            status_color = "#42a5f5"
        elif avg_congestion < 1.5:
            traffic_status = "Moderate"
            status_color = "#ffa726"
        else:
            traffic_status = "Heavy"
            status_color = "#ef5350"
        
        # Total all time count
        total_count = session.query(AnalysisResult).count()
        
        return {
            "routes_today": today_count,
            "total_routes": total_count,
            "avg_congestion": avg_congestion,
            "traffic_status": traffic_status,
            "status_color": status_color
        }
    except Exception as e:
        return {
            "routes_today": 0,
            "total_routes": 0,
            "avg_congestion": 1.0,
            "traffic_status": "Light",
            "status_color": "#42a5f5"
        }
    finally:
        if session is not None:
            session.close()


@app.get("/autocomplete")
async def autocomplete(q: str = Query(..., description="Search query")):
    """
    Get autocomplete suggestions from TomTom API.
    
    Args:
        q: Search query string
        
    Returns:
        List of suggestion objects
    """
    try:
        suggestions = tomtom_autocomplete(q)
        return {"suggestions": suggestions}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Autocomplete failed: {str(e)}")


@app.post("/analyze-route")
async def analyze_route(
    request: RouteAnalysisRequest,
    current_user: Optional[User] = Depends(get_optional_user)
):
    """
    Analyze route alternatives with cost calculation and ML prediction.
    
    Args:
        request: RouteAnalysisRequest with origin, destination, and parameters
        current_user: Optional authenticated user
        
    Returns:
        JSON with analyzed routes and best route recommendation
    """
    try:
        # Parse and validate origin
        if isinstance(request.origin, str):
            o_lat, o_lon = tomtom_geocode(request.origin)
            origin_data = {"name": request.origin, "lat": o_lat, "lon": o_lon}
        else:
            o_lat = float(request.origin.get("lat"))
            o_lon = float(request.origin.get("lon"))
            origin_data = request.origin
        
        # Parse and validate destination
        if isinstance(request.destination, str):
            d_lat, d_lon = tomtom_geocode(request.destination)
            dest_data = {"name": request.destination, "lat": d_lat, "lon": d_lon}
        else:
            d_lat = float(request.destination.get("lat"))
            d_lon = float(request.destination.get("lon"))
            dest_data = request.destination
        
        # Validate coordinates
        if not (-90 <= o_lat <= 90) or not (-180 <= o_lon <= 180):
            raise ValueError("Invalid origin coordinates")
        if not (-90 <= d_lat <= 90) or not (-180 <= d_lon <= 180):
            raise ValueError("Invalid destination coordinates")
        
        # Fetch routes from TomTom
        route_json = tomtom_route(
            o_lat, o_lon, d_lat, d_lon,
            maxAlternatives=request.maxAlternatives
        )
        
        routes = route_json.get("routes", [])
        if not routes:
            raise HTTPException(status_code=404, detail="No routes found")
        
        # Analyze each route
        analyzed_routes = []
        route_id = f"{origin_data.get('name', f'{o_lat},{o_lon}')}→{dest_data.get('name', f'{d_lat},{d_lon}')}"
        
        for idx, route in enumerate(routes):
            summary = summarize_route(route)
            
            if summary["length_m"] == 0:
                summary["length_m"] = haversine_m(o_lat, o_lon, d_lat, d_lon)
            
            cost = compute_route_cost(
                summary["travel_time_s"],
                summary["no_traffic_s"],
                summary["delay_s"],
                summary["length_m"],
                alpha=request.alpha,
                beta=request.beta,
                gamma=request.gamma
            )
            
            ml_predicted = predict_congestion({
                "distance_km": summary["length_m"] / 1000.0,
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": summary["delay_s"]
            })
            
            svr_predicted = None
            try:
                from .svr_model import svr_predict
                svr_predicted = svr_predict({})
            except Exception as e:
                svr_predicted = None
            
            congestion_ratio = (
                summary["travel_time_s"] / summary["no_traffic_s"]
                if summary["no_traffic_s"] and summary["no_traffic_s"] > 0
                else None
            )
            
            geometry = extract_route_geometry(route)
            
            calculated_delay = 0
            if summary["travel_time_s"] and summary["no_traffic_s"]:
                calculated_delay = max(0, summary["travel_time_s"] - summary["no_traffic_s"])
            elif summary.get("delay_s"):
                calculated_delay = summary["delay_s"]
            
            analyzed_route = {
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": calculated_delay,
                "length_m": summary["length_m"],
                "congestion_ratio": congestion_ratio,
                "calculated_cost": cost,
                "ml_predicted_congestion": ml_predicted,
                "svr_predicted_congestion": svr_predicted,
                "geometry": geometry
            }
            analyzed_routes.append(analyzed_route)
            
            # Save to database
            session = None
            try:
                session = get_session()
                save_analysis(session, {
                    "route_id": f"{route_id}_route{idx}",
                    "origin": origin_data,
                    "destination": dest_data,
                    "travel_time_s": summary["travel_time_s"],
                    "no_traffic_s": summary["no_traffic_s"],
                    "delay_s": summary["delay_s"],
                    "length_m": summary["length_m"],
                    "calculated_cost": cost,
                    "ml_predicted": ml_predicted,
                    "raw_json": route,
                    "user_id": current_user.id if current_user else None
                })
            except Exception as e:
                logger.error(f"Database save error: {e}")
            finally:
                if session is not None:
                    session.close()
        
        # Find best route (lowest cost)
        best_route = min(analyzed_routes, key=lambda x: x["calculated_cost"])
        
        return {
            "origin": origin_data,
            "destination": dest_data,
            "analyzed_routes": analyzed_routes,
            "best_route_index": best_route["route_index"],
            "best_route": best_route,
            "timestamp": datetime.now(UTC).isoformat()
        }
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Route analysis error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Route analysis failed: {str(e)}")


@app.get("/api/route-analysis/{route_id}")
async def get_route_analysis(route_id: str, route_index: Optional[int] = None):
    """
    Get detailed analysis data for a specific route.
    
    Args:
        route_id: Route identifier (e.g., "Origin→Destination")
        route_index: Optional route index to filter specific route variant
        
    Returns:
        Analysis data with historical trends and statistics
    """
    session = None
    try:
        session = get_session()
        
        route_id = route_id.replace('%E2%86%92', '→')
        
        if route_index is not None:
            query = session.query(AnalysisResult).filter(
                AnalysisResult.route_id == f"{route_id}_route{route_index}"
            )
        else:
            query = session.query(AnalysisResult).filter(
                AnalysisResult.route_id.like(f"{route_id}_route%")
            )
        
        results = query.order_by(AnalysisResult.timestamp.desc()).all()
        
        if not results:
            raise HTTPException(status_code=404, detail="No analysis data found for this route")
        
        analysis_data = []
        for r in results:
            try:
                origin = json.loads(r.origin) if r.origin and r.origin.startswith('{') else {"name": r.origin}
                dest = json.loads(r.destination) if r.destination and r.destination.startswith('{') else {"name": r.destination}
            except:
                origin = {"name": r.origin}
                dest = {"name": r.destination}
            
            delay_val = r.delay_s
            if delay_val is None or delay_val == 0:
                if r.travel_time_s and r.no_traffic_s:
                    delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                else:
                    delay_val = 0
            
            analysis_data.append({
                "id": r.id,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "route_id": r.route_id,
                "origin": origin,
                "destination": dest,
                "travel_time_s": r.travel_time_s,
                "no_traffic_s": r.no_traffic_s,
                "delay_s": delay_val,
                "length_m": r.length_m,
                "calculated_cost": r.calculated_cost,
                "ml_predicted": r.ml_predicted,
                "congestion_ratio": (r.travel_time_s / r.no_traffic_s) if r.no_traffic_s and r.no_traffic_s > 0 else None
            })
        
        # Calculate statistics
        if analysis_data:
            travel_times = [d["travel_time_s"] for d in analysis_data if d["travel_time_s"]]
            delays = []
            for d in analysis_data:
                delay_val = d.get("delay_s")
                if delay_val is None or delay_val == 0:
                    if d.get("travel_time_s") and d.get("no_traffic_s"):
                        delay_val = max(0, d["travel_time_s"] - d["no_traffic_s"])
                    else:
                        delay_val = 0
                if delay_val > 0:
                    delays.append(delay_val)
            
            costs = [d["calculated_cost"] for d in analysis_data if d.get("calculated_cost")]
            congestion_ratios = [d["congestion_ratio"] for d in analysis_data if d.get("congestion_ratio")]
            
            stats = {
                "avg_travel_time": sum(travel_times) / len(travel_times) if travel_times else 0,
                "avg_delay": sum(delays) / len(delays) if delays else 0,
                "avg_cost": sum(costs) / len(costs) if costs else 0,
                "avg_congestion": sum(congestion_ratios) / len(congestion_ratios) if congestion_ratios else 0,
                "min_travel_time": min(travel_times) if travel_times else 0,
                "max_travel_time": max(travel_times) if travel_times else 0,
                "total_records": len(analysis_data)
            }
        else:
            stats = {}
        
        response_data = {
            "route_id": route_id,
            "route_index": route_index,
            "analysis_data": analysis_data,
            "statistics": stats,
            "latest": analysis_data[0] if analysis_data else None,
            "fetched_at": datetime.now(UTC).isoformat()
        }
        
        response = JSONResponse(content=response_data)
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
        return response
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch analysis: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch analysis: {str(e)}")
    finally:
        if session is not None:
            session.close()


@app.post("/api/refresh-route")
async def refresh_route_analysis(
    request: RouteAnalysisRequest,
    current_user: Optional[User] = Depends(get_optional_user)
):
    """
    Refresh route analysis with latest data from TomTom API.
    Similar to analyze_route but returns data in format suitable for analysis report.
    """
    try:
        # Reuse analyze_route logic
        if isinstance(request.origin, str):
            o_lat, o_lon = tomtom_geocode(request.origin)
            origin_data = {"name": request.origin, "lat": o_lat, "lon": o_lon}
        else:
            o_lat = float(request.origin.get("lat"))
            o_lon = float(request.origin.get("lon"))
            origin_data = request.origin
        
        if isinstance(request.destination, str):
            d_lat, d_lon = tomtom_geocode(request.destination)
            dest_data = {"name": request.destination, "lat": d_lat, "lon": d_lon}
        else:
            d_lat = float(request.destination.get("lat"))
            d_lon = float(request.destination.get("lon"))
            dest_data = request.destination
        
        # Validate coordinates
        if not (-90 <= o_lat <= 90) or not (-180 <= o_lon <= 180):
            raise ValueError("Invalid origin coordinates")
        if not (-90 <= d_lat <= 90) or not (-180 <= d_lon <= 180):
            raise ValueError("Invalid destination coordinates")
        
        route_json = tomtom_route(
            o_lat, o_lon, d_lat, d_lon,
            maxAlternatives=request.maxAlternatives
        )
        
        routes = route_json.get("routes", [])
        if not routes:
            raise HTTPException(status_code=404, detail="No routes found")
        
        analyzed_routes = []
        route_id = f"{origin_data.get('name', f'{o_lat},{o_lon}')}→{dest_data.get('name', f'{d_lat},{d_lon}')}"
        
        for idx, route in enumerate(routes):
            summary = summarize_route(route)
            
            if summary["length_m"] == 0:
                summary["length_m"] = haversine_m(o_lat, o_lon, d_lat, d_lon)
            
            cost = compute_route_cost(
                summary["travel_time_s"],
                summary["no_traffic_s"],
                summary["delay_s"],
                summary["length_m"],
                alpha=request.alpha,
                beta=request.beta,
                gamma=request.gamma
            )
            
            ml_predicted = predict_congestion({
                "distance_km": summary["length_m"] / 1000.0,
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": summary["delay_s"]
            })
            
            congestion_ratio = (
                summary["travel_time_s"] / summary["no_traffic_s"]
                if summary["no_traffic_s"] and summary["no_traffic_s"] > 0
                else None
            )
            
            calculated_delay = 0
            if summary["travel_time_s"] and summary["no_traffic_s"]:
                calculated_delay = max(0, summary["travel_time_s"] - summary["no_traffic_s"])
            elif summary.get("delay_s"):
                calculated_delay = summary["delay_s"]
            
            analyzed_route = {
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": calculated_delay,
                "length_m": summary["length_m"],
                "congestion_ratio": congestion_ratio,
                "calculated_cost": cost,
                "ml_predicted_congestion": ml_predicted
            }
            analyzed_routes.append(analyzed_route)
            
            # Save to database
            session = None
            try:
                session = get_session()
                save_analysis(session, {
                    "route_id": f"{route_id}_route{idx}",
                    "origin": origin_data,
                    "destination": dest_data,
                    "travel_time_s": summary["travel_time_s"],
                    "no_traffic_s": summary["no_traffic_s"],
                    "delay_s": summary["delay_s"],
                    "length_m": summary["length_m"],
                    "calculated_cost": cost,
                    "ml_predicted": ml_predicted,
                    "raw_json": route,
                    "user_id": current_user.id if current_user else None
                })
            except Exception as e:
                logger.error(f"Database save error: {e}")
            finally:
                if session is not None:
                    session.close()
        
        best_route = min(analyzed_routes, key=lambda x: x["calculated_cost"])
        
        return {
            "origin": origin_data,
            "destination": dest_data,
            "route_id": route_id,
            "analyzed_routes": analyzed_routes,
            "best_route_index": best_route["route_index"],
            "best_route": best_route,
            "timestamp": datetime.now(UTC).isoformat()
        }
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Route refresh error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Route refresh failed: {str(e)}")


# ============================================================================
# USER AUTHENTICATION & MANAGEMENT
# ============================================================================

@app.post("/api/auth/register", response_model=UserResponse)
async def register_user(user_data: UserCreate, db: Session = Depends(get_db)):
    """Register a new user."""
    try:
        if len(user_data.password) < 6:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Password must be at least 6 characters long"
            )
        
        MAX_PASSWORD_LENGTH = 10000
        if len(user_data.password) > MAX_PASSWORD_LENGTH:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Password is too long. Maximum {MAX_PASSWORD_LENGTH} characters allowed."
            )

        requested_role = (user_data.role or "user").strip().lower()
        department_map = {
            "user": "general",
            "police_supervisor": "police",
            "logistics_manager": "logistics",
        }
        if requested_role not in department_map:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid role selected during registration",
            )
        
        user_payload = user_data.dict()
        user_payload["department"] = department_map[requested_role]
        user_payload.pop("role", None)

        user = create_user(db, AuthUserCreate(**user_payload))
        return UserResponse.model_validate(user)
    except HTTPException:
        raise
    except Exception as e:
        error_msg = str(e)
        if "password" in error_msg.lower():
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Password validation error: {error_msg}"
            )
        raise HTTPException(status_code=500, detail=f"Registration failed: {error_msg}")


@app.get("/test/create-user")
async def create_test_user(db: Session = Depends(get_db)):
    """Create a test user for development/testing (remove in production)."""
    # Check if test user already exists
    existing_user = get_user_by_username(db, "testuser")
    if existing_user:
        return {
            "status": "user_exists",
            "message": "Test user already exists",
            "username": "testuser",
            "password": "password123"
        }
    
    # Create test user
    try:
        test_user_data = AuthUserCreate(
            email="testuser@example.com",
            username="testuser",
            password="password123",
            full_name="Test User",
            department="general"
        )
        user = create_user(db, test_user_data)
        return {
            "status": "created",
            "message": "Test user created successfully",
            "username": "testuser",
            "password": "password123",
            "email": "testuser@example.com"
        }
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create test user: {str(e)}"
        )


@app.post("/api/auth/login", response_model=RoleToken)
async def login_user(
    form_data: OAuth2PasswordRequestForm = Depends(),
    db: Session = Depends(get_db),
    *,
    response: Response,
):
    """Login and get access token."""
    user = authenticate_user(db, form_data.username, form_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token_expires = timedelta(minutes=30 * 24 * 60)
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    user_role = (getattr(user, "role", None) or "user").strip().lower()
    user_district_id = getattr(user, "district_id", None)

    role = UserRole.user
    district_id = None
    fleet_zone = None
    if user.is_admin:
        role = UserRole.admin
    elif user_role == "police_supervisor":
        role = UserRole.police_supervisor
        district_id = user_district_id or "district_1"
    elif user_department == "police":
        role = UserRole.police_supervisor
        district_id = user_district_id or "district_1"
    elif user_role == "logistics_manager":
        role = UserRole.logistics_manager
        fleet_zone = "zone_default"
    elif user_department == "logistics":
        role = UserRole.logistics_manager
        fleet_zone = "zone_default"

    access_token = create_role_access_token(
        username=user.username,
        role=role,
        district_id=district_id,
        fleet_zone=fleet_zone,
        expires_delta=access_token_expires,
    )
    # Set secure=False for localhost/HTTP development, True for production/HTTPS
    is_secure = os.getenv("SECURE_COOKIES", "false").lower() == "true"
    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=is_secure,
        samesite="Lax",  # Use Lax for both secure and non-secure (None requires Secure=True)
        max_age=30 * 24 * 60 * 60,
        path="/",
    )
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "role": role,
        "district_id": district_id,
        "fleet_zone": fleet_zone,
    }


@app.get("/police/dashboard", response_class=HTMLResponse)
async def police_dashboard(request: Request, current_user: dict = Depends(require_police_department_user())):
    """Render the police supervisor dashboard for the user's district."""
    district_id = _resolve_current_user_district_id(current_user)
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for police dashboard access",
        )

    # Initialize/reset officer statuses for the district
    _ensure_officer_statuses_initialized(district_id)

    context = _build_police_dashboard_context(current_user, district_id)
    context.update({
        "request": request,
        "current_user": current_user,
        "district_id": district_id,
        "district": _format_district_label(district_id),
        "google_maps_api_key": os.getenv("GOOGLE_MAPS_API_KEY", ""),
        "data_error": False,
    })

    return templates.TemplateResponse(
        request=request,
        name="police/supervisor.html",
        context=context,
    )


@app.get("/supervisor/analytics", response_class=HTMLResponse)
async def supervisor_analytics_page(request: Request, current_user: dict = Depends(require_police_department_user())):
    """Render the supervisor analytics dashboard page."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for supervisor analytics",
        )

    return templates.TemplateResponse(
        request=request,
        name="police/supervisor_analytics.html",
        context={
            "request": request,
            "current_user": current_user,
            "district_id": district_id,
        },
    )


@app.get("/api/supervisor/analytics")
async def supervisor_analytics_data(
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_db),
):
    """Return analytics series for supervisor command center dashboards."""
    now = datetime.now(UTC)
    district_id = current_user.get("district_id") or ""

    seven_days_ago = now - timedelta(days=6)
    recent_feedback = (
        _feedback_rows_for_district(db, district_id)
        .filter(MLFeedback.created_at >= seven_days_ago)
        .order_by(MLFeedback.created_at.asc())
        .all()
    )

    day_labels = [(seven_days_ago + timedelta(days=i)).date().isoformat() for i in range(7)]
    incident_volume_map: dict[str, dict[str, int]] = {day: {} for day in day_labels}
    for row in recent_feedback:
        day_key = (row.created_at or now).date().isoformat()
        if day_key not in incident_volume_map:
            continue
        series_key = f"{(row.severity or 'unknown').lower()} | {row.zone or 'Unknown Zone'}"
        incident_volume_map[day_key][series_key] = incident_volume_map[day_key].get(series_key, 0) + 1

    all_volume_keys = sorted({
        key
        for day_data in incident_volume_map.values()
        for key in day_data.keys()
    })
    incident_volume_datasets = [
        {
            "label": key,
            "data": [incident_volume_map[day].get(key, 0) for day in day_labels],
        }
        for key in all_volume_keys
    ]

    weekly_labels: list[str] = []
    weekly_values: list[float] = []
    for week_index in range(3, -1, -1):
        start = (now - timedelta(days=now.weekday())) - timedelta(weeks=week_index)
        end = start + timedelta(days=7)
        rows = (
            _feedback_rows_for_district(db, district_id)
            .filter(MLFeedback.created_at >= start, MLFeedback.created_at < end)
            .all()
        )
        avg_response = round(sum(float(item.response_time_minutes) for item in rows) / len(rows), 2) if rows else 0.0
        weekly_labels.append(start.date().isoformat())
        weekly_values.append(avg_response)

    month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    monthly_workloads = (
        db.query(OfficerIncidentCount)
        .filter(OfficerIncidentCount.last_updated >= month_start)
        .all()
    )
    utilization_map: dict[str, int] = {}
    for item in monthly_workloads:
        total = int(item.incident_count_critical + item.incident_count_high + item.incident_count_medium + item.incident_count_low)
        utilization_map[item.officer_name or item.officer_username] = utilization_map.get(item.officer_name or item.officer_username, 0) + total
    sorted_utilization = sorted(utilization_map.items(), key=lambda pair: pair[1], reverse=True)

    zone_counts: dict[str, int] = {}
    for item in recent_feedback:
        zone_name = item.zone or "Unknown Zone"
        zone_counts[zone_name] = zone_counts.get(zone_name, 0) + 1
    top_zones = sorted(zone_counts.items(), key=lambda pair: pair[1], reverse=True)[:5]

    latest_retrain = (
        db.query(MLRetrainAudit)
        .order_by(MLRetrainAudit.retrained_at.desc())
        .first()
    )

    return {
        "district_id": district_id,
        "incident_volume": {
            "labels": day_labels,
            "datasets": incident_volume_datasets,
        },
        "response_time_trend": {
            "labels": weekly_labels,
            "data": weekly_values,
        },
        "officer_utilization": {
            "labels": [name for name, _ in sorted_utilization],
            "data": [count for _, count in sorted_utilization],
        },
        "top_zones": [
            {"zone": zone_name, "count": count}
            for zone_name, count in top_zones
        ],
        "last_retrained": latest_retrain.retrained_at.isoformat() if latest_retrain else None,
        "updated_at": now.isoformat(),
    }


@app.get("/patrol-units")
@app.get("/police/units/live")
async def live_patrol_units(current_user: dict = Depends(require_police_department_user())):
    """Return patrol units for the dispatch dashboard."""
    district_id = _resolve_current_user_district_id(current_user)
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for patrol unit status",
        )

    _ensure_officer_statuses_initialized(district_id)

    incidents = _load_police_incidents(district_id)
    supervisor_name = current_user.get("username", "Unknown Supervisor")
    assignments = _get_dispatch_assignments(district_id)
    patrol_units = _build_patrol_units(district_id, incidents, supervisor_name, assignments)
    assigned_incident_ids = {assignment.incident_id for assignment in assignments}
    unassigned_incidents = [incident for incident in incidents if incident.get("id") not in assigned_incident_ids]
    available_patrol_units = [unit for unit in patrol_units if unit["status"] == "available"]

    return {
        "patrol_units": patrol_units,
        "available_patrol_units": available_patrol_units,
        "available_count": len(available_patrol_units),
        "unassigned_incidents": unassigned_incidents,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/incidents")
@app.get("/api/incidents")
async def api_incidents(current_user: dict = Depends(require_any_role("police_supervisor", "police_officer"))):
    """Return recent incidents with assignment status for the command center."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incidents",
        )

    raw_incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    incidents = _serialize_recent_incidents(district_id, raw_incidents, assignments)

    return {
        "incidents": incidents,
        "district_id": district_id,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/api/officers/status")
async def api_officers_status(current_user: dict = Depends(require_any_role("police_supervisor", "police_officer", "admin"))):
    """Return officer status rows using command-center JSON contract."""
    district_id = _resolve_current_user_district_id(current_user)
    # For admin users, default to district_1 if not provided
    if not district_id and current_user.get("role") == "admin":
        district_id = "district_1"
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for officer status",
        )

    # Ensure officer statuses are initialized/reset
    _ensure_officer_statuses_initialized(district_id)

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    patrol_units = _build_patrol_units(
        district_id,
        incidents,
        current_user.get("username", "Unknown Supervisor"),
        assignments,
    )

    officers = [
        {
            "patrol_id": unit.get("patrol_id") or unit.get("unit_id"),
            "id": unit.get("unit_id"),
            "unit_name": unit.get("unit_name") or unit.get("unit_id"),
            "name": unit.get("officer_name") or unit.get("officer") or "Unassigned",
            "badge": unit.get("unit_id"),
            "status": str(unit.get("status") or "available").lower(),
            "skills": [],
            "district_id": unit.get("district_id") or district_id,
            "current_area": unit.get("current_area") or unit.get("current_location") or "District duty",
            "current_location": unit.get("current_location") or "District duty",
            "latitude": unit.get("latitude"),
            "longitude": unit.get("longitude"),
            "last_updated": unit.get("last_updated") or "-",
            "updated_at": unit.get("updated_at") or unit.get("last_updated") or "-",
        }
        for unit in patrol_units
    ]

    return {
        "officers": officers,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.post("/api/incident/new")
async def create_incident(
    request: NewIncidentRequest,
    current_user: dict = Depends(require_police_department_user()),
):
    """Create a new district incident and push a real-time update over Socket.IO."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident creation",
        )

    incident_id = f"manual-{uuid.uuid4().hex[:10]}"
    zone_name = request.zone or _infer_zone_name(
        district_id,
        {"latitude": request.latitude, "longitude": request.longitude, "description": request.description},
    )
    incident_record = {
        "id": incident_id,
        "type": request.incident_type,
        "severity": _normalize_severity(request.severity),
        "description": request.description,
        "latitude": request.latitude,
        "longitude": request.longitude,
        "zone": zone_name,
        "start_time": datetime.now(UTC).isoformat(),
        "district_id": district_id,
    }
    _manual_incidents_store.setdefault(district_id, []).append(incident_record)

    await emit_incident_new(
        sio,
        district_id,
        incident_record,
        actor=current_user.get("username", "Unknown Supervisor"),
    )
    await emit_incident_updated(
        sio,
        district_id,
        incident_record,
        update_type="created",
        actor=current_user.get("username", "Unknown Supervisor"),
    )

    # Fire a real-time alert so the supervisor alert sidebar lights up
    severity_label = _normalize_severity(request.severity)
    add_alert(
        district_id=district_id,
        severity=severity_label,
        message=f"New {severity_label} incident reported in {zone_name}: {request.description or request.incident_type}",
        incident_id=incident_id,
    )

    return {
        "status": "created",
        "incident": incident_record,
    }


@app.get("/police/incidents/feed")
async def live_incidents_feed(current_user: dict = Depends(require_police_department_user())):
    """Return district-scoped incident feed sorted by reported time (desc)."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident feed",
        )

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    incidents_feed = _build_incidents_feed(district_id, incidents, assignments)

    return {
        "incidents": incidents_feed,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/police/heatmap/data")
async def police_heatmap_data(current_user: dict = Depends(require_police_department_user())):
    """Return district-scoped incident heatmap points for the supervisor."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for heatmap data",
        )

    incidents = _load_police_incidents(district_id)
    return _build_heatmap_points(incidents)


@app.get("/police/response-times")
@app.get("/police/response_times")
@app.get("/api/police/response-times")
async def police_response_times(current_user: dict = Depends(require_police_department_user())):
    """Return zone-wise average response metrics for today's shift."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for response time metrics",
        )

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    zone_metrics = _build_response_time_by_zone(district_id, incidents, assignments)

    zone_times = sorted(float(zone["avg_response_time"]) for zone in zone_metrics if zone.get("avg_response_time") is not None)
    target_threshold_minutes = 8.0
    if zone_times and all(value > 8.0 for value in zone_times):
        target_threshold_minutes = float(zone_times[len(zone_times) // 2])
    for zone in zone_metrics:
        zone["exceeds_target"] = float(zone.get("avg_response_time") or 0.0) > float(target_threshold_minutes)
        zone["avg_time"] = zone.get("avg_response_time")
        zone["total_units"] = zone.get("total_incidents")

    return {
        "district_id": district_id,
        "target_threshold_minutes": target_threshold_minutes,
        "zones": zone_metrics,
        "zone_count": len(zone_metrics),
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/api/logistics/alerts")
async def logistics_shared_alerts(current_user: dict = Depends(require_role("logistics_manager"))):
    """Return active sanitized police zone alerts for logistics users only."""
    session = get_session()
    try:
        now = datetime.now(UTC)
        rows = (
            session.query(SharedAlert)
            .filter(SharedAlert.expires_at >= now)
            .order_by(SharedAlert.timestamp.desc())
            .limit(25)
            .all()
        )
    finally:
        session.close()

    return {
        "alerts": [
            {
                "alert_id": row.alert_id,
                "zone": row.zone,
                "severity": row.severity,
                "timestamp": row.timestamp.isoformat() if row.timestamp else None,
                "affected_roads": row.affected_roads or [],
                "expires_at": row.expires_at.isoformat() if row.expires_at else None,
            }
            for row in rows
        ],
        "updated_at": datetime.now(UTC).isoformat(),
    }


def _ensure_officer_statuses_initialized(district_id: str):
    """
    Ensure the district's mock patrol slots have dispatch-status rows.
    This keeps the visible patrol cards and dispatch validation aligned.
    """
    session = get_session()
    try:
        # Generate unit IDs based on mock data
        incidents = _load_police_incidents(district_id)
        police_users = (
            session.query(User)
            .filter(
                User.department == "police",
                User.district_id == district_id,
                User.is_active == True,  # noqa: E712
            )
            .count()
        )
        unit_count = min(max(4, police_users), 15)
        
        for index in range(unit_count):
            # Use consistent unit ID format: DISTRICT-X-UXX
            district_prefix = district_id.upper().replace("_", "-")
            unit_id = f"{district_prefix}-U{index + 1:02d}"
            
            # Check if officer status exists in database
            officer_status = session.query(OfficerDispatchStatus).filter(
                OfficerDispatchStatus.officer_id == unit_id,
                OfficerDispatchStatus.district_id == district_id
            ).first()
            
            if officer_status is None:
                officer_status = OfficerDispatchStatus(
                    district_id=district_id,
                    officer_id=unit_id,
                    status="available",
                    assigned_incident_id=None,
                    mobile_token=None
                )
                session.add(officer_status)
            else:
                # Refresh mock patrol slots so the UI and dispatch validator
                # agree on their current availability.
                officer_status.status = "available"
                officer_status.assigned_incident_id = None
                officer_status.updated_at = datetime.now(UTC)
            
        session.commit()
        logger.debug(f"Initialized {unit_count} officer statuses for {district_id}")
    except Exception as e:
        logger.warning(f"Could not initialize officer statuses for {district_id}: {e}")
        session.rollback()
    finally:
        session.close()


@app.post("/dispatch")
async def dispatch_selected_patrol(
    request: PatrolSelectionDispatchRequest,
    current_user: dict = Depends(require_police_department_user()),
):
    """Assign a selected patrol unit and mark it RESPONDING."""
    if request.incident_id:
        return await dispatch_patrol_unit(
            DispatchRequest(
                incident_id=request.incident_id,
                officer_id=request.officer_id or request.patrol_id,
                unit_id=request.unit_id,
            ),
            current_user=current_user,
        )

    district_id = _resolve_current_user_district_id(current_user)
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for dispatch",
        )

    patrol_id = (request.patrol_id or "").strip()
    user_id = (request.user_id or current_user.get("username") or "").strip()
    location = (request.location or "").strip()
    if not user_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="user_id is required",
        )

    caller_id = str(current_user.get("username") or "").strip()
    if caller_id and caller_id != user_id:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="user_id does not match logged-in user",
        )

    incidents = _load_police_incidents(district_id)
    patrol_units = _build_patrol_units(
        district_id,
        incidents,
        current_user.get("username", "Unknown Supervisor"),
        _get_dispatch_assignments(district_id),
    )
    patrol = next((unit for unit in patrol_units if str(unit.get("patrol_id") or unit.get("unit_id")) == patrol_id), None)
    if not patrol:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Patrol unit not found")

    current_status = str(patrol.get("status") or "").lower()
    if current_status in {"busy", "responding", "offline"}:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=f"Patrol is already {current_status.upper()}",
        )

    now = datetime.now(UTC)
    dispatch_target = f"Supervisor dispatch in {_format_district_label(district_id)}"
    dispatch_target_id = f"dispatch:{district_id}:{patrol_id}"
    with _patrol_user_assignments_lock:
        district_assignments = _patrol_user_assignments.setdefault(district_id, {})
        existing = district_assignments.get(patrol_id)
        if existing and str(existing.get("status") or "").lower() in {"responding", "busy"}:
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="Patrol is already assigned",
            )

        district_assignments[patrol_id] = {
            "user_id": user_id,
            "status": "responding",
            "assigned_at": now,
        }

    session = get_session()
    try:
        officer_status = (
            session.query(OfficerDispatchStatus)
            .filter(
                OfficerDispatchStatus.district_id == district_id,
                OfficerDispatchStatus.officer_id == patrol_id,
            )
            .first()
        )
        if officer_status is None:
            officer_status = OfficerDispatchStatus(
                district_id=district_id,
                officer_id=patrol_id,
            )
            session.add(officer_status)

        officer_status.status = "responding"
        officer_status.assigned_incident_id = dispatch_target_id
        officer_status.updated_at = now

        session.add(
            DispatchLog(
                district_id=district_id,
                incident_id=dispatch_target_id,
                officer_id=patrol_id,
                assigned_by=str(current_user.get("username") or user_id),
                assigned_at=now,
                status="dispatched",
            )
        )
        session.commit()
    except Exception:
        session.rollback()
        raise
    finally:
        session.close()

    # Fire a real-time alert for this patrol dispatch
    add_alert(
        district_id=district_id,
        severity="medium",
        message=f"Patrol {patrol_id} dispatched to {dispatch_target}",
    )

    return {
        "success": True,
        "patrol_id": patrol_id,
        "user_id": user_id,
        "status": "RESPONDING",
        "assigned_at": now.isoformat(),
        "message": f"Patrol {patrol_id} dispatched successfully",
    }


@app.post("/api/dispatch")
@app.post("/police/dispatch")
async def dispatch_patrol_unit(
    request: DispatchRequest,
    current_user: dict = Depends(require_police_department_user()),
):
    """Assign a patrol unit to an incident and persist the dispatch.
    
    Validates officer availability, updates dispatch status, records in audit log,
    sends FCM push notification, and broadcasts via SocketIO.
    
    Request:
        incident_id (str): ID of the incident to dispatch
        officer_id (str): ID of the officer/unit to dispatch
    
    Response:
        {
            "success": true,
            "dispatch_id": 42,
            "eta": 5,
            "incident_id": "incident_123",
            "officer_id": "unit_001"
        }
    
    Raises:
        400: Missing district_id, incident_id, or officer_id
        403: User is not a police supervisor
        404: Incident not found
        409: Officer not available or already assigned
        500: Dispatch failed
    """
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for dispatch",
        )

    # Ensure officer statuses are properly initialized before dispatch
    _ensure_officer_statuses_initialized(district_id)

    incidents = _load_police_incidents(district_id)
    incident_map = {incident.get("id"): incident for incident in incidents}
    target_incident = incident_map.get(request.incident_id)
    if not target_incident:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Incident not found",
        )

    officer_id = (request.officer_id or request.unit_id or "").strip()
    if not officer_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="officer_id is required for dispatch",
        )

    session = get_session()
    dispatch_log_record = None
    officer_mobile_token = None
    
    try:
        now = datetime.now(UTC)
        
        # Clear session cache to avoid stale data
        session.expunge_all()
        
        # Check officer status before dispatch
        officer_status = (
            session.query(OfficerDispatchStatus)
            .filter(OfficerDispatchStatus.officer_id == officer_id)
            .with_for_update()
            .first()
        )
        
        if officer_status is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Officer {officer_id} not found in dispatch system",
            )
        
        # Validate officer status is 'available' before dispatch
        officer_status_lower = (officer_status.status or "").lower()
        if officer_status_lower != "available":
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=f"Officer is not available (current status: {officer_status_lower})",
            )
        
        # Store mobile token for FCM notification
        officer_mobile_token = officer_status.mobile_token
        
        # Query fresh data from database and delete ALL existing assignments for this incident and unit
        session.query(PoliceDispatchAssignment).filter(
            PoliceDispatchAssignment.incident_id == request.incident_id
        ).delete(synchronize_session=False)
        
        session.query(PoliceDispatchAssignment).filter(
            PoliceDispatchAssignment.unit_id == officer_id
        ).delete(synchronize_session=False)
        
        session.commit()
        session.expunge_all()
        
        # Create new assignment
        assignment = PoliceDispatchAssignment(
            district_id=district_id,
            incident_id=request.incident_id,
            unit_id=officer_id,
            assigned_by=current_user.get("username", "Unknown Supervisor"),
            assigned_at=now,
            status="active",
        )
        session.add(assignment)

        # Update officer status to 'enroute'
        officer_status.status = "enroute"
        officer_status.assigned_incident_id = request.incident_id
        officer_status.updated_at = now
        session.add(officer_status)

        # Create dispatch log record
        dispatch_log_record = DispatchLog(
            district_id=district_id,
            incident_id=request.incident_id,
            officer_id=officer_id,
            assigned_by=current_user.get("username", "Unknown Supervisor"),
            assigned_at=now,
            status="dispatched",
        )
        session.add(dispatch_log_record)

        shared_alert = _create_shared_alert_for_dispatch(session, district_id, target_incident)
        session.commit()
        session.refresh(assignment)
        session.refresh(dispatch_log_record)
        if shared_alert is not None:
            session.refresh(shared_alert)
            
    except HTTPException:
        session.rollback()
        raise
    except IntegrityError as exc:
        session.rollback()
        existing_assignment = (
            session.query(PoliceDispatchAssignment)
            .filter(PoliceDispatchAssignment.incident_id == request.incident_id)
            .first()
        )
        if existing_assignment and existing_assignment.unit_id == officer_id:
            logger.warning(
                "Dispatch assignment already exists for incident %s and unit %s",
                request.incident_id,
                officer_id,
            )
            assignment = existing_assignment
        else:
            logger.error("Dispatch assignment integrity error: %s", exc)
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="Incident or patrol unit is already assigned",
            )
    except Exception as exc:
        session.rollback()
        logger.error("Dispatch assignment failed: %s", exc)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Dispatch assignment failed",
        )
    finally:
        session.close()

    # Calculate ETA in minutes (simplified: use distance if available)
    eta_minutes = 5  # Default ETA in minutes
    supervisor_id = current_user.get("username", "Unknown Supervisor")
    
    # Send FCM push notification to officer's mobile device
    fcm_result = send_officer_dispatch_notification(
        officer_mobile_token,
        officer_id,
        target_incident,
    )
    logger.info(
        f"Dispatch notification for officer {officer_id}: "
        f"sent={fcm_result.get('sent')}, reason={fcm_result.get('reason')}"
    )
    
    updated_context = _build_police_dashboard_context(current_user, district_id)
    updated_unit = next((unit for unit in updated_context["patrol_units"] if unit["unit_id"] == officer_id), None)
    updated_incident = next((incident for incident in updated_context["incidents"] if incident.get("id") == request.incident_id), None)
    incident_reference = updated_incident or target_incident or {"id": request.incident_id}
    zone_name = (
        incident_reference.get("zone")
        or incident_reference.get("zone_name")
        or _infer_zone_name(district_id, incident_reference)
    )

    # Build dispatch data for SocketIO event
    dispatch_data = {
        "dispatch_id": dispatch_log_record.id if dispatch_log_record else None,
        "officer_id": officer_id,
        "officer_name": (updated_unit or {}).get("officer_name") if isinstance(updated_unit, dict) else None,
        "incident_id": request.incident_id,
        "incident_description": incident_reference.get("description"),
        "zone": zone_name,
        "district_id": district_id,
        "eta": eta_minutes,
        "supervisor_id": supervisor_id,
    }

    # Emit SocketIO event to police namespace
    await emit_officer_dispatched(
        sio,
        district_id,
        dispatch_data,
        actor=supervisor_id,
    )
    
    # Also emit the existing events for backward compatibility

    await emit_incident_updated(
        sio,
        district_id,
        incident_reference,
        update_type="dispatched",
        actor=supervisor_id,
    )
    await emit_officer_status_changed(
        sio,
        district_id,
        {
            "id": officer_id,
            "name": (updated_unit or {}).get("officer_name") if isinstance(updated_unit, dict) else None,
            "badge": officer_id,
            "status": "enroute",
            "district_id": district_id,
            "incident_id": request.incident_id,
        },
        actor=supervisor_id,
    )

    # Fire a real-time alert for this dispatch
    officer_display = (
        (updated_unit or {}).get("officer_name") if isinstance(updated_unit, dict) else None
    ) or officer_id
    incident_severity = (
        (updated_incident or target_incident or {}).get("severity") or "medium"
    )
    add_alert(
        district_id=district_id,
        severity=_normalize_severity(incident_severity),
        message=f"Unit {officer_display} dispatched to {zone_name} — ETA {eta_minutes} min",
        incident_id=request.incident_id,
    )

    # Return simplified response format
    return {
        "success": True,
        "dispatch_id": dispatch_log_record.id if dispatch_log_record else None,
        "eta": eta_minutes,
        "incident_id": request.incident_id,
        "officer_id": officer_id,
        "fcm_sent": fcm_result.get("sent", False),
    }


@app.get("/police/export/pptx")
async def police_export_pptx(current_user: dict = Depends(require_police_department_user())):
    """Generate a police supervisor shift report as a PPTX file."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for police report export",
        )

    try:
        officer_name = current_user.get("username", "Unknown Officer")
        incidents = _load_police_incidents(district_id)
        district_summary = _build_district_summary(incidents, district_id)
        ml_predictions = _predict_police_hotspots(district_id, incidents)
        output = generate_shift_pptx(district_id, officer_name, incidents, district_summary, ml_predictions)

        filename = f"ShiftReport_District{district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except Exception as exc:
        logger.error(f"Police PPTX export failed: {exc}")
        raise HTTPException(
            status_code=500,
            detail="Report generation failed. Please try again or contact admin.",
        )


def generate_detailed_shift_report_pptx(
    shift_id: int,
    shift_data: dict,
    incidents: list[dict],
    officer_workload: list[dict],
    dispatch_assignments: list[dict],
) -> io.BytesIO:
    """Generate a comprehensive shift report PPTX with 5 slides.
    
    Slides:
    1. Shift summary (date, supervisor, total incidents, avg response time)
    2. Incident breakdown table (ID, type, severity, location, officer, resolved)
    3. Officer performance (name, incidents handled, avg response time)
    4. Unresolved/escalated cases
    5. Zone heatmap / incident distribution
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # ========== SLIDE 1: SHIFT SUMMARY ==========
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Shift Summary Report"
    
    shift_start = shift_data.get("start_time")
    shift_end = shift_data.get("end_time")
    if isinstance(shift_start, str):
        shift_start = datetime.fromisoformat(shift_start.replace('Z', '+00:00'))
    if isinstance(shift_end, str):
        shift_end = datetime.fromisoformat(shift_end.replace('Z', '+00:00'))
    
    duration_mins = 0
    if shift_start and shift_end:
        duration_mins = int((shift_end - shift_start).total_seconds() / 60)
    
    avg_response_time = "N/A"
    if dispatch_assignments:
        response_times = [a.get("response_time_s", 0) for a in dispatch_assignments if a.get("response_time_s")]
        if response_times:
            avg_response_time = f"{sum(response_times) / len(response_times):.1f}s"
    
    unresolved_count = sum(1 for inc in incidents if inc.get("status") != "resolved")
    
    summary_text = (
        f"District: {shift_data.get('district_id', 'Unknown')}\n"
        f"Supervisor: {shift_data.get('supervisor_name', 'Unknown')}\n"
        f"Date: {shift_start.strftime('%Y-%m-%d') if shift_start else 'N/A'}\n"
        f"Shift Duration: {duration_mins // 60}h {duration_mins % 60}m\n"
        f"\n"
        f"Total Incidents: {len(incidents)}\n"
        f"Resolved: {len(incidents) - unresolved_count}\n"
        f"Unresolved: {unresolved_count}\n"
        f"Avg Response Time: {avg_response_time}\n"
        f"Officers on Duty: {shift_data.get('officers_on_duty', 0)}"
    )
    
    slide1.placeholders[1].text = summary_text

    # ========== SLIDE 2: INCIDENT BREAKDOWN TABLE ==========
    slide2 = prs.slides.add_slide(blank_layout)
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Incident Breakdown"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True

    # Create incident table
    rows = min(len(incidents) + 1, 12)  # Cap at 12 rows for slide fit
    cols = 6
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.4), Inches(1), Inches(9.2), Inches(6))
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.2)  # ID
    table.columns[1].width = Inches(1.5)  # Type
    table.columns[2].width = Inches(1.2)  # Severity
    table.columns[3].width = Inches(2)    # Location
    table.columns[4].width = Inches(1.5)  # Officer
    table.columns[5].width = Inches(1.8)  # Status
    
    # Header row
    headers = ["ID", "Type", "Severity", "Location", "Officer", "Status"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)

    # Data rows
    for row_idx, incident in enumerate(incidents[:rows - 1], start=1):
        table.cell(row_idx, 0).text = _safe_ppt_text(incident.get("id", "N/A")[:10])
        table.cell(row_idx, 1).text = _safe_ppt_text(incident.get("type", "N/A")[:15])
        
        severity = incident.get("severity", "N/A")
        table.cell(row_idx, 2).text = severity
        
        table.cell(row_idx, 3).text = _safe_ppt_text(incident.get("location", "N/A")[:20])
        
        # Get officer from workload or assignment
        officer = "Unassigned"
        for assign in dispatch_assignments:
            if assign.get("incident_id") == incident.get("id"):
                officer = assign.get("unit_id", "Unassigned")
                break
        table.cell(row_idx, 4).text = officer[:15]
        
        status = incident.get("status", "pending")
        table.cell(row_idx, 5).text = status.capitalize()

    # ========== SLIDE 3: OFFICER PERFORMANCE ==========
    slide3 = prs.slides.add_slide(blank_layout)
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Officer Performance"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True

    # Create officer performance table
    perf_rows = min(len(officer_workload) + 1, 10)
    perf_cols = 5
    perf_table_shape = slide3.shapes.add_table(perf_rows, perf_cols, Inches(0.4), Inches(1), Inches(9.2), Inches(6))
    perf_table = perf_table_shape.table
    
    perf_table.columns[0].width = Inches(2)    # Officer Name
    perf_table.columns[1].width = Inches(1.5)  # Total
    perf_table.columns[2].width = Inches(1.5)  # Critical
    perf_table.columns[3].width = Inches(1.8)  # High
    perf_table.columns[4].width = Inches(1.4)  # Avg Time

    # Header
    perf_headers = ["Officer", "Total", "Critical", "High", "Needs Rotation"]
    for col_idx, header in enumerate(perf_headers):
        cell = perf_table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)

    # Officer data
    for row_idx, officer in enumerate(officer_workload[:perf_rows - 1], start=1):
        perf_table.cell(row_idx, 0).text = _safe_ppt_text(officer.get("officer_name", "N/A")[:20])
        perf_table.cell(row_idx, 1).text = str(officer.get("total_incidents", 0))
        perf_table.cell(row_idx, 2).text = str(officer.get("critical_incidents", 0))
        perf_table.cell(row_idx, 3).text = str(officer.get("high_incidents", 0))
        perf_table.cell(row_idx, 4).text = "⚠️ YES" if officer.get("needs_rotation") else "No"

    # ========== SLIDE 4: UNRESOLVED/ESCALATED CASES ==========
    slide4 = prs.slides.add_slide(blank_layout)
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Unresolved & Escalated Cases"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True

    unresolved_incidents = [inc for inc in incidents if inc.get("status") != "resolved"]
    
    if unresolved_incidents:
        text_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for idx, incident in enumerate(unresolved_incidents[:8]):  # Limit to 8 for readability
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = (
                f"• ID: {incident.get('id', 'N/A')} | "
                f"Type: {incident.get('type', 'N/A')} | "
                f"Severity: {incident.get('severity', 'N/A')} | "
                f"Status: {incident.get('status', 'N/A')}"
            )
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(11)
    else:
        text_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
        text_frame = text_box.text_frame
        text_frame.text = "✓ All incidents have been resolved. No escalations required."
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)

    # ========== SLIDE 5: ZONE HEATMAP / INCIDENT DISTRIBUTION ==========
    slide5 = prs.slides.add_slide(blank_layout)
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Incident Distribution by Zone"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True

    # Group incidents by location/zone
    zone_incidents = {}
    for incident in incidents:
        zone = incident.get("location", "Unknown")
        zone_incidents[zone] = zone_incidents.get(zone, 0) + 1
    
    if zone_incidents:
        text_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for idx, (zone, count) in enumerate(sorted(zone_incidents.items(), key=lambda x: x[1], reverse=True)[:12]):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            bar = "█" * min(count, 20)
            p.text = f"{zone[:30]:.<30} {bar} ({count})"
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.name = "Courier New"
    else:
        text_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.text = "No incident data available."
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@app.post("/api/shift/report")
async def generate_shift_report(
    shift_id: int = Query(..., description="Shift ID to generate report for"),
    current_user: dict = Depends(require_role("police_supervisor")),
    db: Session = Depends(get_db),
):
    """Generate a comprehensive shift report as PPTX with 5 slides.
    
    Pulls data from:
    - Shift (supervisor, district, times)
    - ShiftAttendance (officers on duty)
    - DispatchLog (immutable dispatch audit trail)
    - OfficerIncidentCount (officer workload)
    """
    try:
        # Get shift data
        shift = db.query(Shift).filter(Shift.id == shift_id).first()
        if not shift:
            raise HTTPException(status_code=404, detail="Shift not found")

        # Verify the supervisor owns this shift
        if shift.supervisor_id != current_user.get("username"):
            raise HTTPException(status_code=403, detail="You do not have permission to access this shift's report")

        # Get all incidents for this shift (notifications within shift time window)
        incidents_query = db.query(Notification).filter(
            Notification.created_at >= shift.start_time,
            Notification.created_at <= (shift.end_time or datetime.now(UTC))
        ).all()

        incidents = [
            {
                "id": inc.id,
                "type": inc.type,
                "title": inc.title,
                "message": inc.message,
                "severity": inc.type.split("_")[-1] if "_" in inc.type else "medium",
                "location": inc.message[:50] if inc.message else "Unknown",
                "status": "resolved" if inc.is_read else "pending",
                "timestamp": inc.created_at.isoformat() if inc.created_at else "",
            }
            for inc in incidents_query
        ]

        # Get dispatch records from immutable dispatch log for this shift time window
        dispatch_log_entries = db.query(DispatchLog).filter(
            DispatchLog.assigned_at >= shift.start_time,
            DispatchLog.assigned_at <= (shift.end_time or datetime.now(UTC)),
            DispatchLog.district_id == shift.district_id
        ).all()

        dispatch_data = [
            {
                "incident_id": dl.incident_id,
                "unit_id": dl.officer_id,
                "assigned_by": dl.assigned_by,
                "assigned_at": dl.assigned_at.isoformat() if dl.assigned_at else "",
                "status": dl.status,
                "response_time_s": 0,  # Can be calculated if timestamp data is available
            }
            for dl in dispatch_log_entries
        ]

        # Get officer workload for this shift
        officer_workload = db.query(OfficerIncidentCount).filter(
            OfficerIncidentCount.shift_id == shift_id
        ).all()

        workload_data = [
            {
                "officer_username": ow.officer_username,
                "officer_name": ow.officer_name,
                "total_incidents": ow.incident_count_critical + ow.incident_count_high + 
                                   ow.incident_count_medium + ow.incident_count_low,
                "critical_incidents": ow.incident_count_critical,
                "high_incidents": ow.incident_count_high,
                "medium_incidents": ow.incident_count_medium,
                "low_incidents": ow.incident_count_low,
                "needs_rotation": ow.needs_rotation,
            }
            for ow in officer_workload
        ]

        # Prepare shift data dict
        shift_dict = {
            "shift_id": shift.id,
            "district_id": shift.district_id,
            "supervisor_name": shift.supervisor_name,
            "supervisor_id": shift.supervisor_id,
            "start_time": shift.start_time,
            "end_time": shift.end_time,
            "status": shift.status,
            "officers_on_duty": shift.officers_on_duty,
            "incidents_count": len(incidents),
        }

        # Generate PPTX
        output = generate_detailed_shift_report_pptx(
            shift_id=shift_id,
            shift_data=shift_dict,
            incidents=incidents,
            officer_workload=workload_data,
            dispatch_assignments=dispatch_data,
        )

        filename = f"ShiftReport_{shift.district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"

        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"',
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
                "Expires": "0",
            },
        )
    except HTTPException:
        raise
    except Exception as exc:
        logger.error(f"Shift report generation failed: {exc}")
        raise HTTPException(
            status_code=500,
            detail=f"Report generation failed: {str(exc)}",
        )


@app.get("/police/shift/status")
async def get_shift_status(current_user: dict = Depends(require_role("police_supervisor"))):
    """Get the current active shift status for the supervisor."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            return {
                "status": "no_active_shift",
                "shift": None,
                "message": "No active shift found. Create a new shift to begin.",
            }
        
        attendance_records = (
            session.query(ShiftAttendance)
            .filter(ShiftAttendance.shift_id == active_shift.id)
            .all()
        )
        
        officers_on_duty = [
            {
                "officer_username": record.officer_username,
                "officer_name": record.officer_name,
                "clock_in_time": record.clock_in_time.isoformat() if record.clock_in_time else None,
                "clock_out_time": record.clock_out_time.isoformat() if record.clock_out_time else None,
                "status": record.status,
            }
            for record in attendance_records
            if record.status == "present"
        ]
        
        return {
            "status": "active",
            "shift": {
                "id": active_shift.id,
                "district_id": active_shift.district_id,
                "supervisor_id": active_shift.supervisor_id,
                "supervisor_name": active_shift.supervisor_name,
                "start_time": active_shift.start_time.isoformat(),
                "incidents_count": active_shift.incidents_count,
                "officers_on_duty": len(officers_on_duty),
                "notes": active_shift.notes,
            },
            "officers": officers_on_duty,
            "updated_at": datetime.now(UTC).isoformat(),
        }
    except Exception as exc:
        logger.error(f"Failed to get shift status: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to retrieve shift status",
        )
    finally:
        session.close()


@app.post("/police/shift/attendance")
async def mark_shift_attendance(
    request: ShiftAttendanceRequest,
    current_user: dict = Depends(require_role("police_supervisor")),
):
    """Mark officer attendance for the current shift."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="No active shift found",
            )
        
        for officer_data in request.officers:
            officer_username = officer_data.get("officer_username")
            officer_name = officer_data.get("officer_name")
            status_val = officer_data.get("status", "present")
            
            existing_record = (
                session.query(ShiftAttendance)
                .filter(
                    ShiftAttendance.shift_id == active_shift.id,
                    ShiftAttendance.officer_username == officer_username,
                )
                .first()
            )
            
            if existing_record:
                existing_record.status = status_val
                if status_val == "absent":
                    existing_record.clock_out_time = datetime.now(UTC)
            else:
                attendance_record = ShiftAttendance(
                    shift_id=active_shift.id,
                    officer_username=officer_username,
                    officer_name=officer_name,
                    status=status_val,
                    clock_in_time=datetime.now(UTC) if status_val == "present" else None,
                )
                session.add(attendance_record)
        
        session.commit()
        
        updated_present_count = (
            session.query(ShiftAttendance)
            .filter(
                ShiftAttendance.shift_id == active_shift.id,
                ShiftAttendance.status == "present",
            )
            .count()
        )
        active_shift.officers_on_duty = updated_present_count
        session.commit()
        
        return {
            "message": "Attendance marked successfully",
            "shift_id": active_shift.id,
            "officers_marked": len(request.officers),
            "officers_present": updated_present_count,
            "updated_at": datetime.now(UTC).isoformat(),
        }
    except HTTPException:
        session.rollback()
        raise
    except Exception as exc:
        session.rollback()
        logger.error(f"Failed to mark attendance: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to mark attendance",
        )
    finally:
        session.close()


@app.post("/api/shifts/create")
async def create_shift(current_user: dict = Depends(require_role("police_supervisor"))):
    """Create a new active shift for the current supervisor if one does not exist."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")

    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )

    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        if active_shift:
            return {
                "message": "Active shift already exists",
                "shift_id": active_shift.id,
                "status": "active",
            }

        user = get_user_by_username(session, username)
        supervisor_name = (
            getattr(user, "full_name", None)
            or format_officer_name(username)
        )

        shift = Shift(
            district_id=district_id,
            supervisor_id=username,
            supervisor_name=supervisor_name,
            status="active",
            start_time=datetime.now(UTC),
            officers_on_duty=0,
            incidents_count=0,
        )
        session.add(shift)
        session.commit()
        session.refresh(shift)

        return {
            "message": "Shift created successfully",
            "shift_id": shift.id,
            "status": shift.status,
            "start_time": shift.start_time.isoformat() if shift.start_time else None,
        }
    except HTTPException:
        session.rollback()
        raise
    except Exception as exc:
        session.rollback()
        logger.error(f"Failed to create shift: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to create shift",
        )
    finally:
        session.close()


@app.post("/police/shift/end")
async def end_shift(
    request: ShiftEndRequest,
    current_user: dict = Depends(require_role("police_supervisor")),
):
    """End the current shift, optionally generating and exporting a PPTX report."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="No active shift found to end",
            )
        
        active_shift.end_time = datetime.now(UTC)
        active_shift.status = "completed"
        active_shift.notes = request.notes
        session.commit()
        
        shift_end_result = {
            "message": "Shift ended successfully",
            "shift_id": active_shift.id,
            "start_time": active_shift.start_time.isoformat(),
            "end_time": active_shift.end_time.isoformat(),
            "status": "completed",
            "updated_at": datetime.now(UTC).isoformat(),
        }
        
        if request.export_pptx:
            try:
                incidents = _load_police_incidents(district_id)
                district_summary = _build_district_summary(incidents, district_id)
                ml_predictions = _predict_police_hotspots(district_id, incidents)
                output = generate_shift_pptx(district_id, username, incidents, district_summary, ml_predictions)
                
                filename = f"ShiftReport_District{district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"
                shift_end_result["pptx_filename"] = filename
                shift_end_result["export_status"] = "success"
                
                session.close()
                
                return StreamingResponse(
                    output,
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'},
                )
            except Exception as export_exc:
                logger.warning(f"PPTX export during shift end failed: {export_exc}")
                shift_end_result["export_status"] = "failed"
                shift_end_result["export_error"] = str(export_exc)
                return shift_end_result
        else:
            return shift_end_result
    except HTTPException:
        session.rollback()
        raise
    except Exception as exc:
        session.rollback()
        logger.error(f"Failed to end shift: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to end shift",
        )
    finally:
        session.close()


@app.get("/police/alerts/list")
async def get_alerts_list(current_user: dict = Depends(require_role("police_supervisor"))):
    """Get unread alerts for the supervisor's district."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    alerts = get_unread_alerts(district_id)
    return AlertListResponse(
        alerts=alerts,
        unread_count=len(alerts),
    )


@app.get("/police/alerts/stream")
async def alerts_stream(current_user: dict = Depends(require_role("police_supervisor"))):
    """Server-Sent Events stream for real-time police alerts."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    async def event_generator():
        """Generate SSE events for alerts."""
        # Send initial alert to inform client the stream is connected
        yield f"data: {json.dumps({'type': 'connected', 'message': 'Alert stream connected'})}\n\n"
        
        # Send all existing unread alerts
        existing_alerts = get_unread_alerts(district_id)
        for alert in existing_alerts:
            yield f"data: {json.dumps({'type': 'alert', 'data': alert.model_dump()})}\n\n"
            await asyncio.sleep(0.01)  # Small delay to prevent overwhelming client
        
        # Keep connection open for new alerts (heartbeat)
        last_check = time.time()
        last_alert_count = len(existing_alerts)
        
        while True:
            try:
                await asyncio.sleep(1)
                
                current_alerts = get_unread_alerts(district_id)
                current_count = len(current_alerts)
                
                # Check for new alerts since last check
                if current_count > last_alert_count:
                    new_alerts = current_alerts[last_alert_count:]
                    for alert in new_alerts:
                        yield f"data: {json.dumps({'type': 'alert', 'data': alert.model_dump()})}\n\n"
                    last_alert_count = current_count
                
                # Heartbeat every 30 seconds
                if time.time() - last_check > 30:
                    yield f": heartbeat\n\n"
                    last_check = time.time()
            except asyncio.CancelledError:
                logger.info(f"Alert stream closed for district {district_id}")
                break
            except Exception as exc:
                logger.error(f"Error in alert stream: {exc}")
                break
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.get("/api/police/officer-workload", response_model=OfficerWorkloadResponse)
async def get_officer_workload(
    district_id: str = Query(..., description="District ID"),
    shift_id: Optional[int] = Query(None, description="Optional shift ID (uses current active shift if not provided)"),
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_db),
):
    """Get officer incident counts and rotation status for current shift.
    
    Returns officer workload data with rotation alerts for officers
    handling 3+ critical incidents.
    """
    # Determine shift to use
    if not shift_id:
        # Get current active shift for this district
        shift = db.query(Shift).filter(
            Shift.district_id == district_id,
            Shift.status == "active"
        ).order_by(Shift.start_time.desc()).first()
        
        if not shift:
            raise HTTPException(
                status_code=404,
                detail=f"No active shift found for district {district_id}"
            )
        shift_id = shift.id
    else:
        shift = db.query(Shift).filter(Shift.id == shift_id).first()
        if not shift:
            raise HTTPException(status_code=404, detail="Shift not found")
    
    # Get all officer attendance records for this shift
    attendances = db.query(ShiftAttendance).filter(
        ShiftAttendance.shift_id == shift_id
    ).all()
    
    officers_data = []
    officers_needing_rotation = []
    
    for attendance in attendances:
        # Query OfficerIncidentCount for this officer in this shift
        workload = db.query(OfficerIncidentCount).filter(
            OfficerIncidentCount.shift_id == shift_id,
            OfficerIncidentCount.officer_username == attendance.officer_username
        ).first()
        
        if workload:
            officer_info = OfficerWorkloadData(
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name,
                total_incidents=workload.incident_count_critical + workload.incident_count_high + 
                                workload.incident_count_medium + workload.incident_count_low,
                critical_incidents=workload.incident_count_critical,
                high_incidents=workload.incident_count_high,
                medium_incidents=workload.incident_count_medium,
                low_incidents=workload.incident_count_low,
                needs_rotation=workload.needs_rotation
            )
        else:
            # Create new workload entry if it doesn't exist
            workload = OfficerIncidentCount(
                shift_id=shift_id,
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name
            )
            db.add(workload)
            officer_info = OfficerWorkloadData(
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name,
                total_incidents=0,
                critical_incidents=0,
                high_incidents=0,
                medium_incidents=0,
                low_incidents=0,
                needs_rotation=False
            )
        
        officers_data.append(officer_info)
        if officer_info.needs_rotation:
            officers_needing_rotation.append(attendance.officer_username)
    
    db.commit()
    
    return OfficerWorkloadResponse(
        shift_id=shift_id,
        officers=officers_data,
        officers_needing_rotation=officers_needing_rotation
    )


@app.post("/api/police/incident-handled")
async def log_incident_handled(
    incident_id: str = Query(...),
    officer_username: str = Query(...),
    severity: str = Query(..., pattern="^(critical|high|medium|low)$"),
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_db),
):
    """Log that an officer has handled an incident (increments workload counter).
    
    Called when an officer is dispatched/finishes handling an incident.
    Tracks incident counts and updates rotation status if 3+ critical incidents.
    """
    # Get current active shift for officer
    shift = db.query(Shift).filter(
        Shift.status == "active"
    ).order_by(Shift.start_time.desc()).first()
    
    if not shift:
        raise HTTPException(
            status_code=404,
            detail="No active shift found"
        )
    
    # Get or create workload record for this officer
    workload = db.query(OfficerIncidentCount).filter(
        OfficerIncidentCount.shift_id == shift.id,
        OfficerIncidentCount.officer_username == officer_username
    ).first()
    
    if not workload:
        # Get officer name from attendance record
        attendance = db.query(ShiftAttendance).filter(
            ShiftAttendance.shift_id == shift.id,
            ShiftAttendance.officer_username == officer_username
        ).first()
        
        officer_name = attendance.officer_name if attendance else officer_username
        
        workload = OfficerIncidentCount(
            shift_id=shift.id,
            officer_username=officer_username,
            officer_name=officer_name
        )
        db.add(workload)
    
    # Increment appropriate counter
    if severity == "critical":
        workload.incident_count_critical += 1
    elif severity == "high":
        workload.incident_count_high += 1
    elif severity == "medium":
        workload.incident_count_medium += 1
    else:
        workload.incident_count_low += 1
    
    # Check if needs rotation (3+ critical incidents)
    if workload.incident_count_critical >= 3:
        workload.needs_rotation = True
    
    db.commit()
    
    return {
        "status": "logged",
        "officer_username": officer_username,
        "incident_id": incident_id,
        "severity": severity,
        "total_critical_incidents": workload.incident_count_critical,
        "needs_rotation": workload.needs_rotation
    }


@app.post("/api/incident/resolve")
async def resolve_incident(
    request: IncidentResolveRequest,
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_db),
):
    """Mark incident resolved and persist supervised feedback record for ML training."""
    resolved_at = request.resolved_at.astimezone(UTC) if request.resolved_at else datetime.now(UTC)
    severity_value = _normalize_severity(request.severity)
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident resolution",
        )

    feedback = MLFeedback(
        district_id=district_id,
        incident_type=request.incident_type.strip(),
        zone=request.zone.strip(),
        time_of_day=resolved_at.hour,
        day_of_week=resolved_at.weekday(),
        response_time_minutes=float(request.response_time_minutes),
        severity=severity_value,
        outcome=request.outcome.strip(),
        created_at=resolved_at,
    )
    db.add(feedback)

    # Close any active assignment linked to the resolved incident.
    assignment = (
        db.query(PoliceDispatchAssignment)
        .filter(PoliceDispatchAssignment.incident_id == request.incident_id)
        .first()
    )
    returned_officer_id = None
    if assignment:
        assignment.status = "resolved"
        returned_officer_id = assignment.unit_id

    if returned_officer_id:
        officer_status = (
            db.query(OfficerDispatchStatus)
            .filter(OfficerDispatchStatus.officer_id == returned_officer_id)
            .first()
        )
        if officer_status:
            officer_status.status = "available"
            officer_status.assigned_incident_id = None
            officer_status.updated_at = resolved_at

        dispatch_record = (
            db.query(DispatchLog)
            .filter(
                DispatchLog.incident_id == request.incident_id,
                DispatchLog.officer_id == returned_officer_id,
            )
            .order_by(DispatchLog.assigned_at.desc())
            .first()
        )
        if dispatch_record:
            dispatch_record.status = "returned"

    db.commit()
    db.refresh(feedback)

    await emit_incident_updated(
        sio,
        district_id,
        {
            "id": request.incident_id,
            "severity": severity_value,
            "status": "resolved",
            "resolved_at": resolved_at.isoformat(),
        },
        update_type="resolved",
        actor=current_user.get("username", "Unknown Supervisor"),
    )
    if returned_officer_id:
        await emit_officer_status_changed(
            sio,
            district_id,
            {
                "id": returned_officer_id,
                "badge": returned_officer_id,
                "status": "available",
                "district_id": district_id,
                "incident_id": None,
            },
            actor=current_user.get("username", "Unknown Supervisor"),
        )

    # Fire a real-time alert for the resolution
    add_alert(
        district_id=district_id,
        severity="low",
        message=f"Incident {request.incident_id} resolved"
        + (f" — {returned_officer_id} returned to service" if returned_officer_id else ""),
        incident_id=request.incident_id,
    )

    return {
        "status": "resolved",
        "incident_id": request.incident_id,
        "feedback_id": feedback.id,
        "saved_training_record": True,
        "resolved_at": resolved_at.isoformat(),
    }


@app.post("/api/ml/retrain")
async def retrain_ml_model(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db),
):
    """Export feedback rows to CSV and trigger existing SVR retraining pipeline."""
    feedback_rows = db.query(MLFeedback).order_by(MLFeedback.created_at.asc()).all()
    if not feedback_rows:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No ML feedback available for retraining",
        )

    export_dir = os.path.join("data", "exports")
    os.makedirs(export_dir, exist_ok=True)
    export_name = f"ml_feedback_{datetime.now(UTC).strftime('%Y%m%d_%H%M%S')}.csv"
    export_path = os.path.join(export_dir, export_name)

    with open(export_path, "w", newline="", encoding="utf-8") as csv_file:
        writer = csv.DictWriter(
            csv_file,
            fieldnames=[
                "incident_type",
                "zone",
                "time_of_day",
                "day_of_week",
                "response_time_minutes",
                "severity",
                "outcome",
                "created_at",
            ],
        )
        writer.writeheader()
        for row in feedback_rows:
            writer.writerow(
                {
                    "incident_type": row.incident_type,
                    "zone": row.zone,
                    "time_of_day": row.time_of_day,
                    "day_of_week": row.day_of_week,
                    "response_time_minutes": row.response_time_minutes,
                    "severity": row.severity,
                    "outcome": row.outcome,
                    "created_at": row.created_at.isoformat() if row.created_at else None,
                }
            )

    try:
        from .svr_model import train_svr

        train_svr()
    except Exception as exc:
        logger.error("ML retraining failed: %s", exc)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="ML retraining failed",
        )

    retrain_audit = MLRetrainAudit(
        retrained_at=datetime.now(UTC),
        retrained_by=getattr(current_user, "username", "admin"),
        feedback_rows=len(feedback_rows),
        export_path=export_path,
    )
    db.add(retrain_audit)
    db.commit()

    return {
        "status": "retrained",
        "feedback_rows": len(feedback_rows),
        "export_csv": export_path,
        "last_retrained": retrain_audit.retrained_at.isoformat(),
    }


@app.post("/auth/login", response_model=RoleToken)
async def role_login(
    login_data: RoleLoginRequest,
    db: Session = Depends(get_db),
    *,
    response: Response,
):
    """Login with role and return role-aware JWT. Auto-detects actual user role."""
    user = authenticate_user(db, login_data.username, login_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )

    user_department = (getattr(user, "department", None) or "").strip().lower()

    # Determine actual allowed role based on user's account
    if user.is_admin:
        actual_role = UserRole.admin
    elif user_department == "police":
        actual_role = UserRole.police_supervisor
    elif user_department == "logistics":
        actual_role = UserRole.logistics_manager
    else:
        actual_role = UserRole.user

    # Use actual role instead of requested role for normal users
    # This allows users to select "Normal User" even if system auto-detects their role
    role = actual_role
    district_id = login_data.district_id if login_data.district_id else None
    fleet_zone = login_data.fleet_zone if login_data.fleet_zone else None

    # Validate required fields for specific roles
    if role == UserRole.police_supervisor and not district_id:
        # Default to district_1 for police supervisors if not provided
        district_id = "district_1"

    if role == UserRole.logistics_manager and not fleet_zone:
        # Default to zone_default for logistics managers if not provided
        fleet_zone = "zone_default"

    access_token = create_role_access_token(
        username=user.username,
        role=role,
        district_id=district_id,
        fleet_zone=fleet_zone,
        expires_delta=timedelta(minutes=30 * 24 * 60),
    )
    # Set secure=False for localhost/HTTP development, True for production/HTTPS
    is_secure = os.getenv("SECURE_COOKIES", "false").lower() == "true"
    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=is_secure,
        samesite="Lax",  # Use Lax for both secure and non-secure (None requires Secure=True)
        max_age=30 * 24 * 60 * 60,
        path="/",
    )
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "role": role,
        "district_id": district_id,
        "fleet_zone": fleet_zone,
    }


@app.get("/api/auth/me", response_model=UserResponse)
async def get_current_user_info(current_user: User = Depends(get_current_active_user)):
    """Get current user information."""
    return UserResponse.model_validate(current_user)


# ============================================================================
# SAVED ROUTES
# ============================================================================

@app.post("/api/saved-routes")
@handle_db_errors
async def create_saved_route(
    route_data: SavedRouteCreate,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Save a route for the current user."""
    if not current_user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Please login to save routes."
        )
    
    origin_str = json.dumps(route_data.origin) if isinstance(route_data.origin, dict) else route_data.origin
    dest_str = json.dumps(route_data.destination) if isinstance(route_data.destination, dict) else route_data.destination
    
    saved_route = SavedRoute(
        user_id=current_user.id,
        route_name=route_data.route_name,
        origin=origin_str,
        destination=dest_str,
        route_preferences=route_data.route_preferences,
        is_favorite=False,
        share_token=secrets.token_urlsafe(16)
    )
    db.add(saved_route)
    db.commit()
    db.refresh(saved_route)
    return saved_route


@app.get("/api/saved-routes")
@handle_db_errors
async def get_saved_routes(
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db),
    favorites_only: bool = Query(False)
):
    """Get saved routes for current user."""
    if not current_user:
        return []
    query = db.query(SavedRoute).filter(SavedRoute.user_id == current_user.id)
    if favorites_only:
        query = query.filter(SavedRoute.is_favorite == True)
    routes = query.order_by(SavedRoute.last_used.desc()).all()
    return routes


@app.put("/api/saved-routes/{route_id}/favorite")
@handle_db_errors
async def toggle_favorite(
    route_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Toggle favorite status of a saved route."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required for this feature")
    route = db.query(SavedRoute).filter(
        SavedRoute.id == route_id,
        SavedRoute.user_id == current_user.id
    ).first()
    if not route:
        raise HTTPException(status_code=404, detail="Route not found")
    route.is_favorite = not route.is_favorite
    db.commit()
    return {"is_favorite": route.is_favorite}


@app.delete("/api/saved-routes/{route_id}")
@handle_db_errors
async def delete_saved_route(
    route_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Delete a saved route."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required for this feature")
    route = db.query(SavedRoute).filter(
        SavedRoute.id == route_id,
        SavedRoute.user_id == current_user.id
    ).first()
    if not route:
        raise HTTPException(status_code=404, detail="Route not found")
    db.delete(route)
    db.commit()
    return {"message": "Route deleted"}


@app.get("/api/share-route/{share_token}")
@handle_db_errors
async def get_shared_route(share_token: str, db: Session = Depends(get_db)):
    """Get a shared route by token."""
    route = db.query(SavedRoute).filter(SavedRoute.share_token == share_token).first()
    if not route:
        raise HTTPException(status_code=404, detail="Shared route not found")
    return route


# ============================================================================
# ADVANCED ANALYTICS
# ============================================================================

@app.get("/api/analytics/peak-hours/{route_id}")
async def get_peak_hours(
    route_id: str,
    days: int = Query(30, ge=1, le=365),
    db: Session = Depends(get_db)
):
    """Get peak hours analysis for a route."""
    return get_peak_hours_analysis(db, route_id, days)


@app.get("/api/analytics/day-of-week/{route_id}")
async def get_day_analysis(
    route_id: str,
    days: int = Query(90, ge=1, le=365),
    db: Session = Depends(get_db)
):
    """Get day of week analysis."""
    return get_day_of_week_analysis(db, route_id, days)


@app.get("/api/analytics/seasonal/{route_id}")
async def get_seasonal_analysis(
    route_id: str,
    months: int = Query(12, ge=1, le=24),
    db: Session = Depends(get_db)
):
    """Get seasonal trends."""
    return get_seasonal_trends(db, route_id, months)


@app.get("/api/analytics/reliability/{route_id}")
async def get_reliability(
    route_id: str,
    days: int = Query(30, ge=1, le=365),
    db: Session = Depends(get_db)
):
    """Get route reliability score."""
    return calculate_route_reliability(db, route_id, days)


@app.get("/api/analytics/predict/{route_id}")
async def get_prediction(
    route_id: str,
    hours_ahead: int = Query(24, ge=1, le=168),
    db: Session = Depends(get_db)
):
    """Predict future congestion."""
    return predict_future_congestion(db, route_id, hours_ahead)


# ============================================================================
# EXPORT & REPORTING
# ============================================================================

@app.get("/api/export/csv/{route_id}")
async def export_csv(
    route_id: str,
    db: Session = Depends(get_db)
):
    """Export route data to CSV."""
    csv_content = export_to_csv(db, route_id)
    return StreamingResponse(
        io.BytesIO(csv_content.encode('utf-8')),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename=route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.csv"}
    )


@app.get("/api/export/excel/{route_id}")
async def export_excel(
    route_id: str,
    db: Session = Depends(get_db)
):
    """Export route data to Excel."""
    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        export_to_excel(db, route_id, tmp.name)
        return FileResponse(
            tmp.name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.xlsx"
        )


@app.get("/api/export/pdf/{route_id}")
async def export_pdf(
    route_id: str,
    db: Session = Depends(get_db)
):
    """Export route data to PDF."""
    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
        export_to_pdf(db, route_id, tmp.name)
        return FileResponse(
            tmp.name,
            media_type="application/pdf",
            filename=f"route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.pdf"
        )


# ============================================================================
# NOTIFICATIONS
# ============================================================================

@app.get("/api/notifications")
@handle_db_errors
async def get_notifications(
    unread_only: bool = Query(False),
    limit: int = Query(50, ge=1, le=100),
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Get user notifications."""
    if not current_user:
        return []
    return get_user_notifications(db, current_user.id, unread_only, limit)


@app.put("/api/notifications/{notification_id}/read")
@handle_db_errors
async def mark_read(
    notification_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Mark notification as read."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required")
    success = mark_notification_read(db, notification_id, current_user.id)
    if not success:
        raise HTTPException(status_code=404, detail="Notification not found")
    return {"message": "Notification marked as read"}


@app.post("/api/notifications/check-alerts")
@handle_db_errors
async def check_alerts(
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Check for traffic alerts on saved routes."""
    if not current_user:
        return {"alerts": 0, "notifications": []}
    alerts = check_traffic_alerts(db, current_user.id)
    return {"alerts": len(alerts), "notifications": alerts}


# ============================================================================
# REAL-TIME FEATURES
# ============================================================================

@app.get("/api/realtime/incidents")
async def get_incidents(
    lat: float = Query(..., description="Latitude"),
    lon: float = Query(..., description="Longitude"),
    radius: int = Query(5000, ge=100, le=50000)
):
    """Get traffic incidents near a location."""
    # Validate coordinates
    if not (-90 <= lat <= 90):
        raise HTTPException(status_code=400, detail="Invalid latitude")
    if not (-180 <= lon <= 180):
        raise HTTPException(status_code=400, detail="Invalid longitude")
    
    return get_traffic_incidents(lat, lon, radius)


@app.post("/api/realtime/monitor/{route_id}")
async def monitor_route(
    route_id: str,
    background_tasks: BackgroundTasks,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Monitor route for changes."""
    change = monitor_route_changes(db, route_id)
    if change:
        if current_user:
            from .notifications import create_notification
            create_notification(
                db, current_user.id, 'traffic_alert',
                f"Route Change: {route_id}",
                f"Route travel time changed by {change['change_percent']}%",
                route_id
            )
    return change or {"message": "No significant changes detected"}


# ============================================================================
# ROUTE RATINGS & SOCIAL
# ============================================================================

@app.post("/api/ratings")
@handle_db_errors
async def create_rating(
    rating_data: RouteRatingCreate,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_db)
):
    """Rate a route."""
    user_id = current_user.id if current_user else None
    rating = RouteRating(
        user_id=user_id,
        route_id=rating_data.route_id,
        rating=rating_data.rating,
        review=rating_data.review
    )
    db.add(rating)
    db.commit()
    db.refresh(rating)
    return rating


@app.get("/api/ratings/{route_id}")
@handle_db_errors
async def get_ratings(route_id: str, db: Session = Depends(get_db)):
    """Get ratings for a route."""
    ratings = db.query(RouteRating).filter(RouteRating.route_id == route_id).all()
    if not ratings:
        return {"average_rating": 0, "count": 0, "ratings": []}
    
    avg_rating = sum(r.rating for r in ratings) / len(ratings)
    return {
        "average_rating": round(avg_rating, 2),
        "count": len(ratings),
        "ratings": ratings
    }


# ============================================================================
# ADMIN DASHBOARD
# ============================================================================

@app.get("/api/admin/stats")
@handle_db_errors
async def get_admin_stats(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db)
):
    """Get admin statistics (admin only)."""
    total_users = db.query(User).count()
    total_routes = db.query(AnalysisResult).count()
    total_saved_routes = db.query(SavedRoute).count()
    total_ratings = db.query(RouteRating).count()
    
    recent_routes = db.query(AnalysisResult).order_by(AnalysisResult.timestamp.desc()).limit(10).all()
    recent_users = db.query(User).order_by(User.created_at.desc()).limit(5).all()
    
    return {
        "total_users": total_users,
        "total_route_analyses": total_routes,
        "total_saved_routes": total_saved_routes,
        "total_ratings": total_ratings,
        "cache_stats": get_cache_stats(),
        "recent_activity": {
            "routes": len(recent_routes),
            "new_users": len(recent_users)
        }
    }


@app.get("/api/admin/route-analysis")
async def get_all_route_analyses(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db),
    filter_period: Optional[str] = Query(None, alias="filter"),
    skip: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=1000)
):
    """Get all route analyses with optional filtering (admin only)."""
    import json
    from datetime import datetime, timedelta, UTC
    
    try:
        # Build query
        query = db.query(AnalysisResult)
        
        # Apply time filter if specified
        if filter_period:
            now = datetime.now(UTC)
            if filter_period == "today":
                start_date = now.replace(hour=0, minute=0, second=0, microsecond=0)
                query = query.filter(AnalysisResult.timestamp >= start_date)
            elif filter_period == "week":
                start_date = now - timedelta(days=7)
                query = query.filter(AnalysisResult.timestamp >= start_date)
            elif filter_period == "month":
                start_date = now - timedelta(days=30)
                query = query.filter(AnalysisResult.timestamp >= start_date)
        
        # Get total count before pagination
        total_count = query.count()
        
        # Apply pagination and ordering
        routes = query.order_by(AnalysisResult.timestamp.desc()).offset(skip).limit(limit).all()
        
        # Format response
        route_data = []
        for r in routes:
            try:
                origin = json.loads(r.origin) if isinstance(r.origin, str) and r.origin.startswith('{') else {"name": str(r.origin) if r.origin else ""}
                dest = json.loads(r.destination) if isinstance(r.destination, str) and r.destination.startswith('{') else {"name": str(r.destination) if r.destination else ""}
            except:
                origin = {"name": str(r.origin) if r.origin else ""}
                dest = {"name": str(r.destination) if r.destination else ""}
            
            route_name = f"{origin.get('name', '')} → {dest.get('name', '')}"
            
            delay_val = r.delay_s
            if delay_val is None or delay_val == 0:
                if r.travel_time_s and r.no_traffic_s:
                    delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                else:
                    delay_val = 0
            
            route_data.append({
                "id": r.id,
                "route": route_name,
                "route_id": r.route_id,
                "travel_time_s": r.travel_time_s,
                "delay_s": delay_val,
                "length_m": r.length_m,
                "calculated_cost": r.calculated_cost,
                "ml_predicted": r.ml_predicted,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "origin": origin,
                "destination": dest
            })
        
        # Calculate statistics from filtered results
        all_routes_for_stats = query.all()
        
        if all_routes_for_stats:
            travel_times = [r.travel_time_s for r in all_routes_for_stats if r.travel_time_s is not None]
            delays = []
            for r in all_routes_for_stats:
                delay_val = r.delay_s
                if delay_val is None or delay_val == 0:
                    if r.travel_time_s and r.no_traffic_s:
                        delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                    else:
                        delay_val = 0
                if delay_val > 0:
                    delays.append(delay_val)
            costs = [r.calculated_cost for r in all_routes_for_stats if r.calculated_cost is not None]
            
            stats = {
                "total": total_count,
                "avg_travel_time": sum(travel_times) / len(travel_times) if travel_times else 0,
                "avg_delay": sum(delays) / len(delays) if delays else 0,
                "avg_cost": sum(costs) / len(costs) if costs else 0
            }
        else:
            stats = {
                "total": 0,
                "avg_travel_time": 0,
                "avg_delay": 0,
                "avg_cost": 0
            }
        
        return {
            "routes": route_data,
            "stats": stats,
            "pagination": {
                "skip": skip,
                "limit": limit,
                "total": total_count
            }
        }
        
    except Exception as e:
        logger.error(f"Error fetching route analyses: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch route analyses: {str(e)}")


@app.get("/api/admin/users")
async def get_all_users(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db),
    skip: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=1000)
):
    """Get users filtered by department (admin only - see only users in their department)."""
    # Get admin's department
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    
    # Build query based on admin's department
    query = db.query(User)
    
    # If admin is super admin (department='admin'), show all users
    # Otherwise, show only users in their own department
    if admin_department != "admin":
        # Filter by department - but first check if there are department-specific admins
        # Police department admins see police users
        # Logistics department admins see logistics users
        query = query.filter(User.department == admin_department)
    # Super admins (department='admin') see everyone
    
    users = query.offset(skip).limit(limit).all()
    return [UserResponse.model_validate(u) for u in users]


@app.put("/api/admin/users/{user_id}/activate")
async def toggle_user_status(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db)
):
    """Activate/deactivate a user (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    user.is_active = not user.is_active
    db.commit()
    return {"is_active": user.is_active, "message": f"User {'activated' if user.is_active else 'deactivated'}"}


@app.put("/api/admin/users/{user_id}/admin")
async def toggle_admin_status(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db)
):
    """Grant/revoke admin privileges (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user.id == current_user.id:
        raise HTTPException(status_code=400, detail="Cannot change your own admin status")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    user.is_admin = not user.is_admin
    db.commit()
    return {"is_admin": user.is_admin, "message": f"Admin privileges {'granted' if user.is_admin else 'revoked'}"}


@app.put("/api/admin/users/{user_id}")
async def update_user(
    user_id: int,
    user_update: UserUpdate,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db)
):
    """Update user details (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    # Update username if provided
    if user_update.username is not None:
        existing_user = db.query(User).filter(User.username == user_update.username, User.id != user_id).first()
        if existing_user:
            raise HTTPException(status_code=400, detail="Username already taken")
        user.username = user_update.username
    
    # Update email if provided
    if user_update.email is not None:
        existing_user = db.query(User).filter(User.email == user_update.email, User.id != user_id).first()
        if existing_user:
            raise HTTPException(status_code=400, detail="Email already taken")
        user.email = user_update.email
    
    # Update full name if provided
    if user_update.full_name is not None:
        user.full_name = user_update.full_name
    
    # Update active status if provided
    if user_update.is_active is not None:
        user.is_active = user_update.is_active
    
    # Update admin status if provided
    if user_update.is_admin is not None:
        if user.id == current_user.id and not user_update.is_admin:
            raise HTTPException(status_code=400, detail="Cannot remove your own admin privileges")
        user.is_admin = user_update.is_admin
    
    # Update password if provided and not empty
    if user_update.password is not None and user_update.password.strip() != "":
        if len(user_update.password) < 8:
            raise HTTPException(status_code=400, detail="Password must be at least 8 characters")
        user.hashed_password = get_password_hash(user_update.password)
    
    try:
        db.commit()
        db.refresh(user)
        return UserResponse.model_validate(user)
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating user: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to update user: {str(e)}")


@app.delete("/api/admin/users/{user_id}")
async def delete_user(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_db)
):
    """Delete a user (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user.id == current_user.id:
        raise HTTPException(status_code=400, detail="Cannot delete your own account")
    
    # Check if admin can delete this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only delete users in your department"
        )
    
    db.delete(user)
    db.commit()
    return {"message": "User deleted successfully"}


@app.get("/api/user/stats")
@handle_db_errors
async def get_user_stats(
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_db)
):
    """Get user-specific statistics."""
    saved_routes_count = db.query(SavedRoute).filter(SavedRoute.user_id == current_user.id).count()
    analyses_count = db.query(AnalysisResult).filter(AnalysisResult.user_id == current_user.id).count()
    ratings_count = db.query(RouteRating).filter(RouteRating.user_id == current_user.id).count()
    
    recent_analyses = db.query(AnalysisResult).filter(
        AnalysisResult.user_id == current_user.id
    ).order_by(AnalysisResult.timestamp.desc()).limit(10).all()
    
    return {
        "saved_routes": saved_routes_count,
        "analyses": analyses_count,
        "ratings": ratings_count,
        "recent_analyses": [
            {
                "route_id": r.route_id,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "travel_time": r.travel_time_s,
                "cost": r.calculated_cost
            }
            for r in recent_analyses
        ]
    }


# ============================================================================
# CACHE MANAGEMENT
# ============================================================================

@app.post("/api/cache/clear")
async def clear_route_cache(
    pattern: Optional[str] = None,
    current_user: User = Depends(get_current_admin_user)
):
    """Clear route cache (admin only)."""
    clear_cache(pattern=pattern)
    return {"message": "Cache cleared"}


@app.get("/api/cache/stats")
async def get_cache_statistics():
    """Get cache statistics."""
    return get_cache_stats()


# ============================================================================
# INTEGRATION ENDPOINTS
# ============================================================================

@app.get("/api/integration/navigation/{route_id}")
@handle_db_errors
async def get_navigation_links(
    route_id: str,
    route_index: int = Query(0),
    db: Session = Depends(get_db)
):
    """Get navigation app links (Google Maps, Waze)."""
    result = db.query(AnalysisResult).filter(
        AnalysisResult.route_id.like(f"{route_id}%")
    ).order_by(AnalysisResult.timestamp.desc()).first()
    
    if not result:
        raise HTTPException(status_code=404, detail="Route not found")
    
    try:
        origin = json.loads(result.origin) if result.origin and result.origin.startswith('{') else {"name": result.origin}
        dest = json.loads(result.destination) if result.destination and result.destination.startswith('{') else {"name": result.destination}
    except:
        origin = {"name": result.origin}
        dest = {"name": result.destination}
    
    origin_lat = origin.get('lat', 0)
    origin_lon = origin.get('lon', 0)
    dest_lat = dest.get('lat', 0)
    dest_lon = dest.get('lon', 0)
    
    google_maps = f"https://www.google.com/maps/dir/{origin_lat},{origin_lon}/{dest_lat},{dest_lon}"
    waze = f"https://waze.com/ul?ll={dest_lat},{dest_lon}&navigate=yes"
    
    return {
        "google_maps": google_maps,
        "waze": waze,
        "origin": origin,
        "destination": dest
    }
