"""
FastAPI backend for traffic route analysis.
Provides endpoints for autocomplete, route analysis, and serving the frontend.
"""

import os
import json
import logging
import re
import csv
import socketio
from dotenv import load_dotenv
from typing import Optional, Union, List
from functools import wraps
from fastapi import FastAPI, HTTPException, Query, Depends, status, BackgroundTasks, Request
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse, FileResponse, Response, RedirectResponse
import io
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordRequestForm
from pydantic import BaseModel, Field, EmailStr, field_validator
from jose.exceptions import ExpiredSignatureError
import joblib
from datetime import datetime, UTC, timedelta
import secrets
import uuid
import traceback
import asyncio
import time
from sqlalchemy.exc import IntegrityError

load_dotenv()

# Import logging and rate limiting
from logging_config import setup_logging, get_logger
from rate_limiter import RateLimitMiddleware

# Setup logging
setup_logging()
logger = get_logger(__name__)

from utils import (
    tomtom_geocode,
    tomtom_autocomplete,
    tomtom_route,
    summarize_route,
    compute_route_cost,
    haversine_m
)
from db import (
    init_db, get_session, save_analysis, AnalysisResult,
    User, SavedRoute, RouteRating, Notification, PoliceDispatchAssignment, OfficerDispatchStatus, DispatchLog, SharedAlert, Shift, ShiftAttendance, OfficerIncidentCount, MLFeedback, MLRetrainAudit
)
from sqlalchemy.orm import Session
from auth import (
    verify_password, get_password_hash, create_access_token,
    get_current_user, get_current_active_user, get_current_admin_user,
    authenticate_user, create_user, get_user_by_username, Token, UserCreate as AuthUserCreate, UserResponse,
    get_optional_user, RoleLoginRequest, RoleToken, UserRole, create_role_access_token, require_role, require_any_role,
    require_police_department_user
)
from analytics import (
    get_peak_hours_analysis, get_day_of_week_analysis,
    get_seasonal_trends, calculate_route_reliability, predict_future_congestion,
    get_traffic_hotspots
)
from export_utils import export_to_csv, export_to_excel, export_to_pdf
from notifications import (
    create_notification, check_traffic_alerts,
    suggest_best_time_to_leave, check_congestion_warnings,
    get_user_notifications, mark_notification_read
)
from cache_utils import cached, clear_cache, get_cache_stats
from realtime_utils import get_traffic_incidents, auto_refresh_route, monitor_route_changes
from dispatch_notifications import send_officer_dispatch_notification

# Initialize FastAPI app
app = FastAPI(
    title="Traffic Route Analysis API",
    description="Real-time traffic congestion analysis with ML predictions",
    version="1.0.0"
)

# Socket.IO server for push updates to the command center frontend.
SOCKETIO_REDIS_URL = os.getenv("SOCKETIO_REDIS_URL", "redis://redis:6379/0")

# Try to use Redis manager for distributed systems, fall back to in-memory for development
try:
    socket_client_manager = socketio.AsyncRedisManager(SOCKETIO_REDIS_URL)
    logger.info(f"✅ Socket.IO using Redis manager: {SOCKETIO_REDIS_URL}")
except Exception as e:
    logger.warning(f"⚠️ Redis connection failed ({e}), using in-memory manager for Socket.IO")
    socket_client_manager = None  # Will use in-memory manager by default

sio = socketio.AsyncServer(
    async_mode="asgi",
    cors_allowed_origins="*",
    client_manager=socket_client_manager,  # None = use in-memory manager
    ping_timeout=60,  # Server expects ping from client within 60 seconds
    ping_interval=25,  # Server sends ping every 25 seconds
    max_http_buffer_size=1000000,  # 1MB max message size
    logger=True,  # Enable Socket.IO logger
    engineio_logger=False,  # Don't spam with engine.io logs
)
app.mount("/socket.io", socketio.ASGIApp(sio, socketio_path="/"))

from socketio_events import (
    register_police_socketio_handlers,
    emit_incident_new,
    emit_incident_updated,
    emit_officer_status_changed,
    emit_officer_dispatched,
)

register_police_socketio_handlers(sio, logger)

# ============================================================================
# REAL-TIME ALERTS INFRASTRUCTURE
# ============================================================================
# In-memory alert store: {district_id: [AlertData, ...]}
_alerts_store = {}
_alert_subscribers = {}  # {district_id: [callback_func, ...]}
_manual_incidents_store: dict[str, list[dict]] = {}


def add_alert(district_id: str, severity: str, message: str, incident_id: Optional[str] = None):
    """Add an alert to the store and notify subscribers."""
    from uuid import uuid4
    
    alert = AlertData(
        alert_id=str(uuid4()),
        severity=severity,
        message=message,
        timestamp=datetime.now(UTC).isoformat(),
        district_id=district_id,
        related_incident_id=incident_id,
    )
    
    if district_id not in _alerts_store:
        _alerts_store[district_id] = []
    
    _alerts_store[district_id].append(alert)
    
    # Keep only last 50 alerts per district to prevent memory bloat
    if len(_alerts_store[district_id]) > 50:
        _alerts_store[district_id] = _alerts_store[district_id][-50:]
    
    logger.info(f"Alert added: {severity} - {message} (District: {district_id})")
    return alert


def get_unread_alerts(district_id: str) -> list[AlertData]:
    """Get unread alerts for a district."""
    return _alerts_store.get(district_id, [])


@app.exception_handler(ExpiredSignatureError)
async def expired_signature_exception_handler(request: Request, exc: ExpiredSignatureError):
    accept_header = (request.headers.get("accept") or "").lower()
    if "text/html" in accept_header:
        return RedirectResponse(url="/auth/login?reason=session_expired", status_code=status.HTTP_307_TEMPORARY_REDIRECT)

    return JSONResponse(
        status_code=status.HTTP_401_UNAUTHORIZED,
        content={"detail": "Session expired. Please log in again."},
    )


@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    detail = exc.detail

    if exc.status_code == status.HTTP_403_FORBIDDEN:
        detail = "Access denied. You do not have permission to view this resource."
    elif exc.status_code == status.HTTP_401_UNAUTHORIZED and not detail:
        detail = "Unauthorized access. Please log in to continue."

    return JSONResponse(
        status_code=exc.status_code,
        content={"detail": detail},
        headers=getattr(exc, "headers", None),
    )

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Add rate limiting middleware
app.add_middleware(RateLimitMiddleware)

templates = Jinja2Templates(directory="templates")

# Initialize database
init_db()

# Load ML model if available
ML_MODEL = None
MODEL_PATH = os.getenv("MODEL_PATH", "rf_model.pkl")
if os.path.exists(MODEL_PATH):
    try:
        ML_MODEL = joblib.load(MODEL_PATH)
        logger.info(f"✅ Loaded ML model from {MODEL_PATH}")
    except Exception as e:
        logger.warning(f"⚠️ Failed to load ML model: {e}")

# Mount static files if directory exists
if os.path.exists("static"):
    app.mount("/static", StaticFiles(directory="static"), name="static")

# ============================================================================
# ERROR HANDLING DECORATOR
# ============================================================================

def handle_db_errors(func):
    """Decorator to handle database errors gracefully."""
    @wraps(func)
    async def wrapper(*args, **kwargs):
        try:
            return await func(*args, **kwargs)
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Database error in {func.__name__}: {str(e)}")
            logger.error(traceback.format_exc())
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Database operation failed: {str(e)}"
            )
    return wrapper


DISTRICT_LOCATIONS = {
    "district_1": {"lat": 13.0827, "lon": 80.2707, "name": "Central District"},
    "district_2": {"lat": 13.0545, "lon": 80.2450, "name": "West District"},
    "district_3": {"lat": 12.9716, "lon": 80.1534, "name": "South District"},
    "district_4": {"lat": 13.1278, "lon": 80.2270, "name": "North District"},
}


def _normalize_severity(raw_severity: Optional[str]) -> str:
    value = str(raw_severity or "").lower()
    if value in {"1", "low", "minor"}:
        return "low"
    if value in {"2", "moderate", "medium"}:
        return "moderate"
    if value in {"3", "high"}:
        return "high"
    return "unknown"


def _severity_color(severity: Optional[str]) -> str:
    palette = {
        "low": "#22c55e",
        "moderate": "#f59e0b",
        "high": "#ef4444",
        "unknown": "#64748b",
    }
    return palette.get(_normalize_severity(severity), "#64748b")


def _load_police_incidents(district_id: str) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    incidents = get_traffic_incidents(district["lat"], district["lon"], radius=7000)

    normalized: list[dict] = []
    for index, incident in enumerate(incidents):
        location = incident.get("location") or []
        latitude = None
        longitude = None
        if isinstance(location, (list, tuple)) and len(location) >= 2:
            longitude, latitude = location[0], location[1]
        normalized.append({
            "id": incident.get("id") or f"incident-{index}",
            "type": incident.get("type") or "traffic",
            "severity": _normalize_severity(incident.get("severity")),
            "severity_color": _severity_color(incident.get("severity")),
            "description": incident.get("description") or "Traffic incident",
            "latitude": latitude,
            "longitude": longitude,
            "start_time": incident.get("start_time"),
            "end_time": incident.get("end_time"),
            "district_id": district_id,
            "response_time": incident.get("response_time") or incident.get("responseTime"),
        })

    for manual_incident in _manual_incidents_store.get(district_id, []):
        normalized.append({
            "id": manual_incident.get("id") or f"manual-{uuid.uuid4().hex[:8]}",
            "type": manual_incident.get("type") or "traffic",
            "severity": _normalize_severity(manual_incident.get("severity")),
            "severity_color": _severity_color(manual_incident.get("severity")),
            "description": manual_incident.get("description") or "Manual incident",
            "latitude": manual_incident.get("latitude"),
            "longitude": manual_incident.get("longitude"),
            "start_time": manual_incident.get("start_time") or datetime.now(UTC).isoformat(),
            "end_time": manual_incident.get("end_time"),
            "district_id": district_id,
            "response_time": manual_incident.get("response_time"),
        })

    return normalized


def _build_district_summary(incidents: list[dict]) -> dict:
    today = datetime.now(UTC).date()
    today_count = 0
    response_times: list[float] = []

    for incident in incidents:
        start_time = incident.get("start_time")
        if start_time:
            try:
                parsed = datetime.fromisoformat(str(start_time).replace("Z", "+00:00"))
                if parsed.date() == today:
                    today_count += 1
            except ValueError:
                today_count += 1
        else:
            today_count += 1

        response_time = incident.get("response_time")
        if isinstance(response_time, (int, float)):
            response_times.append(float(response_time))

    avg_response_time = round(sum(response_times) / len(response_times), 2) if response_times else 0.0
    return {
        "total_incidents_today": today_count,
        "avg_response_time": avg_response_time,
    }


def _format_police_timestamp(value: object) -> str:
    if not value:
        return "-"
    if isinstance(value, datetime):
        return value.astimezone(UTC).strftime("%I:%M %p")
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed.astimezone(UTC).strftime("%I:%M %p")
    except ValueError:
        return str(value)


def _incident_sort_value(incident: dict) -> datetime:
    raw_value = incident.get("start_time")
    if isinstance(raw_value, datetime):
        return raw_value if raw_value.tzinfo else raw_value.replace(tzinfo=UTC)
    if raw_value:
        try:
            parsed = datetime.fromisoformat(str(raw_value).replace("Z", "+00:00"))
            return parsed if parsed.tzinfo else parsed.replace(tzinfo=UTC)
        except ValueError:
            return datetime.min.replace(tzinfo=UTC)
    return datetime.min.replace(tzinfo=UTC)


def _build_incidents_feed(district_id: str, incidents: list[dict], assignments: list[PoliceDispatchAssignment]) -> list[dict]:
    assignment_map = {assignment.incident_id: assignment for assignment in assignments}

    feed_items: list[dict] = []
    for incident in incidents:
        incident_id = incident.get("id")
        assignment = assignment_map.get(incident_id)
        feed_items.append({
            "incident_id": incident_id,
            "type": incident.get("type") or "traffic",
            "location_name": incident.get("description") or "Unknown location",
            "severity": incident.get("severity") or "unknown",
            "time_reported": incident.get("start_time") or "-",
            "assigned_unit": assignment.unit_id if assignment else None,
            "is_assigned": assignment is not None,
        })

    feed_items.sort(
        key=lambda item: _incident_sort_value({"start_time": item.get("time_reported")}),
        reverse=True,
    )
    return feed_items


def _build_heatmap_points(incidents: list[dict]) -> list[dict]:
    severity_intensity = {
        "low": 0.35,
        "moderate": 0.65,
        "medium": 0.65,
        "high": 1.0,
        "unknown": 0.45,
    }

    points: list[dict] = []
    for incident in incidents:
        lat = incident.get("latitude")
        lng = incident.get("longitude")
        if lat is None or lng is None:
            continue
        try:
            severity = str(incident.get("severity") or "unknown").lower()
            intensity = severity_intensity.get(severity, 0.45)
            points.append({
                "lat": float(lat),
                "lng": float(lng),
                "intensity": float(intensity),
            })
        except (TypeError, ValueError):
            continue

    return points


def _parse_police_datetime(value: object) -> Optional[datetime]:
    if isinstance(value, datetime):
        return value if value.tzinfo else value.replace(tzinfo=UTC)
    if not value:
        return None
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed if parsed.tzinfo else parsed.replace(tzinfo=UTC)
    except ValueError:
        return None


def _incident_response_minutes(incident: dict) -> float:
    response_value = incident.get("response_time")
    if isinstance(response_value, (int, float)) and float(response_value) > 0:
        return float(response_value)

    severity_default = {
        "low": 6.0,
        "moderate": 8.0,
        "medium": 8.0,
        "high": 11.0,
        "unknown": 9.0,
    }
    severity = str(incident.get("severity") or "unknown").lower()
    return severity_default.get(severity, 9.0)


def _infer_zone_name(district_id: str, incident: dict) -> str:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    lat = incident.get("latitude")
    lon = incident.get("longitude")

    if isinstance(lat, (int, float)) and isinstance(lon, (int, float)):
        delta_lat = float(lat) - float(district["lat"])
        delta_lon = float(lon) - float(district["lon"])
        core_radius = 0.004

        if abs(delta_lat) <= core_radius and abs(delta_lon) <= core_radius:
            return "Central Zone"
        if abs(delta_lat) >= abs(delta_lon):
            return "North Zone" if delta_lat >= 0 else "South Zone"
        return "East Zone" if delta_lon >= 0 else "West Zone"

    description = str(incident.get("description") or "").strip().lower()
    if "north" in description:
        return "North Zone"
    if "south" in description:
        return "South Zone"
    if "east" in description:
        return "East Zone"
    if "west" in description:
        return "West Zone"
    return "Central Zone"


def _extract_affected_roads(incident: dict, zone_name: str) -> list[str]:
    """Extract and sanitize road-impact data without exposing incident details."""
    raw_roads = incident.get("affected_roads")
    if isinstance(raw_roads, list):
        cleaned = []
        for road in raw_roads:
            text_value = str(road or "").strip()
            if text_value:
                cleaned.append(text_value)
        if cleaned:
            return list(dict.fromkeys(cleaned))[:6]

    description = str(incident.get("description") or "")
    road_pattern = re.compile(
        r"([A-Za-z0-9 .'-]{2,80}\\b(?:Road|Rd|Street|St|Avenue|Ave|Highway|Hwy|Expressway|Expwy|Boulevard|Blvd|Lane|Ln|Marg|Flyover))",
        re.IGNORECASE,
    )
    extracted = [" ".join(match.split()) for match in road_pattern.findall(description)]
    deduped = []
    for road_name in extracted:
        if road_name and road_name not in deduped:
            deduped.append(road_name)

    if deduped:
        return deduped[:6]

    return [f"{zone_name} arterial corridors"]


def _shared_alert_expiry(timestamp: datetime, severity: str) -> datetime:
    durations = {
        "critical": 120,
        "high": 90,
    }
    minutes = durations.get(str(severity or "").lower(), 60)
    return timestamp + timedelta(minutes=minutes)


def _create_shared_alert_for_dispatch(session: Session, district_id: str, incident: dict) -> Optional[SharedAlert]:
    """Persist sanitized logistics alert for high-impact dispatches only."""
    severity = _normalize_severity(incident.get("severity"))
    if severity not in {"critical", "high"}:
        return None

    now = datetime.now(UTC)
    zone_name = _infer_zone_name(district_id, incident)
    alert = SharedAlert(
        alert_id=str(uuid.uuid4()),
        zone=zone_name,
        severity=severity,
        timestamp=now,
        affected_roads=_extract_affected_roads(incident, zone_name),
        expires_at=_shared_alert_expiry(now, severity),
    )
    session.add(alert)
    return alert


def _build_response_time_by_zone(
    district_id: str,
    incidents: list[dict],
    assignments: list[PoliceDispatchAssignment],
    target_threshold_minutes: float = 8.0,
) -> list[dict]:
    assignment_map = {assignment.incident_id: assignment for assignment in assignments}
    today = datetime.now(UTC).date()
    def aggregate_rows(source_incidents: list[dict], today_only: bool) -> dict[str, dict]:
        zone_stats: dict[str, dict] = {}
        for incident in source_incidents:
            start_time = _parse_police_datetime(incident.get("start_time"))
            if today_only and start_time and start_time.date() != today:
                continue

            zone_name = _infer_zone_name(district_id, incident)
            bucket = zone_stats.setdefault(
                zone_name,
                {
                    "zone": zone_name,
                    "total_incidents": 0,
                    "response_sum": 0.0,
                    "response_count": 0,
                    "unit_times": {},
                },
            )

            bucket["total_incidents"] += 1
            response_minutes = _incident_response_minutes(incident)
            bucket["response_sum"] += response_minutes
            bucket["response_count"] += 1

            assignment = assignment_map.get(incident.get("id"))
            unit_id = assignment.unit_id if assignment else "Unassigned"
            unit_bucket = bucket["unit_times"].setdefault(unit_id, {"sum": 0.0, "count": 0})
            unit_bucket["sum"] += response_minutes
            unit_bucket["count"] += 1

        return zone_stats

    zone_stats = aggregate_rows(incidents, today_only=True)
    if not zone_stats:
        # Fallback for sparse/legacy datasets that don't carry today's timestamps.
        zone_stats = aggregate_rows(incidents, today_only=False)

    response_rows: list[dict] = []
    for zone_name, stats in zone_stats.items():
        avg_response = round(stats["response_sum"] / max(stats["response_count"], 1), 2)

        unit_averages = []
        for unit_id, values in stats["unit_times"].items():
            if values["count"] <= 0:
                continue
            unit_averages.append((unit_id, values["sum"] / values["count"]))

        unit_averages.sort(key=lambda item: item[1])
        fastest_unit = unit_averages[0][0] if unit_averages else "-"
        slowest_unit = unit_averages[-1][0] if unit_averages else "-"

        response_rows.append({
            "zone": zone_name,
            "avg_response_time": avg_response,
            "fastest_unit": fastest_unit,
            "slowest_unit": slowest_unit,
            "total_incidents": stats["total_incidents"],
            "exceeds_target": avg_response > float(target_threshold_minutes),
        })

    response_rows.sort(key=lambda item: item["avg_response_time"], reverse=True)
    return response_rows


def _normalize_patrol_status(raw_status: object) -> str:
    status_value = str(raw_status or "").strip().lower()
    if status_value in {"responding", "response", "enroute", "en route"}:
        return "Responding"
    if status_value in {"idle", "available", "standby"}:
        return "Idle"
    if status_value == "active":
        return "Active"
    return "Active"


def _get_dispatch_assignments(district_id: str) -> list[PoliceDispatchAssignment]:
    session = get_session()
    try:
        return (
            session.query(PoliceDispatchAssignment)
            .filter(PoliceDispatchAssignment.district_id == district_id)
            .order_by(PoliceDispatchAssignment.assigned_at.desc())
            .all()
        )
    finally:
        session.close()


def _build_patrol_units(
    district_id: str,
    incidents: list[dict],
    supervisor_name: str,
    assignments: Optional[list[PoliceDispatchAssignment]] = None,
) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    session = get_session()
    try:
        police_users = (
            session.query(User)
            .filter(User.department == "police", User.is_active == True)  # noqa: E712
            .order_by(User.id.asc())
            .all()
        )
    finally:
        session.close()

    assignments = assignments or []
    assignment_by_unit = {assignment.unit_id: assignment for assignment in assignments}
    assignment_by_incident = {assignment.incident_id: assignment for assignment in assignments}

    unit_count = max(4, len(incidents), len(police_users))
    patrol_units: list[dict] = []

    for index in range(unit_count):
        incident = incidents[index % len(incidents)] if incidents else None
        officer = police_users[index % len(police_users)] if police_users else None
        unit_id = f"{district_id.upper().replace('_', '-')}-U{index + 1:02d}"
        assignment = assignment_by_unit.get(unit_id)

        if assignment and incident:
            status = "responding"
        elif assignment:
            status = "responding"
        else:
            status = "available" if index % 3 else "responding"

        officer_name = (
            getattr(officer, "full_name", None)
            or getattr(officer, "username", None)
            or supervisor_name
            or f"Officer {index + 1}"
        )

        location = district["name"]
        last_updated = datetime.now(UTC)

        if assignment:
            incident_match = next((item for item in incidents if item.get("id") == assignment.incident_id), None)
            if incident_match:
                location = incident_match.get("description") or incident_match.get("type") or district["name"]
                if incident_match.get("latitude") is not None and incident_match.get("longitude") is not None:
                    location = f"{location} ({incident_match['latitude']:.4f}, {incident_match['longitude']:.4f})"
            else:
                location = f"Assigned to incident {assignment.incident_id}"
            last_updated = assignment.assigned_at or last_updated
        elif incident:
            location = incident.get("description") or incident.get("type") or district["name"]
            if incident.get("latitude") is not None and incident.get("longitude") is not None:
                location = f"{location} ({incident['latitude']:.4f}, {incident['longitude']:.4f})"
            last_updated = incident.get("start_time") or last_updated

        latitude = district["lat"]
        longitude = district["lon"]
        if assignment:
            incident_match = next((item for item in incidents if item.get("id") == assignment.incident_id), None)
            if incident_match and incident_match.get("latitude") is not None and incident_match.get("longitude") is not None:
                latitude = float(incident_match["latitude"])
                longitude = float(incident_match["longitude"])
        elif incident and incident.get("latitude") is not None and incident.get("longitude") is not None:
            latitude = float(incident["latitude"])
            longitude = float(incident["longitude"])

        patrol_units.append({
            "unit_id": unit_id,
            "officer_name": officer_name,
            "status": status,
            "current_location": location,
            "last_updated": _format_police_timestamp(last_updated),
            "district_id": district_id,
            "assigned_incident_id": assignment.incident_id if assignment else None,
            "latitude": latitude,
            "longitude": longitude,
        })

    return patrol_units


def _build_police_dashboard_context(current_user: dict, district_id: str) -> dict:
    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    incidents_feed = _build_incidents_feed(district_id, incidents, assignments)
    district_summary = _build_district_summary(incidents)
    district_info = DISTRICT_LOCATIONS.get(district_id, {"name": district_id or "Unknown District"})
    supervisor_name = current_user.get("username", "Unknown Supervisor")
    patrol_units = _build_patrol_units(district_id, incidents, supervisor_name, assignments)

    assigned_incident_ids = {assignment.incident_id for assignment in assignments}
    unassigned_incidents = [incident for incident in incidents if incident.get("id") not in assigned_incident_ids]
    available_patrol_units = [unit for unit in patrol_units if unit["status"] == "available"]
    active_units = [unit for unit in patrol_units if unit["status"] != "available"]

    return {
        "district_info": district_info,
        "district_name": district_info.get("name", district_id or "Unknown District"),
        "supervisor_name": supervisor_name,
        "shift_time": datetime.now(UTC).strftime("%I:%M %p UTC"),
        "total_active_incidents": len(unassigned_incidents),
        "units_deployed": len(active_units),
        "units_available": len(available_patrol_units),
        "avg_response_time": district_summary.get("avg_response_time", 0.0),
        "patrol_units": patrol_units,
        "available_patrol_units": available_patrol_units,
        "unassigned_incidents": unassigned_incidents,
        "incidents_feed": incidents_feed,
        "incidents": incidents,
        "district_summary": district_summary,
    }


def _district_prediction_candidates(district_id: str, incidents: list[dict]) -> list[dict]:
    district = DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])
    candidates: list[dict] = []

    for incident in incidents:
        if incident.get("latitude") is not None and incident.get("longitude") is not None:
            candidates.append({
                "location": incident.get("description") or "Incident hotspot",
                "latitude": incident["latitude"],
                "longitude": incident["longitude"],
                "base_severity": incident.get("severity", "unknown"),
                "source": "live_incident",
            })

    if not candidates:
        offsets = [
            (0.0000, 0.0000),
            (0.0120, 0.0080),
            (-0.0100, 0.0110),
            (0.0090, -0.0090),
            (-0.0080, -0.0100),
        ]
        for index, (lat_offset, lon_offset) in enumerate(offsets):
            candidates.append({
                "location": f"{district['name']} sector {index + 1}",
                "latitude": district["lat"] + lat_offset,
                "longitude": district["lon"] + lon_offset,
                "base_severity": "unknown",
                "source": "district_grid",
            })

    return candidates[:10]


def _predict_police_hotspots(district_id: str, incidents: list[dict]) -> list[dict]:
    """Use the existing loaded ML model to rank likely hotspot locations."""
    candidates = _district_prediction_candidates(district_id, incidents)
    if not candidates:
        return []

    now = datetime.now(UTC)
    predictions: list[dict] = []

    for candidate in candidates:
        hourly_scores: list[float] = []
        for hours_ahead in range(1, 7):
            future_time = now + timedelta(hours=hours_ahead)
            distance_km = haversine_m(
                candidate["latitude"],
                candidate["longitude"],
                DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])["lat"],
                DISTRICT_LOCATIONS.get(district_id, DISTRICT_LOCATIONS["district_1"])["lon"],
            ) / 1000.0

            severity_bias = {
                "low": 0.05,
                "moderate": 0.15,
                "high": 0.30,
                "unknown": 0.10,
            }.get(candidate.get("base_severity", "unknown"), 0.10)

            feature_frame = {
                "hour": future_time.hour,
                "weekday": future_time.weekday(),
                "is_weekend": 1 if future_time.weekday() >= 5 else 0,
                "distance_km": max(distance_km, 0.1),
                "route_index": 0,
                "travel_time_s": 900 + int(distance_km * 180) + int(severity_bias * 300),
                "no_traffic_s": 750 + int(distance_km * 120),
                "delay_s": 150 + int(severity_bias * 240),
                "rolling_mean_congestion": 1.0 + severity_bias,
                "rolling_std_congestion": 0.05 + severity_bias / 2,
            }

            model_prediction = predict_congestion(feature_frame)
            if model_prediction is None:
                model_prediction = 1.0 + severity_bias

            likelihood_score = max(0.0, min(100.0, ((float(model_prediction) - 0.9) / 1.2) * 100 + severity_bias * 15))
            hourly_scores.append(likelihood_score)

        average_score = round(sum(hourly_scores) / len(hourly_scores), 2)
        confidence = round(max(0.0, min(100.0, 100.0 - (max(hourly_scores) - min(hourly_scores)) * 0.75)), 2)

        if average_score >= 85:
            predicted_type = "critical congestion"
        elif average_score >= 70:
            predicted_type = "high traffic"
        elif average_score >= 50:
            predicted_type = "moderate traffic"
        else:
            predicted_type = "light traffic"

        predictions.append({
            "zone_name": candidate["location"],
            "location": candidate["location"],
            "latitude": candidate["latitude"],
            "longitude": candidate["longitude"],
            "incident_count": 1 if candidate.get("source") == "live_incident" else 0,
            "severity": candidate.get("base_severity", "unknown"),
            "likelihood_score": average_score,
            "confidence": confidence,
            "predicted_type": predicted_type,
        })

    predictions.sort(key=lambda item: item["likelihood_score"], reverse=True)
    return predictions[:5]


def _safe_ppt_text(value: object) -> str:
    if value is None:
        return ""
    return str(value)


def _format_ppt_timestamp(value: Optional[str]) -> str:
    if not value:
        return "-"
    try:
        parsed = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        return parsed.strftime("%Y-%m-%d %H:%M")
    except ValueError:
        return str(value)


def generate_shift_pptx(
    district_id: str,
    officer_name: str,
    incidents: list[dict],
    district_summary: dict,
    ml_predictions: list[dict],
) -> io.BytesIO:
    """Generate a police shift PPTX report and return an in-memory stream."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE

    district_info = DISTRICT_LOCATIONS.get(district_id, {"name": district_id or "Unknown District"})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_body_layout = prs.slide_layouts[5]

    # Slide 1 - Shift summary
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Police Shift Summary"
    slide1.placeholders[1].text = (
        f"District: {district_info.get('name', district_id or 'Unknown District')}\n"
        f"Date: {datetime.now(UTC).strftime('%Y-%m-%d')}\n"
        f"Officer: {officer_name}\n"
        f"Total Incidents: {len(incidents)}\n"
        f"Avg Response Time: {district_summary.get('avg_response_time', 'N/A')}"
    )

    # Slide 2 - Incidents table
    slide2 = prs.slides.add_slide(title_body_layout)
    slide2.shapes.title.text = "Incident Breakdown"
    rows = max(len(incidents), 1) + 1
    cols = 4
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6))
    table = table_shape.table
    headers = ["ID", "Description", "Severity", "Start Time"]
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
    for row_index, incident in enumerate(incidents, start=1):
        table.cell(row_index, 0).text = _safe_ppt_text(incident.get("id"))
        table.cell(row_index, 1).text = _safe_ppt_text(incident.get("description"))
        table.cell(row_index, 2).text = _safe_ppt_text(incident.get("severity"))
        table.cell(row_index, 3).text = _format_ppt_timestamp(incident.get("start_time"))

    # Slide 3 - ML highlights
    slide3 = prs.slides.add_slide(title_body_layout)
    slide3.shapes.title.text = "ML Prediction Highlights"
    text_box = slide3.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
    text_frame = text_box.text_frame
    if ml_predictions:
        for index, prediction in enumerate(ml_predictions, start=1):
            paragraph = text_frame.paragraphs[0] if index == 1 else text_frame.add_paragraph()
            paragraph.text = (
                f"{index}. {prediction['location']} | "
                f"Likelihood: {prediction['likelihood_score']}% | "
                f"Confidence: {prediction['confidence']}% | "
                f"Type: {prediction['predicted_type']}"
            )
            for run in paragraph.runs:
                run.font.size = Pt(18)
    else:
        text_frame.text = "No ML predictions available for this district."

    # Slide 4 - Map placeholder
    slide4 = prs.slides.add_slide(title_body_layout)
    slide4.shapes.title.text = "Map Screenshot Placeholder"
    placeholder = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.1), Inches(1.6), Inches(11.2), Inches(4.6)
    )
    placeholder.text_frame.text = (
        "Insert map screenshot here before final delivery.\n\n"
        f"District: {district_info.get('name', district_id or 'Unknown District')}\n"
        f"Incidents on map: {len(incidents)}"
    )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ============================================================================
# VALIDATION MODELS (Pydantic V2 Compatible)
# ============================================================================

class ValidatedCoordinates(BaseModel):
    """Validated coordinate pair."""
    lat: float
    lon: float
    
    @field_validator('lat')
    @classmethod
    def validate_latitude(cls, v: float) -> float:
        """Validate latitude is within -90 to 90."""
        if v is None:
            raise ValueError('Latitude is required')
        if not -90 <= v <= 90:
            raise ValueError(f'Invalid latitude: {v}. Must be between -90 and 90')
        return v
    
    @field_validator('lon')
    @classmethod
    def validate_longitude(cls, v: float) -> float:
        """Validate longitude is within -180 to 180."""
        if v is None:
            raise ValueError('Longitude is required')
        if not -180 <= v <= 180:
            raise ValueError(f'Invalid longitude: {v}. Must be between -180 and 180')
        return v


class RouteAnalysisRequest(BaseModel):
    """Request model for route analysis with validation."""
    origin: Union[str, dict] = Field(..., description="Origin as place name or {lat, lon}")
    destination: Union[str, dict] = Field(..., description="Destination as place name or {lat, lon}")
    maxAlternatives: int = Field(3, ge=1, le=5, description="Maximum route alternatives")
    alpha: float = Field(1.0, ge=0, le=10, description="Weight for travel time in cost calculation")
    beta: float = Field(0.5, ge=0, le=10, description="Weight for delay in cost calculation")
    gamma: float = Field(0.001, ge=0, le=1, description="Weight for distance in cost calculation")
    avoid_tolls: bool = Field(False, description="Avoid toll roads")
    avoid_ferries: bool = Field(False, description="Avoid ferries")
    avoid_highways: bool = Field(False, description="Avoid highways")
    
    @field_validator('maxAlternatives')
    @classmethod
    def validate_max_alternatives(cls, v: int) -> int:
        if v < 1 or v > 5:
            raise ValueError('maxAlternatives must be between 1 and 5')
        return v
    
    @field_validator('alpha', 'beta', 'gamma')
    @classmethod
    def validate_weights(cls, v: float, info) -> float:
        if v < 0:
            raise ValueError(f'{info.field_name} cannot be negative')
        return v


class UserCreate(BaseModel):
    """User registration model."""
    email: EmailStr
    username: str = Field(..., min_length=3, max_length=50)
    password: str = Field(..., min_length=6)
    full_name: Optional[str] = None
    role: str = Field("user", description="user | police_supervisor | logistics_manager")


class UserLogin(BaseModel):
    """User login model."""
    username: str
    password: str


class SavedRouteCreate(BaseModel):
    """Create saved route model."""
    route_name: str
    origin: Union[str, dict]
    destination: Union[str, dict]
    route_preferences: Optional[dict] = None


class RouteRatingCreate(BaseModel):
    """Create route rating model."""
    route_id: str
    rating: int = Field(..., ge=1, le=5)
    review: Optional[str] = None


class ShareRouteRequest(BaseModel):
    """Share route request model."""
    route_id: str
    route_index: Optional[int] = None


class UserUpdate(BaseModel):
    """User update model for admin editing."""
    username: Optional[str] = None
    email: Optional[EmailStr] = None
    full_name: Optional[str] = None
    is_active: Optional[bool] = None
    is_admin: Optional[bool] = None
    password: Optional[str] = Field(None, min_length=8)


class DispatchRequest(BaseModel):
    """Dispatch assignment request."""
    incident_id: str = Field(..., min_length=1)
    officer_id: Optional[str] = Field(None, min_length=1)
    unit_id: Optional[str] = Field(None, min_length=1)


class ShiftAttendanceRequest(BaseModel):
    """Mark officer attendance during shift."""
    officers: list[dict] = Field(..., description="List of {officer_username, officer_name, status}")


class ShiftEndRequest(BaseModel):
    """End shift and optionally export report."""
    notes: Optional[str] = None
    export_pptx: bool = True


class IncidentResolveRequest(BaseModel):
    """Incident resolution payload used for feedback logging and analytics."""
    incident_id: str = Field(..., min_length=1)
    incident_type: str = Field(..., min_length=1)
    zone: str = Field(..., min_length=1)
    response_time_minutes: float = Field(..., gt=0)
    severity: str = Field(..., pattern="^(critical|high|medium|low|moderate|unknown)$")
    outcome: str = Field(..., min_length=1)
    resolved_at: Optional[datetime] = None


class NewIncidentRequest(BaseModel):
    """Create a new command-center incident for live dispatch workflows."""
    incident_type: str = Field(..., min_length=1)
    severity: str = Field(..., pattern="^(critical|high|medium|low|moderate|unknown)$")
    description: str = Field(..., min_length=1)
    latitude: Optional[float] = None
    longitude: Optional[float] = None
    zone: Optional[str] = None


class AlertData(BaseModel):
    """Real-time alert for police supervisor."""
    alert_id: str
    severity: str  # critical, high, medium, low
    message: str
    timestamp: str  # ISO format
    district_id: Optional[str] = None
    related_incident_id: Optional[str] = None


class PredictedHotspotPolygon(BaseModel):
    """GeoJSON polygon for predicted incident hotspot."""
    zone_name: str
    risk_level: str  # high, medium, low
    prediction_score: float  # 0-100
    coordinates: list[list[list[float]]]  # GeoJSON Polygon coordinates


class PredictedIncidentResponse(BaseModel):
    """Response containing predicted incident hotspots as GeoJSON."""
    type: str = "FeatureCollection"
    features: list[dict]  # GeoJSON features


class OfficerWorkloadData(BaseModel):
    """Officer incident counts and rotation status."""
    officer_username: str
    officer_name: str
    total_incidents: int
    critical_incidents: int
    high_incidents: int
    medium_incidents: int
    low_incidents: int
    needs_rotation: bool


class OfficerWorkloadResponse(BaseModel):
    """Response containing officer workload data for current shift."""
    shift_id: int
    officers: list[OfficerWorkloadData]
    officers_needing_rotation: list[str]  # Usernames of officers needing rotation


class AlertListResponse(BaseModel):
    """Response for unread alerts list."""
    alerts: list[AlertData]
    unread_count: int

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def decode_polyline(encoded: str) -> list[tuple[float, float]]:
    """
    Decode Google-style polyline string to list of (lat, lon) tuples.
    Handles TomTom's coordinate formats properly.
    
    Args:
        encoded: Encoded polyline string or coordinate list
        
    Returns:
        List of (latitude, longitude) tuples
    """
    if not encoded:
        return []
    
    # Handle if encoded is already a list of coordinates
    if isinstance(encoded, list):
        coordinates = []
        for p in encoded:
            try:
                if isinstance(p, dict):
                    if "lat" in p and "lon" in p:
                        coordinates.append((float(p["lat"]), float(p["lon"])))
                    elif "latitude" in p and "longitude" in p:
                        coordinates.append((float(p["latitude"]), float(p["longitude"])))
                elif isinstance(p, (list, tuple)) and len(p) >= 2:
                    coordinates.append((float(p[0]), float(p[1])))
            except (ValueError, TypeError):
                continue
        return coordinates
    
    # Try to use the polyline library if available
    try:
        import polyline
        return polyline.decode(encoded)
    except ImportError:
        # Fallback to custom decoder if polyline not installed
        pass
    except Exception as e:
        logger.warning(f"Polyline library decoding failed: {e}")
    
    # Custom decoder implementation
    coordinates = []
    index = 0
    lat = 0
    lon = 0
    
    try:
        while index < len(encoded):
            # Decode latitude
            shift = 0
            result = 0
            while index < len(encoded):
                b = ord(encoded[index]) - 63
                index += 1
                result |= (b & 0x1f) << shift
                shift += 5
                if b < 0x20:
                    break
            dlat = ~(result >> 1) if (result & 1) else (result >> 1)
            lat += dlat
            
            # Check if we have enough characters for longitude
            if index >= len(encoded):
                break
                
            # Decode longitude
            shift = 0
            result = 0
            while index < len(encoded):
                b = ord(encoded[index]) - 63
                index += 1
                result |= (b & 0x1f) << shift
                shift += 5
                if b < 0x20:
                    break
            dlon = ~(result >> 1) if (result & 1) else (result >> 1)
            lon += dlon
            
            coordinates.append((lat / 1e5, lon / 1e5))
    except (IndexError, ValueError, TypeError) as e:
        logger.error(f"Polyline decoding error: {e}")
        return []
    
    return coordinates


def extract_route_geometry(route_json: dict) -> list[tuple[float, float]]:
    """
    Extract route geometry from TomTom route JSON.
    
    Args:
        route_json: Route object from TomTom API
        
    Returns:
        List of (lat, lon) tuples for route path
    """
    geometry = []
    
    try:
        legs = route_json.get("legs", [])
        for leg in legs:
            points = leg.get("points", [])
            for point in points:
                if isinstance(point, dict):
                    if "latitude" in point and "longitude" in point:
                        try:
                            geometry.append((float(point["latitude"]), float(point["longitude"])))
                        except (ValueError, TypeError):
                            continue
                    elif "lat" in point and "lon" in point:
                        try:
                            geometry.append((float(point["lat"]), float(point["lon"])))
                        except (ValueError, TypeError):
                            continue
        
        if not geometry:
            guidance = route_json.get("guidance", {})
            instructions = guidance.get("instructions", [])
            for instruction in instructions:
                point = instruction.get("point", {})
                if "latitude" in point and "longitude" in point:
                    try:
                        geometry.append((float(point["latitude"]), float(point["longitude"])))
                    except (ValueError, TypeError):
                        continue
    except Exception as e:
        logger.error(f"Error extracting route geometry: {e}")
        return []
    
    return geometry


def predict_congestion(features: dict) -> Optional[float]:
    """
    Predict route score using RF model.
    Lower score = better route.
    """
    if ML_MODEL is None:
        return None
    try:
        import pandas as pd
        now = datetime.now(UTC)

        travel_time_s = features.get("travel_time_s", 0)
        no_traffic_s = features.get("no_traffic_s", 0)
        if travel_time_s is None:
            travel_time_s = 0
        if no_traffic_s is None:
            no_traffic_s = 0

        congestion_ratio = (
            float(travel_time_s) / float(no_traffic_s)
            if no_traffic_s and float(no_traffic_s) > 0
            else 1.0
        )

        hour = int(features.get("hour", now.hour))
        weekday = int(features.get("weekday", now.weekday()))

        base_feature_dict = {
            "hour": hour,
            "weekday": weekday,
            "is_weekend": int(features.get("is_weekend", 1 if weekday >= 5 else 0)),
            "distance_km": float(features.get("distance_km", 0) or 0),
            "route_index": int(features.get("route_index", 0) or 0),
            "travel_time_s": float(travel_time_s),
            "no_traffic_s": float(no_traffic_s),
            "delay_s": float(features.get("delay_s", 0) or 0),
            "rolling_mean_congestion": float(features.get("rolling_mean_congestion", 1.0) or 1.0),
            "rolling_std_congestion": float(features.get("rolling_std_congestion", 0.0) or 0.0),
            "congestion_ratio": float(congestion_ratio),
        }

        model_feature_names = getattr(ML_MODEL, "feature_names_in_", None)
        if model_feature_names is not None:
            expected_columns = [str(col) for col in model_feature_names]
        else:
            expected_columns = list(base_feature_dict.keys())

        aligned_feature_dict = {
            col: base_feature_dict.get(col, 0.0) for col in expected_columns
        }
        X = pd.DataFrame([aligned_feature_dict], columns=expected_columns).apply(pd.to_numeric, errors="coerce").fillna(0.0)
        prediction = ML_MODEL.predict(X)[0]
        return round(float(prediction), 2)
    except Exception as e:
        logger.error(f"ML prediction error: {e}")
        return None


# ============================================================================
# FRONTEND ROUTES
# ============================================================================

@app.get("/", response_class=HTMLResponse)
async def serve_index(request: Request):
    """Serve index page. If not logged in, redirect to login."""
    # Check if user is authenticated
    token = request.cookies.get("token") or request.cookies.get("access_token")
    
    # If no token, redirect to login
    if not token:
        return RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    
    # User is logged in, serve index/route optimizer page
    index_path = os.path.join("templates", "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Traffic Route Analysis API</h1><p>Frontend not found. Please create templates/index.html</p>",
        status_code=404
    )


@app.get("/favicon.ico")
async def serve_favicon():
    """Serve the favicon."""
    favicon_path = os.path.join("static", "favicon.svg")
    if os.path.exists(favicon_path):
        with open(favicon_path, "r", encoding="utf-8") as f:
            return Response(content=f.read(), media_type="image/svg+xml")
    # Return a default empty response to prevent 404 logs
    return Response(content="", status_code=204)


@app.get("/login", response_class=HTMLResponse)
async def serve_login():
    """Serve the login/registration page."""
    login_path = os.path.join("templates", "login.html")
    if os.path.exists(login_path):
        with open(login_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Login</h1><p>Login page not found.</p>",
        status_code=404
    )


@app.get("/logout")
async def logout(response: Response):
    """Clear auth cookies and redirect to login."""
    redirect = RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    redirect.delete_cookie(key="token", path="/")
    redirect.delete_cookie(key="access_token", path="/")
    return redirect


@app.get("/auth/login")
async def serve_auth_login_alias(request: Request):
    """Alias route so auth redirects can target /auth/login for browser clients."""
    reason = request.query_params.get("reason")
    if reason:
        return RedirectResponse(url=f"/login?reason={reason}", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    return RedirectResponse(url="/login", status_code=status.HTTP_307_TEMPORARY_REDIRECT)


@app.get("/admin", response_class=HTMLResponse)
async def serve_admin(current_user: dict = Depends(require_role("admin"))):
    """Serve the admin dashboard page."""
    admin_path = os.path.join("templates", "admin.html")
    if os.path.exists(admin_path):
        with open(admin_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Admin Dashboard</h1><p>Admin page not found.</p>",
        status_code=404
    )


@app.get("/account", response_class=HTMLResponse)
async def serve_account(current_user: dict = Depends(get_current_user)):
    """Serve the user account page."""
    role = current_user.get("role")
    if role == UserRole.police_supervisor.value:
        return RedirectResponse(url="/police/dashboard", status_code=status.HTTP_307_TEMPORARY_REDIRECT)
    if role == UserRole.admin.value:
        return RedirectResponse(url="/admin", status_code=status.HTTP_307_TEMPORARY_REDIRECT)

    account_path = os.path.join("templates", "account.html")
    if os.path.exists(account_path):
        with open(account_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>My Account</h1><p>Account page not found.</p>",
        status_code=404
    )


@app.get("/password-toggle-demo", response_class=HTMLResponse)
async def serve_password_toggle_demo():
    """Serve the password toggle demo page."""
    demo_path = os.path.join("templates", "password_toggle_demo.html")
    if os.path.exists(demo_path):
        with open(demo_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Password Toggle Demo</h1><p>Demo page not found.</p>",
        status_code=404
    )


@app.get("/analysis-report", response_class=HTMLResponse)
async def serve_analysis_report():
    """Serve the analysis report HTML page."""
    report_path = os.path.join("templates", "analysis_report.html")
    if os.path.exists(report_path):
        with open(report_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(
        content="<h1>Analysis Report</h1><p>Report page not found.</p>",
        status_code=404
    )


@app.get("/static/manifest.json")
async def get_manifest():
    """Serve PWA manifest."""
    manifest_path = os.path.join("static", "manifest.json")
    if os.path.exists(manifest_path):
        return FileResponse(manifest_path, media_type="application/json")
    return JSONResponse({"error": "Manifest not found"}, status_code=404)


# ============================================================================
# API ROUTES
# ============================================================================

@app.get("/health")
async def health():
    """Health check endpoint."""
    return {
        "status": "healthy",
        "model_loaded": ML_MODEL is not None,
        "timestamp": datetime.now(UTC).isoformat()
    }


@app.get("/api/stats")
async def get_stats():
    """Get real statistics from database for stats bar."""
    try:
        session = get_session()
        
        # Count routes analyzed today
        today_start = datetime.now(UTC).replace(hour=0, minute=0, second=0, microsecond=0)
        today_count = session.query(AnalysisResult).filter(
            AnalysisResult.timestamp >= today_start
        ).count()
        
        # Get average congestion from last 10 records
        recent_records = session.query(AnalysisResult).order_by(
            AnalysisResult.timestamp.desc()
        ).limit(10).all()
        
        avg_congestion = 1.0
        if recent_records:
            ratios = []
            for r in recent_records:
                if r.travel_time_s and r.no_traffic_s and r.no_traffic_s > 0:
                    ratios.append(r.travel_time_s / r.no_traffic_s)
            if ratios:
                avg_congestion = round(sum(ratios) / len(ratios), 2)
        
        # Traffic status based on congestion
        if avg_congestion < 1.15:
            traffic_status = "Light"
            status_color = "#42a5f5"
        elif avg_congestion < 1.5:
            traffic_status = "Moderate"
            status_color = "#ffa726"
        else:
            traffic_status = "Heavy"
            status_color = "#ef5350"
        
        # Total all time count
        total_count = session.query(AnalysisResult).count()
        session.close()
        
        return {
            "routes_today": today_count,
            "total_routes": total_count,
            "avg_congestion": avg_congestion,
            "traffic_status": traffic_status,
            "status_color": status_color
        }
    except Exception as e:
        return {
            "routes_today": 0,
            "total_routes": 0,
            "avg_congestion": 1.0,
            "traffic_status": "Light",
            "status_color": "#42a5f5"
        }


@app.get("/autocomplete")
async def autocomplete(q: str = Query(..., description="Search query")):
    """
    Get autocomplete suggestions from TomTom API.
    
    Args:
        q: Search query string
        
    Returns:
        List of suggestion objects
    """
    try:
        suggestions = tomtom_autocomplete(q)
        return {"suggestions": suggestions}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Autocomplete failed: {str(e)}")


@app.post("/analyze-route")
async def analyze_route(
    request: RouteAnalysisRequest,
    current_user: Optional[User] = Depends(get_optional_user)
):
    """
    Analyze route alternatives with cost calculation and ML prediction.
    
    Args:
        request: RouteAnalysisRequest with origin, destination, and parameters
        current_user: Optional authenticated user
        
    Returns:
        JSON with analyzed routes and best route recommendation
    """
    try:
        # Parse and validate origin
        if isinstance(request.origin, str):
            o_lat, o_lon = tomtom_geocode(request.origin)
            origin_data = {"name": request.origin, "lat": o_lat, "lon": o_lon}
        else:
            o_lat = float(request.origin.get("lat"))
            o_lon = float(request.origin.get("lon"))
            origin_data = request.origin
        
        # Parse and validate destination
        if isinstance(request.destination, str):
            d_lat, d_lon = tomtom_geocode(request.destination)
            dest_data = {"name": request.destination, "lat": d_lat, "lon": d_lon}
        else:
            d_lat = float(request.destination.get("lat"))
            d_lon = float(request.destination.get("lon"))
            dest_data = request.destination
        
        # Validate coordinates
        if not (-90 <= o_lat <= 90) or not (-180 <= o_lon <= 180):
            raise ValueError("Invalid origin coordinates")
        if not (-90 <= d_lat <= 90) or not (-180 <= d_lon <= 180):
            raise ValueError("Invalid destination coordinates")
        
        # Fetch routes from TomTom
        route_json = tomtom_route(
            o_lat, o_lon, d_lat, d_lon,
            maxAlternatives=request.maxAlternatives
        )
        
        routes = route_json.get("routes", [])
        if not routes:
            raise HTTPException(status_code=404, detail="No routes found")
        
        # Analyze each route
        analyzed_routes = []
        route_id = f"{origin_data.get('name', f'{o_lat},{o_lon}')}→{dest_data.get('name', f'{d_lat},{d_lon}')}"
        
        for idx, route in enumerate(routes):
            summary = summarize_route(route)
            
            if summary["length_m"] == 0:
                summary["length_m"] = haversine_m(o_lat, o_lon, d_lat, d_lon)
            
            cost = compute_route_cost(
                summary["travel_time_s"],
                summary["no_traffic_s"],
                summary["delay_s"],
                summary["length_m"],
                alpha=request.alpha,
                beta=request.beta,
                gamma=request.gamma
            )
            
            ml_predicted = predict_congestion({
                "distance_km": summary["length_m"] / 1000.0,
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": summary["delay_s"]
            })
            
            svr_predicted = None
            try:
                from svr_model import svr_predict
                svr_predicted = svr_predict({})
            except Exception as e:
                svr_predicted = None
            
            congestion_ratio = (
                summary["travel_time_s"] / summary["no_traffic_s"]
                if summary["no_traffic_s"] and summary["no_traffic_s"] > 0
                else None
            )
            
            geometry = extract_route_geometry(route)
            
            calculated_delay = 0
            if summary["travel_time_s"] and summary["no_traffic_s"]:
                calculated_delay = max(0, summary["travel_time_s"] - summary["no_traffic_s"])
            elif summary.get("delay_s"):
                calculated_delay = summary["delay_s"]
            
            analyzed_route = {
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": calculated_delay,
                "length_m": summary["length_m"],
                "congestion_ratio": congestion_ratio,
                "calculated_cost": cost,
                "ml_predicted_congestion": ml_predicted,
                "svr_predicted_congestion": svr_predicted,
                "geometry": geometry
            }
            analyzed_routes.append(analyzed_route)
            
            # Save to database
            try:
                session = get_session()
                save_analysis(session, {
                    "route_id": f"{route_id}_route{idx}",
                    "origin": origin_data,
                    "destination": dest_data,
                    "travel_time_s": summary["travel_time_s"],
                    "no_traffic_s": summary["no_traffic_s"],
                    "delay_s": summary["delay_s"],
                    "length_m": summary["length_m"],
                    "calculated_cost": cost,
                    "ml_predicted": ml_predicted,
                    "raw_json": route,
                    "user_id": current_user.id if current_user else None
                })
                session.close()
            except Exception as e:
                logger.error(f"Database save error: {e}")
        
        # Find best route (lowest cost)
        best_route = min(analyzed_routes, key=lambda x: x["calculated_cost"])
        
        return {
            "origin": origin_data,
            "destination": dest_data,
            "analyzed_routes": analyzed_routes,
            "best_route_index": best_route["route_index"],
            "best_route": best_route,
            "timestamp": datetime.now(UTC).isoformat()
        }
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Route analysis error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Route analysis failed: {str(e)}")


@app.get("/api/route-analysis/{route_id}")
async def get_route_analysis(route_id: str, route_index: Optional[int] = None):
    """
    Get detailed analysis data for a specific route.
    
    Args:
        route_id: Route identifier (e.g., "Origin→Destination")
        route_index: Optional route index to filter specific route variant
        
    Returns:
        Analysis data with historical trends and statistics
    """
    try:
        session = get_session()
        
        route_id = route_id.replace('%E2%86%92', '→')
        
        if route_index is not None:
            query = session.query(AnalysisResult).filter(
                AnalysisResult.route_id == f"{route_id}_route{route_index}"
            )
        else:
            query = session.query(AnalysisResult).filter(
                AnalysisResult.route_id.like(f"{route_id}_route%")
            )
        
        results = query.order_by(AnalysisResult.timestamp.desc()).all()
        session.close()
        
        if not results:
            raise HTTPException(status_code=404, detail="No analysis data found for this route")
        
        analysis_data = []
        for r in results:
            try:
                origin = json.loads(r.origin) if r.origin and r.origin.startswith('{') else {"name": r.origin}
                dest = json.loads(r.destination) if r.destination and r.destination.startswith('{') else {"name": r.destination}
            except:
                origin = {"name": r.origin}
                dest = {"name": r.destination}
            
            delay_val = r.delay_s
            if delay_val is None or delay_val == 0:
                if r.travel_time_s and r.no_traffic_s:
                    delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                else:
                    delay_val = 0
            
            analysis_data.append({
                "id": r.id,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "route_id": r.route_id,
                "origin": origin,
                "destination": dest,
                "travel_time_s": r.travel_time_s,
                "no_traffic_s": r.no_traffic_s,
                "delay_s": delay_val,
                "length_m": r.length_m,
                "calculated_cost": r.calculated_cost,
                "ml_predicted": r.ml_predicted,
                "congestion_ratio": (r.travel_time_s / r.no_traffic_s) if r.no_traffic_s and r.no_traffic_s > 0 else None
            })
        
        # Calculate statistics
        if analysis_data:
            travel_times = [d["travel_time_s"] for d in analysis_data if d["travel_time_s"]]
            delays = []
            for d in analysis_data:
                delay_val = d.get("delay_s")
                if delay_val is None or delay_val == 0:
                    if d.get("travel_time_s") and d.get("no_traffic_s"):
                        delay_val = max(0, d["travel_time_s"] - d["no_traffic_s"])
                    else:
                        delay_val = 0
                if delay_val > 0:
                    delays.append(delay_val)
            
            costs = [d["calculated_cost"] for d in analysis_data if d.get("calculated_cost")]
            congestion_ratios = [d["congestion_ratio"] for d in analysis_data if d.get("congestion_ratio")]
            
            stats = {
                "avg_travel_time": sum(travel_times) / len(travel_times) if travel_times else 0,
                "avg_delay": sum(delays) / len(delays) if delays else 0,
                "avg_cost": sum(costs) / len(costs) if costs else 0,
                "avg_congestion": sum(congestion_ratios) / len(congestion_ratios) if congestion_ratios else 0,
                "min_travel_time": min(travel_times) if travel_times else 0,
                "max_travel_time": max(travel_times) if travel_times else 0,
                "total_records": len(analysis_data)
            }
        else:
            stats = {}
        
        response_data = {
            "route_id": route_id,
            "route_index": route_index,
            "analysis_data": analysis_data,
            "statistics": stats,
            "latest": analysis_data[0] if analysis_data else None,
            "fetched_at": datetime.now(UTC).isoformat()
        }
        
        response = JSONResponse(content=response_data)
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["Pragma"] = "no-cache"
        response.headers["Expires"] = "0"
        return response
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch analysis: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch analysis: {str(e)}")


@app.post("/api/refresh-route")
async def refresh_route_analysis(
    request: RouteAnalysisRequest,
    current_user: Optional[User] = Depends(get_optional_user)
):
    """
    Refresh route analysis with latest data from TomTom API.
    Similar to analyze_route but returns data in format suitable for analysis report.
    """
    try:
        # Reuse analyze_route logic
        if isinstance(request.origin, str):
            o_lat, o_lon = tomtom_geocode(request.origin)
            origin_data = {"name": request.origin, "lat": o_lat, "lon": o_lon}
        else:
            o_lat = float(request.origin.get("lat"))
            o_lon = float(request.origin.get("lon"))
            origin_data = request.origin
        
        if isinstance(request.destination, str):
            d_lat, d_lon = tomtom_geocode(request.destination)
            dest_data = {"name": request.destination, "lat": d_lat, "lon": d_lon}
        else:
            d_lat = float(request.destination.get("lat"))
            d_lon = float(request.destination.get("lon"))
            dest_data = request.destination
        
        # Validate coordinates
        if not (-90 <= o_lat <= 90) or not (-180 <= o_lon <= 180):
            raise ValueError("Invalid origin coordinates")
        if not (-90 <= d_lat <= 90) or not (-180 <= d_lon <= 180):
            raise ValueError("Invalid destination coordinates")
        
        route_json = tomtom_route(
            o_lat, o_lon, d_lat, d_lon,
            maxAlternatives=request.maxAlternatives
        )
        
        routes = route_json.get("routes", [])
        if not routes:
            raise HTTPException(status_code=404, detail="No routes found")
        
        analyzed_routes = []
        route_id = f"{origin_data.get('name', f'{o_lat},{o_lon}')}→{dest_data.get('name', f'{d_lat},{d_lon}')}"
        
        for idx, route in enumerate(routes):
            summary = summarize_route(route)
            
            if summary["length_m"] == 0:
                summary["length_m"] = haversine_m(o_lat, o_lon, d_lat, d_lon)
            
            cost = compute_route_cost(
                summary["travel_time_s"],
                summary["no_traffic_s"],
                summary["delay_s"],
                summary["length_m"],
                alpha=request.alpha,
                beta=request.beta,
                gamma=request.gamma
            )
            
            ml_predicted = predict_congestion({
                "distance_km": summary["length_m"] / 1000.0,
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": summary["delay_s"]
            })
            
            congestion_ratio = (
                summary["travel_time_s"] / summary["no_traffic_s"]
                if summary["no_traffic_s"] and summary["no_traffic_s"] > 0
                else None
            )
            
            calculated_delay = 0
            if summary["travel_time_s"] and summary["no_traffic_s"]:
                calculated_delay = max(0, summary["travel_time_s"] - summary["no_traffic_s"])
            elif summary.get("delay_s"):
                calculated_delay = summary["delay_s"]
            
            analyzed_route = {
                "route_index": idx,
                "travel_time_s": summary["travel_time_s"],
                "no_traffic_s": summary["no_traffic_s"],
                "delay_s": calculated_delay,
                "length_m": summary["length_m"],
                "congestion_ratio": congestion_ratio,
                "calculated_cost": cost,
                "ml_predicted_congestion": ml_predicted
            }
            analyzed_routes.append(analyzed_route)
            
            # Save to database
            try:
                session = get_session()
                save_analysis(session, {
                    "route_id": f"{route_id}_route{idx}",
                    "origin": origin_data,
                    "destination": dest_data,
                    "travel_time_s": summary["travel_time_s"],
                    "no_traffic_s": summary["no_traffic_s"],
                    "delay_s": summary["delay_s"],
                    "length_m": summary["length_m"],
                    "calculated_cost": cost,
                    "ml_predicted": ml_predicted,
                    "raw_json": route,
                    "user_id": current_user.id if current_user else None
                })
                session.close()
            except Exception as e:
                logger.error(f"Database save error: {e}")
        
        best_route = min(analyzed_routes, key=lambda x: x["calculated_cost"])
        
        return {
            "origin": origin_data,
            "destination": dest_data,
            "route_id": route_id,
            "analyzed_routes": analyzed_routes,
            "best_route_index": best_route["route_index"],
            "best_route": best_route,
            "timestamp": datetime.now(UTC).isoformat()
        }
        
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        logger.error(f"Route refresh error: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Route refresh failed: {str(e)}")


# ============================================================================
# USER AUTHENTICATION & MANAGEMENT
# ============================================================================

@app.post("/api/auth/register", response_model=UserResponse)
async def register_user(user_data: UserCreate, db: Session = Depends(get_session)):
    """Register a new user."""
    try:
        if len(user_data.password) < 6:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Password must be at least 6 characters long"
            )
        
        MAX_PASSWORD_LENGTH = 10000
        if len(user_data.password) > MAX_PASSWORD_LENGTH:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Password is too long. Maximum {MAX_PASSWORD_LENGTH} characters allowed."
            )

        requested_role = (user_data.role or "user").strip().lower()
        department_map = {
            "user": "general",
            "police_supervisor": "police",
            "logistics_manager": "logistics",
        }
        if requested_role not in department_map:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Invalid role selected during registration",
            )
        
        user_payload = user_data.dict()
        user_payload["department"] = department_map[requested_role]
        user_payload.pop("role", None)

        user = create_user(db, AuthUserCreate(**user_payload))
        return UserResponse.model_validate(user)
    except HTTPException:
        raise
    except Exception as e:
        error_msg = str(e)
        if "password" in error_msg.lower():
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Password validation error: {error_msg}"
            )
        raise HTTPException(status_code=500, detail=f"Registration failed: {error_msg}")


@app.get("/test/create-user")
async def create_test_user(db: Session = Depends(get_session)):
    """Create a test user for development/testing (remove in production)."""
    # Check if test user already exists
    existing_user = get_user_by_username(db, "testuser")
    if existing_user:
        return {
            "status": "user_exists",
            "message": "Test user already exists",
            "username": "testuser",
            "password": "password123"
        }
    
    # Create test user
    try:
        test_user_data = AuthUserCreate(
            email="testuser@example.com",
            username="testuser",
            password="password123",
            full_name="Test User",
            department="general"
        )
        user = create_user(db, test_user_data)
        return {
            "status": "created",
            "message": "Test user created successfully",
            "username": "testuser",
            "password": "password123",
            "email": "testuser@example.com"
        }
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create test user: {str(e)}"
        )


@app.post("/api/auth/login", response_model=RoleToken)
async def login_user(
    form_data: OAuth2PasswordRequestForm = Depends(),
    db: Session = Depends(get_session),
    *,
    response: Response,
):
    """Login and get access token."""
    user = authenticate_user(db, form_data.username, form_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token_expires = timedelta(minutes=30 * 24 * 60)
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    user_role = (getattr(user, "role", None) or "user").strip().lower()
    user_district_id = getattr(user, "district_id", None)

    role = UserRole.user
    district_id = None
    fleet_zone = None
    if user.is_admin:
        role = UserRole.admin
    elif user_role == "police_supervisor":
        role = UserRole.police_supervisor
        district_id = user_district_id or "district_1"
    elif user_department == "police":
        role = UserRole.police_supervisor
        district_id = user_district_id or "district_1"
    elif user_role == "logistics_manager":
        role = UserRole.logistics_manager
        fleet_zone = "zone_default"
    elif user_department == "logistics":
        role = UserRole.logistics_manager
        fleet_zone = "zone_default"

    access_token = create_role_access_token(
        username=user.username,
        role=role,
        district_id=district_id,
        fleet_zone=fleet_zone,
        expires_delta=access_token_expires,
    )
    # Set secure=False for localhost/HTTP development, True for production/HTTPS
    is_secure = os.getenv("SECURE_COOKIES", "false").lower() == "true"
    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=is_secure,
        samesite="Lax",  # Use Lax for both secure and non-secure (None requires Secure=True)
        max_age=30 * 24 * 60 * 60,
        path="/",
    )
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "role": role,
        "district_id": district_id,
        "fleet_zone": fleet_zone,
    }


@app.get("/police/dashboard", response_class=HTMLResponse)
async def police_dashboard(request: Request, current_user: dict = Depends(require_police_department_user())):
    """Render the police supervisor dashboard for the user's district."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for police dashboard access",
        )

    context = _build_police_dashboard_context(current_user, district_id)
    context.update({
        "request": request,
        "current_user": current_user,
        "district_id": district_id,
        "google_maps_api_key": os.getenv("GOOGLE_MAPS_API_KEY", ""),
        "ml_predictions": _predict_police_hotspots(district_id, context["incidents"]),
        "data_error": False,
    })

    return templates.TemplateResponse(
        request=request,
        name="police/supervisor.html",
        context=context,
    )


@app.get("/supervisor/analytics", response_class=HTMLResponse)
async def supervisor_analytics_page(request: Request, current_user: dict = Depends(require_police_department_user())):
    """Render the supervisor analytics dashboard page."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for supervisor analytics",
        )

    return templates.TemplateResponse(
        request=request,
        name="police/supervisor_analytics.html",
        context={
            "request": request,
            "current_user": current_user,
            "district_id": district_id,
        },
    )


@app.get("/api/supervisor/analytics")
async def supervisor_analytics_data(
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_session),
):
    """Return analytics series for supervisor command center dashboards."""
    now = datetime.now(UTC)
    district_id = current_user.get("district_id") or ""

    seven_days_ago = now - timedelta(days=6)
    recent_feedback = (
        db.query(MLFeedback)
        .filter(MLFeedback.created_at >= seven_days_ago)
        .order_by(MLFeedback.created_at.asc())
        .all()
    )

    day_labels = [(seven_days_ago + timedelta(days=i)).date().isoformat() for i in range(7)]
    incident_volume_map: dict[str, dict[str, int]] = {day: {} for day in day_labels}
    for row in recent_feedback:
        day_key = (row.created_at or now).date().isoformat()
        if day_key not in incident_volume_map:
            continue
        series_key = f"{(row.severity or 'unknown').lower()} | {row.zone or 'Unknown Zone'}"
        incident_volume_map[day_key][series_key] = incident_volume_map[day_key].get(series_key, 0) + 1

    all_volume_keys = sorted({
        key
        for day_data in incident_volume_map.values()
        for key in day_data.keys()
    })
    incident_volume_datasets = [
        {
            "label": key,
            "data": [incident_volume_map[day].get(key, 0) for day in day_labels],
        }
        for key in all_volume_keys
    ]

    weekly_labels: list[str] = []
    weekly_values: list[float] = []
    for week_index in range(3, -1, -1):
        start = (now - timedelta(days=now.weekday())) - timedelta(weeks=week_index)
        end = start + timedelta(days=7)
        rows = (
            db.query(MLFeedback)
            .filter(MLFeedback.created_at >= start, MLFeedback.created_at < end)
            .all()
        )
        avg_response = round(sum(float(item.response_time_minutes) for item in rows) / len(rows), 2) if rows else 0.0
        weekly_labels.append(start.date().isoformat())
        weekly_values.append(avg_response)

    month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    monthly_workloads = (
        db.query(OfficerIncidentCount)
        .filter(OfficerIncidentCount.last_updated >= month_start)
        .all()
    )
    utilization_map: dict[str, int] = {}
    for item in monthly_workloads:
        total = int(item.incident_count_critical + item.incident_count_high + item.incident_count_medium + item.incident_count_low)
        utilization_map[item.officer_name or item.officer_username] = utilization_map.get(item.officer_name or item.officer_username, 0) + total
    sorted_utilization = sorted(utilization_map.items(), key=lambda pair: pair[1], reverse=True)

    zone_counts: dict[str, int] = {}
    for item in recent_feedback:
        zone_name = item.zone or "Unknown Zone"
        zone_counts[zone_name] = zone_counts.get(zone_name, 0) + 1
    top_zones = sorted(zone_counts.items(), key=lambda pair: pair[1], reverse=True)[:5]

    latest_retrain = (
        db.query(MLRetrainAudit)
        .order_by(MLRetrainAudit.retrained_at.desc())
        .first()
    )

    return {
        "district_id": district_id,
        "incident_volume": {
            "labels": day_labels,
            "datasets": incident_volume_datasets,
        },
        "response_time_trend": {
            "labels": weekly_labels,
            "data": weekly_values,
        },
        "officer_utilization": {
            "labels": [name for name, _ in sorted_utilization],
            "data": [count for _, count in sorted_utilization],
        },
        "top_zones": [
            {"zone": zone_name, "count": count}
            for zone_name, count in top_zones
        ],
        "last_retrained": latest_retrain.retrained_at.isoformat() if latest_retrain else None,
        "updated_at": now.isoformat(),
    }


@app.get("/police/units/live")
async def live_patrol_units(current_user: dict = Depends(require_police_department_user())):
    """Return the latest patrol unit status data for the supervisor dashboard."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for patrol unit status",
        )

    incidents = _load_police_incidents(district_id)
    supervisor_name = current_user.get("username", "Unknown Supervisor")
    assignments = _get_dispatch_assignments(district_id)
    patrol_units = _build_patrol_units(district_id, incidents, supervisor_name, assignments)
    assigned_incident_ids = {assignment.incident_id for assignment in assignments}
    unassigned_incidents = [incident for incident in incidents if incident.get("id") not in assigned_incident_ids]
    available_patrol_units = [unit for unit in patrol_units if unit["status"] == "available"]

    return {
        "patrol_units": patrol_units,
        "available_patrol_units": available_patrol_units,
        "unassigned_incidents": unassigned_incidents,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/api/incidents")
async def api_incidents(current_user: dict = Depends(require_any_role("police_supervisor", "police_officer"))):
    """Return active incidents using command-center JSON contract."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incidents",
        )

    raw_incidents = _load_police_incidents(district_id)
    incidents = [
        {
            "id": incident.get("id"),
            "title": incident.get("type") or "traffic",
            "location": incident.get("description") or "Unknown location",
            "severity": incident.get("severity") or "unknown",
            "time": incident.get("start_time") or datetime.now(UTC).isoformat(),
            "lat": incident.get("latitude"),
            "lng": incident.get("longitude"),
            "backup_required": str(incident.get("severity") or "").lower() in {"critical", "high"},
            "notes": "",
        }
        for incident in raw_incidents
    ]

    return {
        "incidents": incidents,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/api/officers/status")
async def api_officers_status(current_user: dict = Depends(require_any_role("police_supervisor", "police_officer", "admin"))):
    """Return officer status rows using command-center JSON contract."""
    district_id = current_user.get("district_id")
    # For admin users, default to district_1 if not provided
    if not district_id and current_user.get("role") == "admin":
        district_id = "district_1"
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for officer status",
        )

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    patrol_units = _build_patrol_units(
        district_id,
        incidents,
        current_user.get("username", "Unknown Supervisor"),
        assignments,
    )

    officers = [
        {
            "id": unit.get("unit_id"),
            "name": unit.get("officer_name") or "Unknown Officer",
            "badge": unit.get("unit_id"),
            "status": str(unit.get("status") or "available").lower(),
            "skills": [],
            "district_id": unit.get("district_id") or district_id,
            "current_location": unit.get("current_location") or "District duty",
            "latitude": unit.get("latitude"),
            "longitude": unit.get("longitude"),
            "last_updated": unit.get("last_updated") or "-",
        }
        for unit in patrol_units
    ]

    return {
        "officers": officers,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.post("/api/incident/new")
async def create_incident(
    request: NewIncidentRequest,
    current_user: dict = Depends(require_police_department_user()),
):
    """Create a new district incident and push a real-time update over Socket.IO."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident creation",
        )

    incident_id = f"manual-{uuid.uuid4().hex[:10]}"
    zone_name = request.zone or _infer_zone_name(
        district_id,
        {"latitude": request.latitude, "longitude": request.longitude, "description": request.description},
    )
    incident_record = {
        "id": incident_id,
        "type": request.incident_type,
        "severity": _normalize_severity(request.severity),
        "description": request.description,
        "latitude": request.latitude,
        "longitude": request.longitude,
        "zone": zone_name,
        "start_time": datetime.now(UTC).isoformat(),
        "district_id": district_id,
    }
    _manual_incidents_store.setdefault(district_id, []).append(incident_record)

    await emit_incident_new(
        sio,
        district_id,
        incident_record,
        actor=current_user.get("username", "Unknown Supervisor"),
    )
    await emit_incident_updated(
        sio,
        district_id,
        incident_record,
        update_type="created",
        actor=current_user.get("username", "Unknown Supervisor"),
    )

    return {
        "status": "created",
        "incident": incident_record,
    }


@app.get("/police/incidents/feed")
async def live_incidents_feed(current_user: dict = Depends(require_police_department_user())):
    """Return district-scoped incident feed sorted by reported time (desc)."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident feed",
        )

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    incidents_feed = _build_incidents_feed(district_id, incidents, assignments)

    return {
        "incidents": incidents_feed,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/police/heatmap/data")
async def police_heatmap_data(current_user: dict = Depends(require_police_department_user())):
    """Return district-scoped incident heatmap points for the supervisor."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for heatmap data",
        )

    incidents = _load_police_incidents(district_id)
    return _build_heatmap_points(incidents)


@app.get("/police/response-times")
@app.get("/police/response_times")
@app.get("/api/police/response-times")
async def police_response_times(current_user: dict = Depends(require_police_department_user())):
    """Return zone-wise average response metrics for today's shift."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for response time metrics",
        )

    incidents = _load_police_incidents(district_id)
    assignments = _get_dispatch_assignments(district_id)
    target_threshold_minutes = 8.0
    zone_metrics = _build_response_time_by_zone(
        district_id=district_id,
        incidents=incidents,
        assignments=assignments,
        target_threshold_minutes=target_threshold_minutes,
    )

    return {
        "district_id": district_id,
        "target_threshold_minutes": target_threshold_minutes,
        "zones": zone_metrics,
        "updated_at": datetime.now(UTC).isoformat(),
    }


@app.get("/api/logistics/alerts")
async def logistics_shared_alerts(current_user: dict = Depends(require_role("logistics_manager"))):
    """Return active sanitized police zone alerts for logistics users only."""
    session = get_session()
    try:
        now = datetime.now(UTC)
        rows = (
            session.query(SharedAlert)
            .filter(SharedAlert.expires_at >= now)
            .order_by(SharedAlert.timestamp.desc())
            .limit(25)
            .all()
        )
    finally:
        session.close()

    return {
        "alerts": [
            {
                "alert_id": row.alert_id,
                "zone": row.zone,
                "severity": row.severity,
                "timestamp": row.timestamp.isoformat() if row.timestamp else None,
                "affected_roads": row.affected_roads or [],
                "expires_at": row.expires_at.isoformat() if row.expires_at else None,
            }
            for row in rows
        ],
        "updated_at": datetime.now(UTC).isoformat(),
    }


def _ensure_officer_statuses_initialized(district_id: str):
    """
    Ensure all mock officers in a district have 'available' status in database.
    This syncs the OfficerDispatchStatus table with the mock units being displayed.
    Resets all officers to 'available' on each request to ensure consistency with mock data.
    """
    session = get_session()
    try:
        # Generate unit IDs based on mock data
        incidents = _load_police_incidents(district_id)
        police_users = session.query(User).filter(User.department == "police").count()
        unit_count = max(4, len(incidents), police_users)
        
        for index in range(unit_count):
            unit_id = f"{district_id.upper().replace('_', '-')}-U{index + 1:02d}"
            
            # Check if officer status exists in database
            officer_status = session.query(OfficerDispatchStatus).filter(
                OfficerDispatchStatus.officer_id == unit_id,
                OfficerDispatchStatus.district_id == district_id
            ).first()
            
            if officer_status is None:
                # Create new entry with 'available' status
                officer_status = OfficerDispatchStatus(
                    district_id=district_id,
                    officer_id=unit_id,
                    status="available",
                    assigned_incident_id=None,
                    mobile_token=None
                )
                session.add(officer_status)
            else:
                # Reset all officers to 'available' to match mock data display
                # This ensures what's shown in the UI matches what the backend validates
                officer_status.status = "available"
                officer_status.assigned_incident_id = None
                officer_status.updated_at = datetime.now(UTC)
            
        session.commit()
        logger.debug(f"Initialized {unit_count} officer statuses for {district_id}")
    except Exception as e:
        logger.warning(f"Could not initialize officer statuses for {district_id}: {e}")
        session.rollback()
    finally:
        session.close()


@app.post("/api/dispatch")
@app.post("/police/dispatch")
async def dispatch_patrol_unit(
    request: DispatchRequest,
    current_user: dict = Depends(require_police_department_user()),
):
    """Assign a patrol unit to an incident and persist the dispatch.
    
    Validates officer availability, updates dispatch status, records in audit log,
    sends FCM push notification, and broadcasts via SocketIO.
    
    Request:
        incident_id (str): ID of the incident to dispatch
        officer_id (str): ID of the officer/unit to dispatch
    
    Response:
        {
            "success": true,
            "dispatch_id": 42,
            "eta": 5,
            "incident_id": "incident_123",
            "officer_id": "unit_001"
        }
    
    Raises:
        400: Missing district_id, incident_id, or officer_id
        403: User is not a police supervisor
        404: Incident not found
        409: Officer not available or already assigned
        500: Dispatch failed
    """
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for dispatch",
        )

    # Ensure officer statuses are properly initialized before dispatch
    _ensure_officer_statuses_initialized(district_id)

    incidents = _load_police_incidents(district_id)
    incident_map = {incident.get("id"): incident for incident in incidents}
    target_incident = incident_map.get(request.incident_id)
    if not target_incident:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Incident not found",
        )

    officer_id = (request.officer_id or request.unit_id or "").strip()
    if not officer_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="officer_id is required for dispatch",
        )

    session = get_session()
    dispatch_log_record = None
    officer_mobile_token = None
    
    try:
        now = datetime.now(UTC)
        
        # Clear session cache to avoid stale data
        session.expunge_all()
        
        # Check officer status before dispatch
        officer_status = (
            session.query(OfficerDispatchStatus)
            .filter(OfficerDispatchStatus.officer_id == officer_id)
            .with_for_update()
            .first()
        )
        
        if officer_status is None:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Officer {officer_id} not found in dispatch system",
            )
        
        # Validate officer status is 'available' before dispatch
        officer_status_lower = (officer_status.status or "").lower()
        if officer_status_lower != "available":
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=f"Officer is not available (current status: {officer_status_lower})",
            )
        
        # Store mobile token for FCM notification
        officer_mobile_token = officer_status.mobile_token
        
        # Query fresh data from database and delete ALL existing assignments for this incident and unit
        session.query(PoliceDispatchAssignment).filter(
            PoliceDispatchAssignment.incident_id == request.incident_id
        ).delete(synchronize_session=False)
        
        session.query(PoliceDispatchAssignment).filter(
            PoliceDispatchAssignment.unit_id == officer_id
        ).delete(synchronize_session=False)
        
        session.commit()
        session.expunge_all()
        
        # Create new assignment
        assignment = PoliceDispatchAssignment(
            district_id=district_id,
            incident_id=request.incident_id,
            unit_id=officer_id,
            assigned_by=current_user.get("username", "Unknown Supervisor"),
            assigned_at=now,
            status="active",
        )
        session.add(assignment)

        # Update officer status to 'enroute'
        officer_status.status = "enroute"
        officer_status.assigned_incident_id = request.incident_id
        officer_status.updated_at = now
        session.add(officer_status)

        # Create dispatch log record
        dispatch_log_record = DispatchLog(
            district_id=district_id,
            incident_id=request.incident_id,
            officer_id=officer_id,
            assigned_by=current_user.get("username", "Unknown Supervisor"),
            assigned_at=now,
            status="dispatched",
        )
        session.add(dispatch_log_record)

        shared_alert = _create_shared_alert_for_dispatch(session, district_id, target_incident)
        session.commit()
        session.refresh(assignment)
        session.refresh(dispatch_log_record)
        if shared_alert is not None:
            session.refresh(shared_alert)
            
    except HTTPException:
        session.rollback()
        raise
    except IntegrityError as exc:
        session.rollback()
        existing_assignment = (
            session.query(PoliceDispatchAssignment)
            .filter(PoliceDispatchAssignment.incident_id == request.incident_id)
            .first()
        )
        if existing_assignment and existing_assignment.unit_id == officer_id:
            logger.warning(
                "Dispatch assignment already exists for incident %s and unit %s",
                request.incident_id,
                officer_id,
            )
            assignment = existing_assignment
        else:
            logger.error("Dispatch assignment integrity error: %s", exc)
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail="Incident or patrol unit is already assigned",
            )
    except Exception as exc:
        session.rollback()
        logger.error("Dispatch assignment failed: %s", exc)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Dispatch assignment failed",
        )
    finally:
        session.close()

    # Calculate ETA in minutes (simplified: use distance if available)
    eta_minutes = 5  # Default ETA in minutes
    supervisor_id = current_user.get("username", "Unknown Supervisor")
    
    # Send FCM push notification to officer's mobile device
    fcm_result = send_officer_dispatch_notification(
        officer_mobile_token,
        officer_id,
        target_incident,
    )
    logger.info(
        f"Dispatch notification for officer {officer_id}: "
        f"sent={fcm_result.get('sent')}, reason={fcm_result.get('reason')}"
    )
    
    # Build dispatch data for SocketIO event
    dispatch_data = {
        "dispatch_id": dispatch_log_record.id if dispatch_log_record else None,
        "officer_id": officer_id,
        "incident_id": request.incident_id,
        "eta": eta_minutes,
        "supervisor_id": supervisor_id,
    }
    
    # Emit SocketIO event to police namespace
    await emit_officer_dispatched(
        sio,
        district_id,
        dispatch_data,
        actor=supervisor_id,
    )
    
    # Also emit the existing events for backward compatibility
    updated_context = _build_police_dashboard_context(current_user, district_id)
    updated_unit = next((unit for unit in updated_context["patrol_units"] if unit["unit_id"] == officer_id), None)
    updated_incident = next((incident for incident in updated_context["incidents"] if incident.get("id") == request.incident_id), None)

    await emit_incident_updated(
        sio,
        district_id,
        updated_incident or {"id": request.incident_id},
        update_type="dispatched",
        actor=supervisor_id,
    )
    await emit_officer_status_changed(
        sio,
        district_id,
        {
            "id": officer_id,
            "name": (updated_unit or {}).get("officer_name") if isinstance(updated_unit, dict) else None,
            "badge": officer_id,
            "status": "enroute",
            "district_id": district_id,
            "incident_id": request.incident_id,
        },
        actor=supervisor_id,
    )

    # Return simplified response format
    return {
        "success": True,
        "dispatch_id": dispatch_log_record.id if dispatch_log_record else None,
        "eta": eta_minutes,
        "incident_id": request.incident_id,
        "officer_id": officer_id,
        "fcm_sent": fcm_result.get("sent", False),
    }


@app.get("/police/export/pptx")
async def police_export_pptx(current_user: dict = Depends(require_police_department_user())):
    """Generate a police supervisor shift report as a PPTX file."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for police report export",
        )

    try:
        officer_name = current_user.get("username", "Unknown Officer")
        incidents = _load_police_incidents(district_id)
        district_summary = _build_district_summary(incidents)
        ml_predictions = _predict_police_hotspots(district_id, incidents)
        output = generate_shift_pptx(district_id, officer_name, incidents, district_summary, ml_predictions)

        filename = f"ShiftReport_District{district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except Exception as exc:
        logger.error(f"Police PPTX export failed: {exc}")
        raise HTTPException(
            status_code=500,
            detail="Report generation failed. Please try again or contact admin.",
        )


def generate_detailed_shift_report_pptx(
    shift_id: int,
    shift_data: dict,
    incidents: list[dict],
    officer_workload: list[dict],
    dispatch_assignments: list[dict],
) -> io.BytesIO:
    """Generate a comprehensive shift report PPTX with 5 slides.
    
    Slides:
    1. Shift summary (date, supervisor, total incidents, avg response time)
    2. Incident breakdown table (ID, type, severity, location, officer, resolved)
    3. Officer performance (name, incidents handled, avg response time)
    4. Unresolved/escalated cases
    5. Zone heatmap / incident distribution
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    # ========== SLIDE 1: SHIFT SUMMARY ==========
    slide1 = prs.slides.add_slide(title_layout)
    slide1.shapes.title.text = "Shift Summary Report"
    
    shift_start = shift_data.get("start_time")
    shift_end = shift_data.get("end_time")
    if isinstance(shift_start, str):
        shift_start = datetime.fromisoformat(shift_start.replace('Z', '+00:00'))
    if isinstance(shift_end, str):
        shift_end = datetime.fromisoformat(shift_end.replace('Z', '+00:00'))
    
    duration_mins = 0
    if shift_start and shift_end:
        duration_mins = int((shift_end - shift_start).total_seconds() / 60)
    
    avg_response_time = "N/A"
    if dispatch_assignments:
        response_times = [a.get("response_time_s", 0) for a in dispatch_assignments if a.get("response_time_s")]
        if response_times:
            avg_response_time = f"{sum(response_times) / len(response_times):.1f}s"
    
    unresolved_count = sum(1 for inc in incidents if inc.get("status") != "resolved")
    
    summary_text = (
        f"District: {shift_data.get('district_id', 'Unknown')}\n"
        f"Supervisor: {shift_data.get('supervisor_name', 'Unknown')}\n"
        f"Date: {shift_start.strftime('%Y-%m-%d') if shift_start else 'N/A'}\n"
        f"Shift Duration: {duration_mins // 60}h {duration_mins % 60}m\n"
        f"\n"
        f"Total Incidents: {len(incidents)}\n"
        f"Resolved: {len(incidents) - unresolved_count}\n"
        f"Unresolved: {unresolved_count}\n"
        f"Avg Response Time: {avg_response_time}\n"
        f"Officers on Duty: {shift_data.get('officers_on_duty', 0)}"
    )
    
    slide1.placeholders[1].text = summary_text

    # ========== SLIDE 2: INCIDENT BREAKDOWN TABLE ==========
    slide2 = prs.slides.add_slide(blank_layout)
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Incident Breakdown"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True

    # Create incident table
    rows = min(len(incidents) + 1, 12)  # Cap at 12 rows for slide fit
    cols = 6
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.4), Inches(1), Inches(9.2), Inches(6))
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.2)  # ID
    table.columns[1].width = Inches(1.5)  # Type
    table.columns[2].width = Inches(1.2)  # Severity
    table.columns[3].width = Inches(2)    # Location
    table.columns[4].width = Inches(1.5)  # Officer
    table.columns[5].width = Inches(1.8)  # Status
    
    # Header row
    headers = ["ID", "Type", "Severity", "Location", "Officer", "Status"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)

    # Data rows
    for row_idx, incident in enumerate(incidents[:rows - 1], start=1):
        table.cell(row_idx, 0).text = _safe_ppt_text(incident.get("id", "N/A")[:10])
        table.cell(row_idx, 1).text = _safe_ppt_text(incident.get("type", "N/A")[:15])
        
        severity = incident.get("severity", "N/A")
        table.cell(row_idx, 2).text = severity
        
        table.cell(row_idx, 3).text = _safe_ppt_text(incident.get("location", "N/A")[:20])
        
        # Get officer from workload or assignment
        officer = "Unassigned"
        for assign in dispatch_assignments:
            if assign.get("incident_id") == incident.get("id"):
                officer = assign.get("unit_id", "Unassigned")
                break
        table.cell(row_idx, 4).text = officer[:15]
        
        status = incident.get("status", "pending")
        table.cell(row_idx, 5).text = status.capitalize()

    # ========== SLIDE 3: OFFICER PERFORMANCE ==========
    slide3 = prs.slides.add_slide(blank_layout)
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Officer Performance"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True

    # Create officer performance table
    perf_rows = min(len(officer_workload) + 1, 10)
    perf_cols = 5
    perf_table_shape = slide3.shapes.add_table(perf_rows, perf_cols, Inches(0.4), Inches(1), Inches(9.2), Inches(6))
    perf_table = perf_table_shape.table
    
    perf_table.columns[0].width = Inches(2)    # Officer Name
    perf_table.columns[1].width = Inches(1.5)  # Total
    perf_table.columns[2].width = Inches(1.5)  # Critical
    perf_table.columns[3].width = Inches(1.8)  # High
    perf_table.columns[4].width = Inches(1.4)  # Avg Time

    # Header
    perf_headers = ["Officer", "Total", "Critical", "High", "Needs Rotation"]
    for col_idx, header in enumerate(perf_headers):
        cell = perf_table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)

    # Officer data
    for row_idx, officer in enumerate(officer_workload[:perf_rows - 1], start=1):
        perf_table.cell(row_idx, 0).text = _safe_ppt_text(officer.get("officer_name", "N/A")[:20])
        perf_table.cell(row_idx, 1).text = str(officer.get("total_incidents", 0))
        perf_table.cell(row_idx, 2).text = str(officer.get("critical_incidents", 0))
        perf_table.cell(row_idx, 3).text = str(officer.get("high_incidents", 0))
        perf_table.cell(row_idx, 4).text = "⚠️ YES" if officer.get("needs_rotation") else "No"

    # ========== SLIDE 4: UNRESOLVED/ESCALATED CASES ==========
    slide4 = prs.slides.add_slide(blank_layout)
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Unresolved & Escalated Cases"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True

    unresolved_incidents = [inc for inc in incidents if inc.get("status") != "resolved"]
    
    if unresolved_incidents:
        text_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for idx, incident in enumerate(unresolved_incidents[:8]):  # Limit to 8 for readability
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = (
                f"• ID: {incident.get('id', 'N/A')} | "
                f"Type: {incident.get('type', 'N/A')} | "
                f"Severity: {incident.get('severity', 'N/A')} | "
                f"Status: {incident.get('status', 'N/A')}"
            )
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(11)
    else:
        text_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
        text_frame = text_box.text_frame
        text_frame.text = "✓ All incidents have been resolved. No escalations required."
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)

    # ========== SLIDE 5: ZONE HEATMAP / INCIDENT DISTRIBUTION ==========
    slide5 = prs.slides.add_slide(blank_layout)
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Incident Distribution by Zone"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True

    # Group incidents by location/zone
    zone_incidents = {}
    for incident in incidents:
        zone = incident.get("location", "Unknown")
        zone_incidents[zone] = zone_incidents.get(zone, 0) + 1
    
    if zone_incidents:
        text_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for idx, (zone, count) in enumerate(sorted(zone_incidents.items(), key=lambda x: x[1], reverse=True)[:12]):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            bar = "█" * min(count, 20)
            p.text = f"{zone[:30]:.<30} {bar} ({count})"
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.name = "Courier New"
    else:
        text_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        text_frame = text_box.text_frame
        text_frame.text = "No incident data available."
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(18)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@app.post("/api/shift/report")
async def generate_shift_report(
    shift_id: int = Query(..., description="Shift ID to generate report for"),
    current_user: dict = Depends(require_role("police_supervisor")),
    db: Session = Depends(get_session),
):
    """Generate a comprehensive shift report as PPTX with 5 slides.
    
    Pulls data from:
    - Shift (supervisor, district, times)
    - ShiftAttendance (officers on duty)
    - DispatchLog (immutable dispatch audit trail)
    - OfficerIncidentCount (officer workload)
    """
    try:
        # Get shift data
        shift = db.query(Shift).filter(Shift.id == shift_id).first()
        if not shift:
            raise HTTPException(status_code=404, detail="Shift not found")

        # Verify the supervisor owns this shift
        if shift.supervisor_id != current_user.get("username"):
            raise HTTPException(status_code=403, detail="You do not have permission to access this shift's report")

        # Get all incidents for this shift (notifications within shift time window)
        incidents_query = db.query(Notification).filter(
            Notification.timestamp >= shift.start_time,
            Notification.timestamp <= (shift.end_time or datetime.now(UTC))
        ).all()

        incidents = [
            {
                "id": inc.id,
                "type": inc.type,
                "title": inc.title,
                "message": inc.message,
                "severity": inc.type.split("_")[-1] if "_" in inc.type else "medium",
                "location": inc.message[:50] if inc.message else "Unknown",
                "status": "resolved" if inc.is_read else "pending",
                "timestamp": inc.created_at.isoformat() if inc.created_at else "",
            }
            for inc in incidents_query
        ]

        # Get dispatch records from immutable dispatch log for this shift time window
        dispatch_log_entries = db.query(DispatchLog).filter(
            DispatchLog.assigned_at >= shift.start_time,
            DispatchLog.assigned_at <= (shift.end_time or datetime.now(UTC)),
            DispatchLog.district_id == shift.district_id
        ).all()

        dispatch_data = [
            {
                "incident_id": dl.incident_id,
                "unit_id": dl.officer_id,
                "assigned_by": dl.assigned_by,
                "assigned_at": dl.assigned_at.isoformat() if dl.assigned_at else "",
                "status": dl.status,
                "response_time_s": 0,  # Can be calculated if timestamp data is available
            }
            for dl in dispatch_log_entries
        ]

        # Get officer workload for this shift
        officer_workload = db.query(OfficerIncidentCount).filter(
            OfficerIncidentCount.shift_id == shift_id
        ).all()

        workload_data = [
            {
                "officer_username": ow.officer_username,
                "officer_name": ow.officer_name,
                "total_incidents": ow.incident_count_critical + ow.incident_count_high + 
                                   ow.incident_count_medium + ow.incident_count_low,
                "critical_incidents": ow.incident_count_critical,
                "high_incidents": ow.incident_count_high,
                "medium_incidents": ow.incident_count_medium,
                "low_incidents": ow.incident_count_low,
                "needs_rotation": ow.needs_rotation,
            }
            for ow in officer_workload
        ]

        # Prepare shift data dict
        shift_dict = {
            "shift_id": shift.id,
            "district_id": shift.district_id,
            "supervisor_name": shift.supervisor_name,
            "supervisor_id": shift.supervisor_id,
            "start_time": shift.start_time,
            "end_time": shift.end_time,
            "status": shift.status,
            "officers_on_duty": shift.officers_on_duty,
            "incidents_count": len(incidents),
        }

        # Generate PPTX
        output = generate_detailed_shift_report_pptx(
            shift_id=shift_id,
            shift_data=shift_dict,
            incidents=incidents,
            officer_workload=workload_data,
            dispatch_assignments=dispatch_data,
        )

        filename = f"ShiftReport_{shift.district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"

        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"',
                "Cache-Control": "no-cache, no-store, must-revalidate",
                "Pragma": "no-cache",
                "Expires": "0",
            },
        )
    except HTTPException:
        raise
    except Exception as exc:
        logger.error(f"Shift report generation failed: {exc}")
        raise HTTPException(
            status_code=500,
            detail=f"Report generation failed: {str(exc)}",
        )


@app.get("/police/shift/status")
async def get_shift_status(current_user: dict = Depends(require_role("police_supervisor"))):
    """Get the current active shift status for the supervisor."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            return {
                "status": "no_active_shift",
                "shift": None,
                "message": "No active shift found. Create a new shift to begin.",
            }
        
        attendance_records = (
            session.query(ShiftAttendance)
            .filter(ShiftAttendance.shift_id == active_shift.id)
            .all()
        )
        
        officers_on_duty = [
            {
                "officer_username": record.officer_username,
                "officer_name": record.officer_name,
                "clock_in_time": record.clock_in_time.isoformat() if record.clock_in_time else None,
                "clock_out_time": record.clock_out_time.isoformat() if record.clock_out_time else None,
                "status": record.status,
            }
            for record in attendance_records
            if record.status == "present"
        ]
        
        return {
            "status": "active",
            "shift": {
                "id": active_shift.id,
                "district_id": active_shift.district_id,
                "supervisor_id": active_shift.supervisor_id,
                "supervisor_name": active_shift.supervisor_name,
                "start_time": active_shift.start_time.isoformat(),
                "incidents_count": active_shift.incidents_count,
                "officers_on_duty": len(officers_on_duty),
                "notes": active_shift.notes,
            },
            "officers": officers_on_duty,
            "updated_at": datetime.now(UTC).isoformat(),
        }
    except Exception as exc:
        logger.error(f"Failed to get shift status: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to retrieve shift status",
        )
    finally:
        session.close()


@app.post("/police/shift/attendance")
async def mark_shift_attendance(
    request: ShiftAttendanceRequest,
    current_user: dict = Depends(require_role("police_supervisor")),
):
    """Mark officer attendance for the current shift."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="No active shift found",
            )
        
        for officer_data in request.officers:
            officer_username = officer_data.get("officer_username")
            officer_name = officer_data.get("officer_name")
            status_val = officer_data.get("status", "present")
            
            existing_record = (
                session.query(ShiftAttendance)
                .filter(
                    ShiftAttendance.shift_id == active_shift.id,
                    ShiftAttendance.officer_username == officer_username,
                )
                .first()
            )
            
            if existing_record:
                existing_record.status = status_val
                if status_val == "absent":
                    existing_record.clock_out_time = datetime.now(UTC)
            else:
                attendance_record = ShiftAttendance(
                    shift_id=active_shift.id,
                    officer_username=officer_username,
                    officer_name=officer_name,
                    status=status_val,
                    clock_in_time=datetime.now(UTC) if status_val == "present" else None,
                )
                session.add(attendance_record)
        
        session.commit()
        
        updated_present_count = (
            session.query(ShiftAttendance)
            .filter(
                ShiftAttendance.shift_id == active_shift.id,
                ShiftAttendance.status == "present",
            )
            .count()
        )
        active_shift.officers_on_duty = updated_present_count
        session.commit()
        
        return {
            "message": "Attendance marked successfully",
            "shift_id": active_shift.id,
            "officers_marked": len(request.officers),
            "officers_present": updated_present_count,
            "updated_at": datetime.now(UTC).isoformat(),
        }
    except HTTPException:
        session.rollback()
        raise
    except Exception as exc:
        session.rollback()
        logger.error(f"Failed to mark attendance: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to mark attendance",
        )
    finally:
        session.close()


@app.post("/police/shift/end")
async def end_shift(
    request: ShiftEndRequest,
    current_user: dict = Depends(require_role("police_supervisor")),
):
    """End the current shift, optionally generating and exporting a PPTX report."""
    district_id = current_user.get("district_id")
    username = current_user.get("username", "Unknown")
    
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    session = get_session()
    try:
        active_shift = (
            session.query(Shift)
            .filter(
                Shift.district_id == district_id,
                Shift.supervisor_id == username,
                Shift.status == "active",
            )
            .order_by(Shift.start_time.desc())
            .first()
        )
        
        if not active_shift:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="No active shift found to end",
            )
        
        active_shift.end_time = datetime.now(UTC)
        active_shift.status = "completed"
        active_shift.notes = request.notes
        session.commit()
        
        shift_end_result = {
            "message": "Shift ended successfully",
            "shift_id": active_shift.id,
            "start_time": active_shift.start_time.isoformat(),
            "end_time": active_shift.end_time.isoformat(),
            "status": "completed",
            "updated_at": datetime.now(UTC).isoformat(),
        }
        
        if request.export_pptx:
            try:
                incidents = _load_police_incidents(district_id)
                district_summary = _build_district_summary(incidents)
                ml_predictions = _predict_police_hotspots(district_id, incidents)
                output = generate_shift_pptx(district_id, username, incidents, district_summary, ml_predictions)
                
                filename = f"ShiftReport_District{district_id}_{datetime.now(UTC).strftime('%Y%m%d_%H%M')}.pptx"
                shift_end_result["pptx_filename"] = filename
                shift_end_result["export_status"] = "success"
                
                session.close()
                
                return StreamingResponse(
                    output,
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": f'attachment; filename="{filename}"'},
                )
            except Exception as export_exc:
                logger.warning(f"PPTX export during shift end failed: {export_exc}")
                shift_end_result["export_status"] = "failed"
                shift_end_result["export_error"] = str(export_exc)
                return shift_end_result
        else:
            return shift_end_result
    except HTTPException:
        session.rollback()
        raise
    except Exception as exc:
        session.rollback()
        logger.error(f"Failed to end shift: {exc}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Failed to end shift",
        )
    finally:
        session.close()


@app.get("/police/alerts/list")
async def get_alerts_list(current_user: dict = Depends(require_role("police_supervisor"))):
    """Get unread alerts for the supervisor's district."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    alerts = get_unread_alerts(district_id)
    return AlertListResponse(
        alerts=alerts,
        unread_count=len(alerts),
    )


@app.get("/police/alerts/stream")
async def alerts_stream(current_user: dict = Depends(require_role("police_supervisor"))):
    """Server-Sent Events stream for real-time police alerts."""
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required",
        )
    
    async def event_generator():
        """Generate SSE events for alerts."""
        # Send initial alert to inform client the stream is connected
        yield f"data: {json.dumps({'type': 'connected', 'message': 'Alert stream connected'})}\n\n"
        
        # Send all existing unread alerts
        existing_alerts = get_unread_alerts(district_id)
        for alert in existing_alerts:
            yield f"data: {json.dumps({'type': 'alert', 'data': alert.model_dump()})}\n\n"
            await asyncio.sleep(0.01)  # Small delay to prevent overwhelming client
        
        # Keep connection open for new alerts (heartbeat)
        last_check = time.time()
        last_alert_count = len(existing_alerts)
        
        while True:
            try:
                await asyncio.sleep(1)
                
                current_alerts = get_unread_alerts(district_id)
                current_count = len(current_alerts)
                
                # Check for new alerts since last check
                if current_count > last_alert_count:
                    new_alerts = current_alerts[last_alert_count:]
                    for alert in new_alerts:
                        yield f"data: {json.dumps({'type': 'alert', 'data': alert.model_dump()})}\n\n"
                    last_alert_count = current_count
                
                # Heartbeat every 30 seconds
                if time.time() - last_check > 30:
                    yield f": heartbeat\n\n"
                    last_check = time.time()
            except asyncio.CancelledError:
                logger.info(f"Alert stream closed for district {district_id}")
                break
            except Exception as exc:
                logger.error(f"Error in alert stream: {exc}")
                break
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.get("/api/incidents/predicted", response_model=PredictedIncidentResponse)
async def get_predicted_incidents(
    district_id: str = Query(..., description="District ID"),
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_session),
):
    """Get predicted high-risk incident zones as GeoJSON.
    
    Uses ML model predictions to identify zones likely to experience incidents
    based on time-of-day, day-of-week, and historical patterns.
    
    Returns GeoJSON FeatureCollection with risk zones.
    """
    user_district_id = current_user.get("district_id")
    if user_district_id and district_id != user_district_id:
        district_id = user_district_id

    incident_data = _load_police_incidents(district_id)

    # Get predictions from ML model
    try:
        hotspots = _predict_police_hotspots(district_id, incident_data)
    except Exception as e:
        logger.error(f"Error predicting hotspots: {e}")
        hotspots = []
    
    # Convert to GeoJSON
    features = []
    for hotspot in hotspots:
        # Classify risk level based on prediction score
        score = float(hotspot.get("likelihood_score", 0))
        if score >= 75:
            risk_level = "high"
        elif score >= 40:
            risk_level = "medium"
        else:
            risk_level = "low"
        
        # Get district center coordinates for creating polygon
        district_coords = DISTRICT_LOCATIONS.get(district_id, {})
        center_lat = district_coords.get("lat", 28.6139)
        center_lon = district_coords.get("lon", 77.2090)
        
        # Create circular polygon (approximate circle with 8 points)
        radius = 0.05 * (100 - score) / 100  # Radius decreases with higher risk
        lat = hotspot.get("latitude", center_lat)
        lon = hotspot.get("longitude", center_lon)
        
        # Generate 8-point circle polygon
        import math
        polygon_coords = []
        for i in range(8):
            angle = (i / 8) * 2 * math.pi
            poly_lat = lat + radius * math.cos(angle)
            poly_lon = lon + radius * math.sin(angle)
            polygon_coords.append([poly_lon, poly_lat])
        polygon_coords.append(polygon_coords[0])  # Close the polygon
        
        feature = {
            "type": "Feature",
            "geometry": {
                "type": "Polygon",
                "coordinates": [polygon_coords]
            },
            "properties": {
                "zone_name": hotspot.get("zone_name", hotspot.get("location", f"Zone_{lat:.4f}_{lon:.4f}")),
                "risk_level": risk_level,
                "prediction_score": round(score, 2),
                "incident_count": hotspot.get("incident_count", 0),
                "severity": hotspot.get("severity", "medium")
            }
        }
        features.append(feature)
    
    return PredictedIncidentResponse(
        type="FeatureCollection",
        features=features
    )


@app.get("/api/police/officer-workload", response_model=OfficerWorkloadResponse)
async def get_officer_workload(
    district_id: str = Query(..., description="District ID"),
    shift_id: Optional[int] = Query(None, description="Optional shift ID (uses current active shift if not provided)"),
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_session),
):
    """Get officer incident counts and rotation status for current shift.
    
    Returns officer workload data with rotation alerts for officers
    handling 3+ critical incidents.
    """
    # Determine shift to use
    if not shift_id:
        # Get current active shift for this district
        shift = db.query(Shift).filter(
            Shift.district_id == district_id,
            Shift.status == "active"
        ).order_by(Shift.start_time.desc()).first()
        
        if not shift:
            raise HTTPException(
                status_code=404,
                detail=f"No active shift found for district {district_id}"
            )
        shift_id = shift.id
    else:
        shift = db.query(Shift).filter(Shift.id == shift_id).first()
        if not shift:
            raise HTTPException(status_code=404, detail="Shift not found")
    
    # Get all officer attendance records for this shift
    attendances = db.query(ShiftAttendance).filter(
        ShiftAttendance.shift_id == shift_id
    ).all()
    
    officers_data = []
    officers_needing_rotation = []
    
    for attendance in attendances:
        # Query OfficerIncidentCount for this officer in this shift
        workload = db.query(OfficerIncidentCount).filter(
            OfficerIncidentCount.shift_id == shift_id,
            OfficerIncidentCount.officer_username == attendance.officer_username
        ).first()
        
        if workload:
            officer_info = OfficerWorkloadData(
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name,
                total_incidents=workload.incident_count_critical + workload.incident_count_high + 
                                workload.incident_count_medium + workload.incident_count_low,
                critical_incidents=workload.incident_count_critical,
                high_incidents=workload.incident_count_high,
                medium_incidents=workload.incident_count_medium,
                low_incidents=workload.incident_count_low,
                needs_rotation=workload.needs_rotation
            )
        else:
            # Create new workload entry if it doesn't exist
            workload = OfficerIncidentCount(
                shift_id=shift_id,
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name
            )
            db.add(workload)
            officer_info = OfficerWorkloadData(
                officer_username=attendance.officer_username,
                officer_name=attendance.officer_name,
                total_incidents=0,
                critical_incidents=0,
                high_incidents=0,
                medium_incidents=0,
                low_incidents=0,
                needs_rotation=False
            )
        
        officers_data.append(officer_info)
        if officer_info.needs_rotation:
            officers_needing_rotation.append(attendance.officer_username)
    
    db.commit()
    
    return OfficerWorkloadResponse(
        shift_id=shift_id,
        officers=officers_data,
        officers_needing_rotation=officers_needing_rotation
    )


@app.post("/api/police/incident-handled")
async def log_incident_handled(
    incident_id: str = Query(...),
    officer_username: str = Query(...),
    severity: str = Query(..., pattern="^(critical|high|medium|low)$"),
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_session),
):
    """Log that an officer has handled an incident (increments workload counter).
    
    Called when an officer is dispatched/finishes handling an incident.
    Tracks incident counts and updates rotation status if 3+ critical incidents.
    """
    # Get current active shift for officer
    shift = db.query(Shift).filter(
        Shift.status == "active"
    ).order_by(Shift.start_time.desc()).first()
    
    if not shift:
        raise HTTPException(
            status_code=404,
            detail="No active shift found"
        )
    
    # Get or create workload record for this officer
    workload = db.query(OfficerIncidentCount).filter(
        OfficerIncidentCount.shift_id == shift.id,
        OfficerIncidentCount.officer_username == officer_username
    ).first()
    
    if not workload:
        # Get officer name from attendance record
        attendance = db.query(ShiftAttendance).filter(
            ShiftAttendance.shift_id == shift.id,
            ShiftAttendance.officer_username == officer_username
        ).first()
        
        officer_name = attendance.officer_name if attendance else officer_username
        
        workload = OfficerIncidentCount(
            shift_id=shift.id,
            officer_username=officer_username,
            officer_name=officer_name
        )
        db.add(workload)
    
    # Increment appropriate counter
    if severity == "critical":
        workload.incident_count_critical += 1
    elif severity == "high":
        workload.incident_count_high += 1
    elif severity == "medium":
        workload.incident_count_medium += 1
    else:
        workload.incident_count_low += 1
    
    # Check if needs rotation (3+ critical incidents)
    if workload.incident_count_critical >= 3:
        workload.needs_rotation = True
    
    db.commit()
    
    return {
        "status": "logged",
        "officer_username": officer_username,
        "incident_id": incident_id,
        "severity": severity,
        "total_critical_incidents": workload.incident_count_critical,
        "needs_rotation": workload.needs_rotation
    }


@app.post("/api/incident/resolve")
async def resolve_incident(
    request: IncidentResolveRequest,
    current_user: dict = Depends(require_police_department_user()),
    db: Session = Depends(get_session),
):
    """Mark incident resolved and persist supervised feedback record for ML training."""
    resolved_at = request.resolved_at.astimezone(UTC) if request.resolved_at else datetime.now(UTC)
    severity_value = _normalize_severity(request.severity)
    district_id = current_user.get("district_id")
    if not district_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="district_id is required for incident resolution",
        )

    feedback = MLFeedback(
        incident_type=request.incident_type.strip(),
        zone=request.zone.strip(),
        time_of_day=resolved_at.hour,
        day_of_week=resolved_at.weekday(),
        response_time_minutes=float(request.response_time_minutes),
        severity=severity_value,
        outcome=request.outcome.strip(),
        created_at=resolved_at,
    )
    db.add(feedback)

    # Close any active assignment linked to the resolved incident.
    assignment = (
        db.query(PoliceDispatchAssignment)
        .filter(PoliceDispatchAssignment.incident_id == request.incident_id)
        .first()
    )
    returned_officer_id = None
    if assignment:
        assignment.status = "resolved"
        returned_officer_id = assignment.unit_id

    if returned_officer_id:
        officer_status = (
            db.query(OfficerDispatchStatus)
            .filter(OfficerDispatchStatus.officer_id == returned_officer_id)
            .first()
        )
        if officer_status:
            officer_status.status = "available"
            officer_status.assigned_incident_id = None
            officer_status.updated_at = resolved_at

        dispatch_record = (
            db.query(DispatchLog)
            .filter(
                DispatchLog.incident_id == request.incident_id,
                DispatchLog.officer_id == returned_officer_id,
            )
            .order_by(DispatchLog.assigned_at.desc())
            .first()
        )
        if dispatch_record:
            dispatch_record.status = "returned"

    db.commit()
    db.refresh(feedback)

    await emit_incident_updated(
        sio,
        district_id,
        {
            "id": request.incident_id,
            "severity": severity_value,
            "status": "resolved",
            "resolved_at": resolved_at.isoformat(),
        },
        update_type="resolved",
        actor=current_user.get("username", "Unknown Supervisor"),
    )
    if returned_officer_id:
        await emit_officer_status_changed(
            sio,
            district_id,
            {
                "id": returned_officer_id,
                "badge": returned_officer_id,
                "status": "available",
                "district_id": district_id,
                "incident_id": None,
            },
            actor=current_user.get("username", "Unknown Supervisor"),
        )

    return {
        "status": "resolved",
        "incident_id": request.incident_id,
        "feedback_id": feedback.id,
        "saved_training_record": True,
        "resolved_at": resolved_at.isoformat(),
    }


@app.post("/api/ml/retrain")
async def retrain_ml_model(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session),
):
    """Export feedback rows to CSV and trigger existing SVR retraining pipeline."""
    feedback_rows = db.query(MLFeedback).order_by(MLFeedback.created_at.asc()).all()
    if not feedback_rows:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No ML feedback available for retraining",
        )

    export_dir = os.path.join("data", "exports")
    os.makedirs(export_dir, exist_ok=True)
    export_name = f"ml_feedback_{datetime.now(UTC).strftime('%Y%m%d_%H%M%S')}.csv"
    export_path = os.path.join(export_dir, export_name)

    with open(export_path, "w", newline="", encoding="utf-8") as csv_file:
        writer = csv.DictWriter(
            csv_file,
            fieldnames=[
                "incident_type",
                "zone",
                "time_of_day",
                "day_of_week",
                "response_time_minutes",
                "severity",
                "outcome",
                "created_at",
            ],
        )
        writer.writeheader()
        for row in feedback_rows:
            writer.writerow(
                {
                    "incident_type": row.incident_type,
                    "zone": row.zone,
                    "time_of_day": row.time_of_day,
                    "day_of_week": row.day_of_week,
                    "response_time_minutes": row.response_time_minutes,
                    "severity": row.severity,
                    "outcome": row.outcome,
                    "created_at": row.created_at.isoformat() if row.created_at else None,
                }
            )

    try:
        from svr_model import train_svr

        train_svr()
    except Exception as exc:
        logger.error("ML retraining failed: %s", exc)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="ML retraining failed",
        )

    retrain_audit = MLRetrainAudit(
        retrained_at=datetime.now(UTC),
        retrained_by=getattr(current_user, "username", "admin"),
        feedback_rows=len(feedback_rows),
        export_path=export_path,
    )
    db.add(retrain_audit)
    db.commit()

    return {
        "status": "retrained",
        "feedback_rows": len(feedback_rows),
        "export_csv": export_path,
        "last_retrained": retrain_audit.retrained_at.isoformat(),
    }


@app.post("/auth/login", response_model=RoleToken)
async def role_login(
    login_data: RoleLoginRequest,
    db: Session = Depends(get_session),
    *,
    response: Response,
):
    """Login with role and return role-aware JWT. Auto-detects actual user role."""
    user = authenticate_user(db, login_data.username, login_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )

    user_department = (getattr(user, "department", None) or "").strip().lower()

    # Determine actual allowed role based on user's account
    if user.is_admin:
        actual_role = UserRole.admin
    elif user_department == "police":
        actual_role = UserRole.police_supervisor
    elif user_department == "logistics":
        actual_role = UserRole.logistics_manager
    else:
        actual_role = UserRole.user

    # Use actual role instead of requested role for normal users
    # This allows users to select "Normal User" even if system auto-detects their role
    role = actual_role
    district_id = login_data.district_id if login_data.district_id else None
    fleet_zone = login_data.fleet_zone if login_data.fleet_zone else None

    # Validate required fields for specific roles
    if role == UserRole.police_supervisor and not district_id:
        # Default to district_1 for police supervisors if not provided
        district_id = "district_1"

    if role == UserRole.logistics_manager and not fleet_zone:
        # Default to zone_default for logistics managers if not provided
        fleet_zone = "zone_default"

    access_token = create_role_access_token(
        username=user.username,
        role=role,
        district_id=district_id,
        fleet_zone=fleet_zone,
        expires_delta=timedelta(minutes=30 * 24 * 60),
    )
    # Set secure=False for localhost/HTTP development, True for production/HTTPS
    is_secure = os.getenv("SECURE_COOKIES", "false").lower() == "true"
    response.set_cookie(
        key="token",
        value=access_token,
        httponly=True,
        secure=is_secure,
        samesite="Lax",  # Use Lax for both secure and non-secure (None requires Secure=True)
        max_age=30 * 24 * 60 * 60,
        path="/",
    )
    return {
        "access_token": access_token,
        "token_type": "bearer",
        "role": role,
        "district_id": district_id,
        "fleet_zone": fleet_zone,
    }


@app.get("/api/auth/me", response_model=UserResponse)
async def get_current_user_info(current_user: User = Depends(get_current_active_user)):
    """Get current user information."""
    return UserResponse.model_validate(current_user)


# ============================================================================
# SAVED ROUTES
# ============================================================================

@app.post("/api/saved-routes")
@handle_db_errors
async def create_saved_route(
    route_data: SavedRouteCreate,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Save a route for the current user."""
    if not current_user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Please login to save routes."
        )
    
    origin_str = json.dumps(route_data.origin) if isinstance(route_data.origin, dict) else route_data.origin
    dest_str = json.dumps(route_data.destination) if isinstance(route_data.destination, dict) else route_data.destination
    
    saved_route = SavedRoute(
        user_id=current_user.id,
        route_name=route_data.route_name,
        origin=origin_str,
        destination=dest_str,
        route_preferences=route_data.route_preferences,
        is_favorite=False,
        share_token=secrets.token_urlsafe(16)
    )
    db.add(saved_route)
    db.commit()
    db.refresh(saved_route)
    return saved_route


@app.get("/api/saved-routes")
@handle_db_errors
async def get_saved_routes(
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session),
    favorites_only: bool = Query(False)
):
    """Get saved routes for current user."""
    if not current_user:
        return []
    query = db.query(SavedRoute).filter(SavedRoute.user_id == current_user.id)
    if favorites_only:
        query = query.filter(SavedRoute.is_favorite == True)
    routes = query.order_by(SavedRoute.last_used.desc()).all()
    return routes


@app.put("/api/saved-routes/{route_id}/favorite")
@handle_db_errors
async def toggle_favorite(
    route_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Toggle favorite status of a saved route."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required for this feature")
    route = db.query(SavedRoute).filter(
        SavedRoute.id == route_id,
        SavedRoute.user_id == current_user.id
    ).first()
    if not route:
        raise HTTPException(status_code=404, detail="Route not found")
    route.is_favorite = not route.is_favorite
    db.commit()
    return {"is_favorite": route.is_favorite}


@app.delete("/api/saved-routes/{route_id}")
@handle_db_errors
async def delete_saved_route(
    route_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Delete a saved route."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required for this feature")
    route = db.query(SavedRoute).filter(
        SavedRoute.id == route_id,
        SavedRoute.user_id == current_user.id
    ).first()
    if not route:
        raise HTTPException(status_code=404, detail="Route not found")
    db.delete(route)
    db.commit()
    return {"message": "Route deleted"}


@app.get("/api/share-route/{share_token}")
@handle_db_errors
async def get_shared_route(share_token: str, db: Session = Depends(get_session)):
    """Get a shared route by token."""
    route = db.query(SavedRoute).filter(SavedRoute.share_token == share_token).first()
    if not route:
        raise HTTPException(status_code=404, detail="Shared route not found")
    return route


# ============================================================================
# ADVANCED ANALYTICS
# ============================================================================

@app.get("/api/analytics/peak-hours/{route_id}")
async def get_peak_hours(
    route_id: str,
    days: int = Query(30, ge=1, le=365),
    db: Session = Depends(get_session)
):
    """Get peak hours analysis for a route."""
    return get_peak_hours_analysis(db, route_id, days)


@app.get("/api/analytics/day-of-week/{route_id}")
async def get_day_analysis(
    route_id: str,
    days: int = Query(90, ge=1, le=365),
    db: Session = Depends(get_session)
):
    """Get day of week analysis."""
    return get_day_of_week_analysis(db, route_id, days)


@app.get("/api/analytics/seasonal/{route_id}")
async def get_seasonal_analysis(
    route_id: str,
    months: int = Query(12, ge=1, le=24),
    db: Session = Depends(get_session)
):
    """Get seasonal trends."""
    return get_seasonal_trends(db, route_id, months)


@app.get("/api/analytics/reliability/{route_id}")
async def get_reliability(
    route_id: str,
    days: int = Query(30, ge=1, le=365),
    db: Session = Depends(get_session)
):
    """Get route reliability score."""
    return calculate_route_reliability(db, route_id, days)


@app.get("/api/analytics/predict/{route_id}")
async def get_prediction(
    route_id: str,
    hours_ahead: int = Query(24, ge=1, le=168),
    db: Session = Depends(get_session)
):
    """Predict future congestion."""
    return predict_future_congestion(db, route_id, hours_ahead)


@app.get("/api/analytics/hotspots")
async def get_hotspots(
    days: int = Query(7, ge=1, le=30),
    db: Session = Depends(get_session)
):
    """Get traffic hotspots."""
    return get_traffic_hotspots(db, days)


# ============================================================================
# EXPORT & REPORTING
# ============================================================================

@app.get("/api/export/csv/{route_id}")
async def export_csv(
    route_id: str,
    db: Session = Depends(get_session)
):
    """Export route data to CSV."""
    csv_content = export_to_csv(db, route_id)
    return StreamingResponse(
        io.BytesIO(csv_content.encode('utf-8')),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename=route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.csv"}
    )


@app.get("/api/export/excel/{route_id}")
async def export_excel(
    route_id: str,
    db: Session = Depends(get_session)
):
    """Export route data to Excel."""
    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        export_to_excel(db, route_id, tmp.name)
        return FileResponse(
            tmp.name,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.xlsx"
        )


@app.get("/api/export/pdf/{route_id}")
async def export_pdf(
    route_id: str,
    db: Session = Depends(get_session)
):
    """Export route data to PDF."""
    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
        export_to_pdf(db, route_id, tmp.name)
        return FileResponse(
            tmp.name,
            media_type="application/pdf",
            filename=f"route_{route_id}_{datetime.now(UTC).strftime('%Y%m%d')}.pdf"
        )


# ============================================================================
# NOTIFICATIONS
# ============================================================================

@app.get("/api/notifications")
@handle_db_errors
async def get_notifications(
    unread_only: bool = Query(False),
    limit: int = Query(50, ge=1, le=100),
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Get user notifications."""
    if not current_user:
        return []
    return get_user_notifications(db, current_user.id, unread_only, limit)


@app.put("/api/notifications/{notification_id}/read")
@handle_db_errors
async def mark_read(
    notification_id: int,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Mark notification as read."""
    if not current_user:
        raise HTTPException(status_code=401, detail="Login required")
    success = mark_notification_read(db, notification_id, current_user.id)
    if not success:
        raise HTTPException(status_code=404, detail="Notification not found")
    return {"message": "Notification marked as read"}


@app.post("/api/notifications/check-alerts")
@handle_db_errors
async def check_alerts(
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Check for traffic alerts on saved routes."""
    if not current_user:
        return {"alerts": 0, "notifications": []}
    alerts = check_traffic_alerts(db, current_user.id)
    return {"alerts": len(alerts), "notifications": alerts}


# ============================================================================
# REAL-TIME FEATURES
# ============================================================================

@app.get("/api/realtime/incidents")
async def get_incidents(
    lat: float = Query(..., description="Latitude"),
    lon: float = Query(..., description="Longitude"),
    radius: int = Query(5000, ge=100, le=50000)
):
    """Get traffic incidents near a location."""
    # Validate coordinates
    if not (-90 <= lat <= 90):
        raise HTTPException(status_code=400, detail="Invalid latitude")
    if not (-180 <= lon <= 180):
        raise HTTPException(status_code=400, detail="Invalid longitude")
    
    return get_traffic_incidents(lat, lon, radius)


@app.post("/api/realtime/monitor/{route_id}")
async def monitor_route(
    route_id: str,
    background_tasks: BackgroundTasks,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Monitor route for changes."""
    change = monitor_route_changes(db, route_id)
    if change:
        if current_user:
            from notifications import create_notification
            create_notification(
                db, current_user.id, 'traffic_alert',
                f"Route Change: {route_id}",
                f"Route travel time changed by {change['change_percent']}%",
                route_id
            )
    return change or {"message": "No significant changes detected"}


# ============================================================================
# ROUTE RATINGS & SOCIAL
# ============================================================================

@app.post("/api/ratings")
@handle_db_errors
async def create_rating(
    rating_data: RouteRatingCreate,
    current_user: Optional[User] = Depends(get_optional_user),
    db: Session = Depends(get_session)
):
    """Rate a route."""
    user_id = current_user.id if current_user else None
    rating = RouteRating(
        user_id=user_id,
        route_id=rating_data.route_id,
        rating=rating_data.rating,
        review=rating_data.review
    )
    db.add(rating)
    db.commit()
    db.refresh(rating)
    return rating


@app.get("/api/ratings/{route_id}")
@handle_db_errors
async def get_ratings(route_id: str, db: Session = Depends(get_session)):
    """Get ratings for a route."""
    ratings = db.query(RouteRating).filter(RouteRating.route_id == route_id).all()
    if not ratings:
        return {"average_rating": 0, "count": 0, "ratings": []}
    
    avg_rating = sum(r.rating for r in ratings) / len(ratings)
    return {
        "average_rating": round(avg_rating, 2),
        "count": len(ratings),
        "ratings": ratings
    }


# ============================================================================
# ADMIN DASHBOARD
# ============================================================================

@app.get("/api/admin/stats")
@handle_db_errors
async def get_admin_stats(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session)
):
    """Get admin statistics (admin only)."""
    total_users = db.query(User).count()
    total_routes = db.query(AnalysisResult).count()
    total_saved_routes = db.query(SavedRoute).count()
    total_ratings = db.query(RouteRating).count()
    
    recent_routes = db.query(AnalysisResult).order_by(AnalysisResult.timestamp.desc()).limit(10).all()
    recent_users = db.query(User).order_by(User.created_at.desc()).limit(5).all()
    
    return {
        "total_users": total_users,
        "total_route_analyses": total_routes,
        "total_saved_routes": total_saved_routes,
        "total_ratings": total_ratings,
        "cache_stats": get_cache_stats(),
        "recent_activity": {
            "routes": len(recent_routes),
            "new_users": len(recent_users)
        }
    }


@app.get("/api/admin/route-analysis")
async def get_all_route_analyses(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session),
    filter_period: Optional[str] = Query(None, alias="filter"),
    skip: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=1000)
):
    """Get all route analyses with optional filtering (admin only)."""
    import json
    from datetime import datetime, timedelta, UTC
    
    try:
        # Build query
        query = db.query(AnalysisResult)
        
        # Apply time filter if specified
        if filter_period:
            now = datetime.now(UTC)
            if filter_period == "today":
                start_date = now.replace(hour=0, minute=0, second=0, microsecond=0)
                query = query.filter(AnalysisResult.timestamp >= start_date)
            elif filter_period == "week":
                start_date = now - timedelta(days=7)
                query = query.filter(AnalysisResult.timestamp >= start_date)
            elif filter_period == "month":
                start_date = now - timedelta(days=30)
                query = query.filter(AnalysisResult.timestamp >= start_date)
        
        # Get total count before pagination
        total_count = query.count()
        
        # Apply pagination and ordering
        routes = query.order_by(AnalysisResult.timestamp.desc()).offset(skip).limit(limit).all()
        
        # Format response
        route_data = []
        for r in routes:
            try:
                origin = json.loads(r.origin) if isinstance(r.origin, str) and r.origin.startswith('{') else {"name": str(r.origin) if r.origin else ""}
                dest = json.loads(r.destination) if isinstance(r.destination, str) and r.destination.startswith('{') else {"name": str(r.destination) if r.destination else ""}
            except:
                origin = {"name": str(r.origin) if r.origin else ""}
                dest = {"name": str(r.destination) if r.destination else ""}
            
            route_name = f"{origin.get('name', '')} → {dest.get('name', '')}"
            
            delay_val = r.delay_s
            if delay_val is None or delay_val == 0:
                if r.travel_time_s and r.no_traffic_s:
                    delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                else:
                    delay_val = 0
            
            route_data.append({
                "id": r.id,
                "route": route_name,
                "route_id": r.route_id,
                "travel_time_s": r.travel_time_s,
                "delay_s": delay_val,
                "length_m": r.length_m,
                "calculated_cost": r.calculated_cost,
                "ml_predicted": r.ml_predicted,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "origin": origin,
                "destination": dest
            })
        
        # Calculate statistics from filtered results
        all_routes_for_stats = query.all()
        
        if all_routes_for_stats:
            travel_times = [r.travel_time_s for r in all_routes_for_stats if r.travel_time_s is not None]
            delays = []
            for r in all_routes_for_stats:
                delay_val = r.delay_s
                if delay_val is None or delay_val == 0:
                    if r.travel_time_s and r.no_traffic_s:
                        delay_val = max(0, r.travel_time_s - r.no_traffic_s)
                    else:
                        delay_val = 0
                if delay_val > 0:
                    delays.append(delay_val)
            costs = [r.calculated_cost for r in all_routes_for_stats if r.calculated_cost is not None]
            
            stats = {
                "total": total_count,
                "avg_travel_time": sum(travel_times) / len(travel_times) if travel_times else 0,
                "avg_delay": sum(delays) / len(delays) if delays else 0,
                "avg_cost": sum(costs) / len(costs) if costs else 0
            }
        else:
            stats = {
                "total": 0,
                "avg_travel_time": 0,
                "avg_delay": 0,
                "avg_cost": 0
            }
        
        return {
            "routes": route_data,
            "stats": stats,
            "pagination": {
                "skip": skip,
                "limit": limit,
                "total": total_count
            }
        }
        
    except Exception as e:
        logger.error(f"Error fetching route analyses: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch route analyses: {str(e)}")


@app.get("/api/admin/users")
async def get_all_users(
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session),
    skip: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=1000)
):
    """Get users filtered by department (admin only - see only users in their department)."""
    # Get admin's department
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    
    # Build query based on admin's department
    query = db.query(User)
    
    # If admin is super admin (department='admin'), show all users
    # Otherwise, show only users in their own department
    if admin_department != "admin":
        # Filter by department - but first check if there are department-specific admins
        # Police department admins see police users
        # Logistics department admins see logistics users
        query = query.filter(User.department == admin_department)
    # Super admins (department='admin') see everyone
    
    users = query.offset(skip).limit(limit).all()
    return [UserResponse.model_validate(u) for u in users]


@app.put("/api/admin/users/{user_id}/activate")
async def toggle_user_status(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session)
):
    """Activate/deactivate a user (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    user.is_active = not user.is_active
    db.commit()
    return {"is_active": user.is_active, "message": f"User {'activated' if user.is_active else 'deactivated'}"}


@app.put("/api/admin/users/{user_id}/admin")
async def toggle_admin_status(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session)
):
    """Grant/revoke admin privileges (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user.id == current_user.id:
        raise HTTPException(status_code=400, detail="Cannot change your own admin status")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    user.is_admin = not user.is_admin
    db.commit()
    return {"is_admin": user.is_admin, "message": f"Admin privileges {'granted' if user.is_admin else 'revoked'}"}


@app.put("/api/admin/users/{user_id}")
async def update_user(
    user_id: int,
    user_update: UserUpdate,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session)
):
    """Update user details (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    # Check if admin can modify this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only manage users in your department"
        )
    
    # Update username if provided
    if user_update.username is not None:
        existing_user = db.query(User).filter(User.username == user_update.username, User.id != user_id).first()
        if existing_user:
            raise HTTPException(status_code=400, detail="Username already taken")
        user.username = user_update.username
    
    # Update email if provided
    if user_update.email is not None:
        existing_user = db.query(User).filter(User.email == user_update.email, User.id != user_id).first()
        if existing_user:
            raise HTTPException(status_code=400, detail="Email already taken")
        user.email = user_update.email
    
    # Update full name if provided
    if user_update.full_name is not None:
        user.full_name = user_update.full_name
    
    # Update active status if provided
    if user_update.is_active is not None:
        user.is_active = user_update.is_active
    
    # Update admin status if provided
    if user_update.is_admin is not None:
        if user.id == current_user.id and not user_update.is_admin:
            raise HTTPException(status_code=400, detail="Cannot remove your own admin privileges")
        user.is_admin = user_update.is_admin
    
    # Update password if provided and not empty
    if user_update.password is not None and user_update.password.strip() != "":
        if len(user_update.password) < 8:
            raise HTTPException(status_code=400, detail="Password must be at least 8 characters")
        user.hashed_password = get_password_hash(user_update.password)
    
    try:
        db.commit()
        db.refresh(user)
        return UserResponse.model_validate(user)
    except Exception as e:
        db.rollback()
        logger.error(f"Error updating user: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to update user: {str(e)}")


@app.delete("/api/admin/users/{user_id}")
async def delete_user(
    user_id: int,
    current_user: User = Depends(get_current_admin_user),
    db: Session = Depends(get_session)
):
    """Delete a user (admin only - department-restricted)."""
    user = db.query(User).filter(User.id == user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    
    if user.id == current_user.id:
        raise HTTPException(status_code=400, detail="Cannot delete your own account")
    
    # Check if admin can delete this user (same department or super admin)
    admin_department = (getattr(current_user, "department", None) or "admin").strip().lower()
    user_department = (getattr(user, "department", None) or "general").strip().lower()
    
    if admin_department != "admin" and admin_department != user_department:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You can only delete users in your department"
        )
    
    db.delete(user)
    db.commit()
    return {"message": "User deleted successfully"}


@app.get("/api/user/stats")
@handle_db_errors
async def get_user_stats(
    current_user: User = Depends(get_current_active_user),
    db: Session = Depends(get_session)
):
    """Get user-specific statistics."""
    saved_routes_count = db.query(SavedRoute).filter(SavedRoute.user_id == current_user.id).count()
    analyses_count = db.query(AnalysisResult).filter(AnalysisResult.user_id == current_user.id).count()
    ratings_count = db.query(RouteRating).filter(RouteRating.user_id == current_user.id).count()
    
    recent_analyses = db.query(AnalysisResult).filter(
        AnalysisResult.user_id == current_user.id
    ).order_by(AnalysisResult.timestamp.desc()).limit(10).all()
    
    return {
        "saved_routes": saved_routes_count,
        "analyses": analyses_count,
        "ratings": ratings_count,
        "recent_analyses": [
            {
                "route_id": r.route_id,
                "timestamp": r.timestamp.isoformat() if r.timestamp else None,
                "travel_time": r.travel_time_s,
                "cost": r.calculated_cost
            }
            for r in recent_analyses
        ]
    }


# ============================================================================
# CACHE MANAGEMENT
# ============================================================================

@app.post("/api/cache/clear")
async def clear_route_cache(
    pattern: Optional[str] = None,
    current_user: User = Depends(get_current_admin_user)
):
    """Clear route cache (admin only)."""
    clear_cache(pattern=pattern)
    return {"message": "Cache cleared"}


@app.get("/api/cache/stats")
async def get_cache_statistics():
    """Get cache statistics."""
    return get_cache_stats()


# ============================================================================
# INTEGRATION ENDPOINTS
# ============================================================================

@app.get("/api/integration/navigation/{route_id}")
@handle_db_errors
async def get_navigation_links(
    route_id: str,
    route_index: int = Query(0),
    db: Session = Depends(get_session)
):
    """Get navigation app links (Google Maps, Waze)."""
    result = db.query(AnalysisResult).filter(
        AnalysisResult.route_id.like(f"{route_id}%")
    ).order_by(AnalysisResult.timestamp.desc()).first()
    
    if not result:
        raise HTTPException(status_code=404, detail="Route not found")
    
    try:
        origin = json.loads(result.origin) if result.origin and result.origin.startswith('{') else {"name": result.origin}
        dest = json.loads(result.destination) if result.destination and result.destination.startswith('{') else {"name": result.destination}
    except:
        origin = {"name": result.origin}
        dest = {"name": result.destination}
    
    origin_lat = origin.get('lat', 0)
    origin_lon = origin.get('lon', 0)
    dest_lat = dest.get('lat', 0)
    dest_lon = dest.get('lon', 0)
    
    google_maps = f"https://www.google.com/maps/dir/{origin_lat},{origin_lon}/{dest_lat},{dest_lon}"
    waze = f"https://waze.com/ul?ll={dest_lat},{dest_lon}&navigate=yes"
    
    return {
        "google_maps": google_maps,
        "waze": waze,
        "origin": origin,
        "destination": dest
    }
